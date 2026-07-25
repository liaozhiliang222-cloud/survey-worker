"""RenderAudit 编排与 Final Validation。

编排完整 QA 管线：
  1. 内容完整性检查
  2. 文本容量检查
  3. 元素边界检查
  4. 图表空间检查
  5. 页面密度检查
  6. 模板重复检查
  7. Final Validation

Final Validation 确认：
  - 内容：SlideBrief 核心内容存在、evidence 未丢失
  - 布局：无越界、无空核心区域、无文本溢出
  - 图表：categories/series/values 长度一致
"""

from __future__ import annotations

from typing import Any

from pptx import Presentation
from pptx.util import Emu

from .models import (
    FinalValidationResult,
    QAIssue,
    IssueType,
    RenderAction,
    Severity,
    SlideQAReport,
    PageDensityScore,
)
from .auto_fix import run_auto_fix, check_layout_repeat
from .text_checks import TextBoxInfo
from .layout_checks import ShapeInfo, PageInfo
from .chart_checks import ChartGeometry


def _extract_slide_info(slide, slide_width: int, slide_height: int, index: int) -> dict[str, Any]:
    """从 python-pptx Slide 对象提取 QA 所需的几何信息。"""
    slide_id = f"slide_{index + 1:02d}"
    textboxes: list[TextBoxInfo] = []
    shapes: list[ShapeInfo] = []
    charts: list[ChartGeometry] = []

    for shape in slide.shapes:
        left = int(shape.left) if shape.left is not None else 0
        top = int(shape.top) if shape.top is not None else 0
        width = int(shape.width) if shape.width is not None else 0
        height = int(shape.height) if shape.height is not None else 0
        shape_id = shape.shape_id
        shape_name = shape.name or f"shape_{shape_id}"

        # 分类
        if shape.has_chart:
            chart = shape.chart
            categories = []
            series_names = []
            try:
                plot = chart.plots[0]
                categories = [str(cat) for cat in plot.categories]
                series_names = [s.name or f"Series {i}" for i, s in enumerate(plot.series)]
            except (IndexError, AttributeError):
                pass
            chart_type_str = "column"
            try:
                ct = chart.chart_type
                if ct is not None:
                    ct_name = str(ct).lower()
                    if "bar" in ct_name:
                        chart_type_str = "bar"
                    elif "line" in ct_name:
                        chart_type_str = "line"
                    elif "pie" in ct_name:
                        chart_type_str = "pie"
            except Exception:
                pass
            charts.append(ChartGeometry(
                element_id=shape_name,
                chart_type=chart_type_str,
                categories=categories,
                series_names=series_names,
                width_emu=width,
                height_emu=height,
            ))
            shapes.append(ShapeInfo(
                element_id=shape_name, shape_type="chart",
                left=left, top=top, width=width, height=height,
            ))
        elif shape.has_text_frame:
            text = shape.text_frame.text or ""
            is_title = "title" in shape_name.lower() or top < int(slide_height * 0.15)
            font_size_pt = 18.0
            try:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            font_size_pt = run.font.size.pt
                            break
                    break
            except (AttributeError, TypeError):
                pass
            textboxes.append(TextBoxInfo(
                element_id=shape_name,
                text=text,
                font_size_pt=font_size_pt,
                box_width_emu=width,
                box_height_emu=height,
                is_title=is_title,
            ))
            shapes.append(ShapeInfo(
                element_id=shape_name, shape_type="textbox",
                left=left, top=top, width=width, height=height,
                text=text,
            ))
        elif shape.has_table:
            shapes.append(ShapeInfo(
                element_id=shape_name, shape_type="table",
                left=left, top=top, width=width, height=height,
            ))
        else:
            shapes.append(ShapeInfo(
                element_id=shape_name, shape_type="other",
                left=left, top=top, width=width, height=height,
            ))

    return {
        "slide_id": slide_id,
        "textboxes": textboxes,
        "shapes": shapes,
        "charts": charts,
    }


def run_render_qa(
    prs: Presentation,
    *,
    slide_briefs: list[dict[str, Any]] | None = None,
) -> FinalValidationResult:
    """对已渲染的 Presentation 执行完整 QA 管线。

    Args:
        prs: python-pptx Presentation 对象（已渲染完成）。
        slide_briefs: 可选的 SlideBrief 字典列表，用于内容完整性校验。

    Returns:
        FinalValidationResult 包含所有 issues、actions 和评分。
    """
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    result = FinalValidationResult()

    briefs_by_id: dict[str, dict] = {}
    if slide_briefs:
        for brief in slide_briefs:
            sid = brief.get("slide_id", "")
            if sid:
                briefs_by_id[sid] = brief

    slide_type_sequence: list[tuple[str, str, str]] = []

    for index, slide in enumerate(prs.slides):
        info = _extract_slide_info(slide, slide_width, slide_height, index)
        slide_id = info["slide_id"]

        # 从 brief 获取元数据
        brief = briefs_by_id.get(slide_id, {})
        title = brief.get("title", "")
        claim = brief.get("claim", "")
        evidence_fact_ids = brief.get("evidence_fact_ids", [])
        source_references = brief.get("source_references", [])
        slide_type = brief.get("slide_type", "")
        layout_family = brief.get("layout_family", "")

        if slide_type:
            slide_type_sequence.append((slide_id, slide_type, layout_family))

        # 构建 PageInfo
        page_info = PageInfo(
            slide_id=slide_id,
            shapes=info["shapes"],
            slide_width=slide_width,
            slide_height=slide_height,
        )

        # 构建 charts_data 用于内容完整性检查
        charts_data = [
            {"categories": c.categories, "series": c.series_names}
            for c in info["charts"]
        ]

        # 执行自动修复管线
        slide_report = run_auto_fix(
            slide_id=slide_id,
            textboxes=info["textboxes"],
            shapes=info["shapes"],
            charts=info["charts"],
            page_info=page_info,
            title=title,
            claim=claim,
            evidence_fact_ids=evidence_fact_ids,
            source_references=source_references,
            charts_data=charts_data if charts_data else None,
            slide_width=slide_width,
            slide_height=slide_height,
        )
        result.slide_reports.append(slide_report)
        result.issues.extend(slide_report.issues)
        result.actions.extend(slide_report.actions)

    # ─── 6. 模板重复检查 ────────────────────────────────────
    if slide_type_sequence:
        repeat_issues, repeat_actions = check_layout_repeat(slide_type_sequence)
        result.issues.extend(repeat_issues)
        result.actions.extend(repeat_actions)

    # ─── 7. Final Validation ────────────────────────────────
    result = final_validation(result)
    return result


def final_validation(result: FinalValidationResult) -> FinalValidationResult:
    """最终验证：汇总所有检查结果，计算评分。

    确认：
      - 内容：核心内容存在、evidence 未丢失
      - 布局：无越界、无空核心区域、无文本溢出
      - 图表：数据一致性
    """
    high_issues = [i for i in result.issues if i.severity == Severity.HIGH]
    medium_issues = [i for i in result.issues if i.severity == Severity.MEDIUM]

    # 分类检查
    content_types = {IssueType.CONTENT_MISSING, IssueType.EMPTY_ELEMENT}
    layout_types = {IssueType.OUT_OF_BOUND, IssueType.TEXT_OVERFLOW, IssueType.PAGE_OVERLOAD}
    chart_types = {IssueType.CHART_DENSITY_HIGH, IssueType.CHART_TOO_SMALL, IssueType.LEGEND_TOO_LONG}

    content_high = [i for i in high_issues if i.issue_type in content_types]
    layout_high = [i for i in high_issues if i.issue_type in layout_types]
    chart_high = [i for i in high_issues if i.issue_type in chart_types]

    result.content_ok = not content_high
    result.layout_ok = not layout_high
    result.chart_ok = not chart_high
    result.passed = result.content_ok and result.layout_ok and result.chart_ok

    # 评分：100 分起，每个 HIGH -15，MEDIUM -5，LOW -1
    score = 100
    score -= len(high_issues) * 15
    score -= len(medium_issues) * 5
    score -= len([i for i in result.issues if i.severity == Severity.LOW]) * 1
    # 自动修复加分（每个修复 +2，最多恢复 20 分）
    fix_bonus = min(20, len(result.actions) * 2)
    score += fix_bonus
    result.score = max(0, min(100, score))

    return result
