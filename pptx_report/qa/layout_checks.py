"""页面边界检查与密度评分。

边界检查：确保所有 Shape/TextBox/Chart/Image 在页面范围内。
密度评分：基于字符数、文本框数、图表数、元素数、空白比例计算页面密度。
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import Any

from pptx.util import Emu, Inches

from .models import (
    QAIssue,
    IssueType,
    PageDensityScore,
    RenderAction,
    Severity,
)

# ─── 默认页面尺寸（16:9 宽屏）───────────────────────────────
DEFAULT_SLIDE_WIDTH = Inches(13.333)
DEFAULT_SLIDE_HEIGHT = Inches(7.5)

# ─── 密度阈值 ───────────────────────────────────────────────
DENSITY_THRESHOLDS = {
    "char_count": {"LOW": 200, "MEDIUM": 600, "HIGH": 1200},
    "element_count": {"LOW": 5, "MEDIUM": 12, "HIGH": 20},
    "text_box_count": {"LOW": 3, "MEDIUM": 7, "HIGH": 12},
    "chart_count": {"LOW": 1, "MEDIUM": 3, "HIGH": 5},
}


@dataclass
class ShapeInfo:
    """Shape 几何信息。"""

    element_id: str
    shape_type: str  # textbox, chart, image, table, group
    left: int        # EMU
    top: int         # EMU
    width: int       # EMU
    height: int      # EMU
    text: str = ""
    is_core: bool = True  # 是否为核心内容元素


@dataclass
class PageInfo:
    """页面信息，用于密度计算。"""

    slide_id: str
    shapes: list[ShapeInfo]
    slide_width: int = int(DEFAULT_SLIDE_WIDTH)
    slide_height: int = int(DEFAULT_SLIDE_HEIGHT)


def check_boundaries(
    slide_id: str,
    shapes: list[ShapeInfo],
    slide_width: int = int(DEFAULT_SLIDE_WIDTH),
    slide_height: int = int(DEFAULT_SLIDE_HEIGHT),
) -> tuple[list[QAIssue], list[RenderAction], list[ShapeInfo]]:
    """检查所有元素是否在页面边界内，越界则自动修复。

    修复策略：
      1. 缩小（等比缩放到页面内）
      2. 居中（移回页面中心）
      3. 调整 margin（贴边放置）

    Returns:
        (issues, actions, fixed_shapes)
    """
    issues: list[QAIssue] = []
    actions: list[RenderAction] = []
    fixed_shapes: list[ShapeInfo] = []

    for shape in shapes:
        out_of_bound = False
        reasons: list[str] = []

        if shape.left < 0:
            out_of_bound = True
            reasons.append(f"left={shape.left}")
        if shape.top < 0:
            out_of_bound = True
            reasons.append(f"top={shape.top}")
        if shape.left + shape.width > slide_width:
            out_of_bound = True
            reasons.append(f"right={shape.left + shape.width} > {slide_width}")
        if shape.top + shape.height > slide_height:
            out_of_bound = True
            reasons.append(f"bottom={shape.top + shape.height} > {slide_height}")

        if not out_of_bound:
            fixed_shapes.append(shape)
            continue

        # ─── 自动修复 ───────────────────────────────────────
        new_left = shape.left
        new_top = shape.top
        new_width = shape.width
        new_height = shape.height
        fix_action = "adjust_position"

        # 策略 1：如果尺寸超出页面，等比缩小
        if shape.width > slide_width or shape.height > slide_height:
            scale = min(
                slide_width / max(1, shape.width),
                slide_height / max(1, shape.height),
            ) * 0.95
            new_width = int(shape.width * scale)
            new_height = int(shape.height * scale)
            fix_action = "scale_down"

        # 策略 2：确保位置在页面内
        if new_left < 0:
            new_left = 0
        if new_top < 0:
            new_top = 0
        if new_left + new_width > slide_width:
            new_left = max(0, slide_width - new_width)
        if new_top + new_height > slide_height:
            new_top = max(0, slide_height - new_height)

        issues.append(QAIssue(
            slide_id=slide_id,
            issue_type=IssueType.OUT_OF_BOUND,
            severity=Severity.MEDIUM,
            message=f"元素越界：{', '.join(reasons)}",
            element_id=shape.element_id,
            details={"reasons": reasons},
        ))
        actions.append(RenderAction(
            slide_id=slide_id,
            issue_type=IssueType.OUT_OF_BOUND.value,
            severity=Severity.MEDIUM.value,
            action=fix_action,
            before=f"({shape.left}, {shape.top}, {shape.width}x{shape.height})",
            after=f"({new_left}, {new_top}, {new_width}x{new_height})",
            reason="元素超出页面边界，自动调整位置和尺寸",
        ))
        fixed_shapes.append(ShapeInfo(
            element_id=shape.element_id,
            shape_type=shape.shape_type,
            left=new_left,
            top=new_top,
            width=new_width,
            height=new_height,
            text=shape.text,
            is_core=shape.is_core,
        ))

    return issues, actions, fixed_shapes


def compute_page_density(page: PageInfo) -> PageDensityScore:
    """计算页面密度评分。

    输入：字符数量、文本框数量、图表数量、元素数量、空白比例。
    输出：LOW / MEDIUM / HIGH / OVERLOAD
    """
    total_chars = sum(len(s.text) for s in page.shapes if s.text)
    text_box_count = sum(1 for s in page.shapes if s.shape_type == "textbox")
    chart_count = sum(1 for s in page.shapes if s.shape_type == "chart")
    element_count = len(page.shapes)

    # 计算空白比例
    total_area = page.slide_width * page.slide_height
    occupied_area = sum(s.width * s.height for s in page.shapes)
    whitespace_ratio = max(0.0, 1.0 - (occupied_area / max(1, total_area)))

    # 综合评分
    score = 0
    thresholds = DENSITY_THRESHOLDS

    if total_chars > thresholds["char_count"]["HIGH"]:
        score += 3
    elif total_chars > thresholds["char_count"]["MEDIUM"]:
        score += 2
    elif total_chars > thresholds["char_count"]["LOW"]:
        score += 1

    if element_count > thresholds["element_count"]["HIGH"]:
        score += 3
    elif element_count > thresholds["element_count"]["MEDIUM"]:
        score += 2
    elif element_count > thresholds["element_count"]["LOW"]:
        score += 1

    if text_box_count > thresholds["text_box_count"]["HIGH"]:
        score += 2
    elif text_box_count > thresholds["text_box_count"]["MEDIUM"]:
        score += 1

    if chart_count > thresholds["chart_count"]["HIGH"]:
        score += 2
    elif chart_count > thresholds["chart_count"]["MEDIUM"]:
        score += 1

    # 空白比例过低（< 15%）加分
    if whitespace_ratio < 0.15:
        score += 2
    elif whitespace_ratio < 0.25:
        score += 1

    if score >= 8:
        return PageDensityScore.OVERLOAD
    elif score >= 5:
        return PageDensityScore.HIGH
    elif score >= 2:
        return PageDensityScore.MEDIUM
    return PageDensityScore.LOW


def check_page_density(page: PageInfo) -> list[QAIssue]:
    """检查页面密度，OVERLOAD 时生成 issue。"""
    density = compute_page_density(page)
    if density == PageDensityScore.OVERLOAD:
        return [QAIssue(
            slide_id=page.slide_id,
            issue_type=IssueType.PAGE_OVERLOAD,
            severity=Severity.HIGH,
            message="页面内容过载，建议调整布局或拆页",
            details={
                "density": density.value,
                "element_count": len(page.shapes),
                "total_chars": sum(len(s.text) for s in page.shapes if s.text),
            },
        )]
    elif density == PageDensityScore.HIGH:
        return [QAIssue(
            slide_id=page.slide_id,
            issue_type=IssueType.PAGE_OVERLOAD,
            severity=Severity.LOW,
            message="页面内容较多，建议关注排版质量",
            details={"density": density.value},
        )]
    return []
