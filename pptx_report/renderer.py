"""渲染器与模板支持（统一入口）。

:class:`ReportRenderer` 接收 :class:`~pptx_report.model.ReportSpec`（数据），
输出 .pptx 文件。支持三种用法：
  1. 无模板直接绘制（默认，最灵活，支持全部图表 / 布局）；
  2. 加载现有 .pptx 作为设计底（继承母版背景与配色）；
  3. 在模板的预设占位符上按名称填充文本（见 :meth:`fill_named_placeholders`）。
"""

from __future__ import annotations

import io
import os
from typing import Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

from .exceptions import RenderingError, TemplateNotFoundError
from .model import ReportSpec, MultiGroupBarPageContent
from .pages import (
    build_appendix,
    build_chart_page,
    build_cover,
    build_exec_summary,
    build_multi_group_bar_page,
    build_toc,
)
from .theme import Theme
from .template import build_template_mapping
from .utils import set_slide_background


class ReportRenderer:
    """把 ReportSpec 渲染为 PowerPoint 文件。"""

    def __init__(self, theme: Optional[Theme] = None,
                 template_path: Optional[str] = None,
                 progress_callback=None):
        """
        Args:
            theme: 视觉主题；为空则用 :class:`~pptx_report.theme.Theme` 默认主题。
            template_path: 可选 .pptx 模板路径，作为设计底（继承母版 / 背景）。
        """
        self.theme = theme or Theme()
        self.template_path = template_path
        self.progress_callback = progress_callback
        self._prs = None
        self._slide_w = None
        self._slide_h = None
        self._template_mapping = None

    def render(self, spec: ReportSpec, output_path: str) -> str:
        """渲染并保存到 output_path，返回该路径。"""
        # 目录为空时自动提取章节标题
        if not spec.toc.sections:
            spec.toc.sections = spec.auto_toc()
        spec.validate()

        self._prs = self._build_presentation()
        dims = (self._slide_w, self._slide_h)
        total_units = 3 + len(spec.chart_pages) + (1 if spec.appendix is not None else 0)
        completed = 0

        def report_page(message):
            nonlocal completed
            completed += 1
            if self.progress_callback:
                percent = 48 + round(42 * completed / max(1, total_units))
                self.progress_callback(percent, message)
        try:
            self._add_cover(spec.cover, dims)
            report_page("正在绘制封面")
            self._add_toc(spec.toc, dims)
            report_page("正在绘制目录")
            self._add_exec_summary(spec.executive_summary, dims)
            report_page("正在绘制执行摘要")
            for index, page in enumerate(spec.chart_pages, 1):
                self._add_chart_page(page, dims)
                report_page(f"正在绘制数据页 {index}/{len(spec.chart_pages)}")
            if spec.appendix is not None:
                self._add_appendix(spec.appendix, dims)
                report_page("正在绘制附录")
            if self.progress_callback:
                self.progress_callback(94, "正在打包演示文稿")
            self._prs.save(output_path)
            if self.progress_callback:
                self.progress_callback(97, "演示文稿已生成")
        except RenderingError:
            raise
        except Exception as exc:  # noqa: BLE001
            raise RenderingError(f"渲染失败: {exc}") from exc
        return output_path

    # ------------------------- 内部 -------------------------
    def _build_presentation(self) -> Presentation:
        if self.template_path:
            if not os.path.exists(self.template_path):
                raise TemplateNotFoundError(self.template_path)
            prs = Presentation(self.template_path)
            self._template_mapping = build_template_mapping(prs)
            # 上传模板只提供母版、版式、背景和主题。移除模板中的示例页，避免旧内容
            # 出现在新报告前面；版式与母版关系仍会保留。
            for slide_id in list(prs.slides._sldIdLst):
                relationship_id = slide_id.rId
                prs.slides._sldIdLst.remove(slide_id)
                prs.part.drop_rel(relationship_id)
            self._slide_w = prs.slide_width
            self._slide_h = prs.slide_height
        else:
            prs = Presentation()
            # 默认 4:3，改为 16:9 宽屏
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            self._slide_w = prs.slide_width
            self._slide_h = prs.slide_height
        return prs

    def _blank_slide(self, role="chart"):
        prs = self._prs
        layout = None
        if self.template_path and self._template_mapping:
            role_map = self._template_mapping.get("roles", {}).get(role) or {}
            layout_index = role_map.get("layout_index")
            if isinstance(layout_index, int) and 0 <= layout_index < len(prs.slide_layouts):
                layout = prs.slide_layouts[layout_index]
        # 未使用语义映射时，优先选择名为 Blank / 空白 的版式。
        if layout is None:
            for name in ("Blank", "空白", "Office 主题"):  # noqa: S105
                for lay in prs.slide_layouts:
                    if lay.name == name:
                        layout = lay
                        break
                if layout is not None:
                    break
        if layout is None:
            try:
                layout = prs.slide_layouts[6]  # 多数模板的「空白」版式
            except IndexError:
                layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        # 版式可能带标题/正文占位符。当前报告由原生形状重新绘制，因此清除占位符，
        # 但保留母版背景、Logo 和页脚等母版级元素。
        if self.template_path:
            for shape in list(slide.placeholders):
                element = shape._element
                element.getparent().remove(element)
            # Some imported layouts contain a user-drawn logo picture that
            # PowerPoint renders inconsistently across appended slides. Add an
            # identical slide-level copy so every generated page keeps branding.
            header_picture_sources = list(layout.shapes)
            if role != "chart":
                chart_map = self._template_mapping.get("roles", {}).get("chart") or {}
                chart_layout_index = chart_map.get("layout_index")
                if (
                    isinstance(chart_layout_index, int)
                    and 0 <= chart_layout_index < len(prs.slide_layouts)
                ):
                    header_picture_sources.extend(
                        prs.slide_layouts[chart_layout_index].shapes
                    )
            copied_header_pictures = set()
            for layout_shape in header_picture_sources:
                if (
                    layout_shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                    and int(layout_shape.left) > int(self._slide_w) * 0.55
                    and int(layout_shape.top) < int(self._slide_h) * 0.18
                ):
                    picture_key = (
                        int(layout_shape.left),
                        int(layout_shape.top),
                        int(layout_shape.width),
                        int(layout_shape.height),
                    )
                    if picture_key in copied_header_pictures:
                        continue
                    copied_header_pictures.add(picture_key)
                    slide.shapes.add_picture(
                        io.BytesIO(layout_shape.image.blob),
                        layout_shape.left,
                        layout_shape.top,
                        layout_shape.width,
                        layout_shape.height,
                    )
        # 仅无模板时强制背景色，避免覆盖模板设计
        if not self.template_path:
            set_slide_background(slide, self.theme.background)
        return slide

    def _apply_template_geometry(self, slide, start_index: int, role: str) -> None:
        """把新绘制内容映射到模板识别出的标题/内容/页脚区域。"""
        if not self.template_path or not self._template_mapping:
            return
        role_map = self._template_mapping.get("roles", {}).get(role) or {}
        zones = role_map.get("zones") or {}
        template_title_size = role_map.get("title_font_size_pt")
        header_art_bottom = role_map.get("header_art_bottom")
        header_art_left = role_map.get("header_art_left")
        width, height = int(self._slide_w), int(self._slide_h)
        source_zones = {
            "title": (0.035, 0.01, 0.93, 0.14),
            "content": (0.035, 0.14, 0.93, 0.79),
            "footer": (0.035, 0.93, 0.93, 0.055),
        }
        if role == "matrix":
            # Matrix pages place their insight bullets above the table, earlier
            # than standard chart pages. Include that band in the source content
            # zone so the table keeps its intended distance below the insights.
            source_zones["content"] = (0.035, 0.10, 0.93, 0.83)
        def scale_parts(parts, target_total):
            """Scale table grid lengths while keeping their sum equal to the frame."""
            values = [max(1, int(value)) for value in parts]
            source_total = sum(values)
            if not values or source_total <= 0:
                return values
            scaled = [max(1, round(value / source_total * target_total)) for value in values]
            scaled[-1] += int(target_total) - sum(scaled)
            if scaled[-1] < 1:
                deficit = 1 - scaled[-1]
                scaled[-1] = 1
                for index in range(len(scaled) - 2, -1, -1):
                    available = max(0, scaled[index] - 1)
                    taken = min(available, deficit)
                    scaled[index] -= taken
                    deficit -= taken
                    if deficit == 0:
                        break
            return scaled

        def transform(shape, source, target):
            sx, sy, sw, sh = source
            tx, ty, tw, th = target["x"], target["y"], target["w"], target["h"]
            nx = int(shape.left) / width
            ny = int(shape.top) / height
            nw = int(shape.width) / width
            nh = int(shape.height) / height
            new_left = int((tx + (nx - sx) / max(sw, 0.001) * tw) * width)
            new_top = int((ty + (ny - sy) / max(sh, 0.001) * th) * height)
            new_width = max(1, int(nw / max(sw, 0.001) * tw * width))
            new_height = max(1, int(nh / max(sh, 0.001) * th * height))
            shape.left = new_left
            shape.top = new_top

            # Resizing a table graphic frame does not resize its internal grid.
            # Its columns/rows would therefore keep the pre-template dimensions,
            # while the overlaid charts are mapped correctly. Scale the grid and
            # the frame from the same target geometry so both layers stay aligned.
            if getattr(shape, "has_table", False):
                table = shape.table
                column_widths = scale_parts(
                    [column.width for column in table.columns], new_width
                )
                row_heights = scale_parts(
                    [row.height for row in table.rows], new_height
                )
                for column, column_width in zip(table.columns, column_widths):
                    column.width = column_width
                for row, row_height in zip(table.rows, row_heights):
                    row.height = row_height

            shape.width = new_width
            shape.height = new_height
        new_shapes = list(slide.shapes)[start_index:]
        for shape in new_shapes:
            center_y = (int(shape.top) + int(shape.height) / 2) / height
            shape_text = str(getattr(shape, "text", "") or "").strip()
            is_title_text = (
                bool(shape_text)
                and getattr(shape, "has_text_frame", False)
                and int(shape.top) < height * 0.09
                and int(shape.width) >= width * 0.35
            )
            is_title_divider = (
                not shape_text
                and not getattr(shape, "has_chart", False)
                and not getattr(shape, "has_table", False)
                and int(shape.width) >= width * 0.45
                and abs(int(shape.height)) <= height * 0.01
                and int(shape.top) < height * 0.18
            )
            zone_name = (
                "title" if is_title_text or is_title_divider
                else "footer" if center_y > 0.92
                else "content"
            )
            target = zones.get(zone_name)
            # 极小或越界区域通常是装饰物误识别，继续使用安全的系统布局。
            if not target or target.get("w", 0) < 0.25 or target.get("h", 0) < 0.04:
                continue
            mapped_target = target
            if zone_name == "title":
                # Keep imported title geometry above the body even when an old
                # or external mapping contains an over-tall title rectangle.
                content_zone = zones.get("content") or {}
                safe_bottom = min(
                    0.17,
                    float(content_zone.get("y", 0.18)) - 0.01,
                )
                target_top = max(0.01, min(float(target["y"]), 0.12))
                target_bottom = min(
                    float(target["y"]) + float(target["h"]),
                    safe_bottom,
                )
                if target_bottom - target_top >= 0.04:
                    mapped_target = {
                        **target,
                        "y": target_top,
                        "h": target_bottom - target_top,
                    }
                    if header_art_left and is_title_text:
                        safe_title_right = float(header_art_left) - 0.02
                        safe_title_width = safe_title_right - float(mapped_target["x"])
                        if safe_title_width >= 0.45:
                            mapped_target = {
                                **mapped_target,
                                "w": min(float(mapped_target["w"]), safe_title_width),
                            }
            elif zone_name == "content":
                # Template content unions can include edge decorations and may
                # extend almost to the slide boundary. Keep a print-safe right
                # margin and use the same adjusted zone for every content shape
                # so tables and their chart overlays remain synchronized.
                safe_left, safe_right = 0.025, 0.975
                target_left = max(safe_left, float(target["x"]))
                target_right = min(safe_right, float(target["x"]) + float(target["w"]))
                if target_right - target_left >= 0.25:
                    mapped_target = {
                        **target,
                        "x": target_left,
                        "w": target_right - target_left,
                    }
            transform(shape, source_zones[zone_name], mapped_target)
            if zone_name == "content":
                shape.top = max(int(shape.top), int(float(mapped_target["y"]) * height))
            if (
                zone_name == "title"
                and template_title_size
                and getattr(shape, "has_text_frame", False)
                and str(getattr(shape, "text", "") or "").strip()
            ):
                title_units = sum(2 if ord(char) > 127 else 1 for char in shape.text)
                template_cap = min(28.0, float(template_title_size))
                title_size = (
                    min(template_cap, 16.0) if title_units > 95
                    else min(template_cap, 18.0) if title_units > 70
                    else min(template_cap, 20.0) if title_units > 52
                    else template_cap
                )
                shape.top = int(float(mapped_target["y"]) * height)
                shape.height = int(float(mapped_target["h"]) * height)
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is None or run.font.size.pt > title_size:
                            run.font.size = Pt(title_size)

        # Preserve a small visual gap between the title glyphs and the generated
        # divider while still keeping the line comfortably above the insight box.
        title_text_shapes = [
            shape for shape in new_shapes
            if getattr(shape, "has_text_frame", False)
            and str(getattr(shape, "text", "") or "").strip()
            and int(shape.top) < height * 0.13
            and int(shape.width) >= width * 0.35
        ]
        divider_shapes = [
            shape for shape in new_shapes
            if not str(getattr(shape, "text", "") or "").strip()
            and not getattr(shape, "has_chart", False)
            and not getattr(shape, "has_table", False)
            and int(shape.width) >= width * 0.45
            and abs(int(shape.height)) <= height * 0.01
        ]
        body_text_tops = [
            int(shape.top) for shape in new_shapes
            if getattr(shape, "has_text_frame", False)
            and str(getattr(shape, "text", "") or "").strip()
            and height * 0.12 <= int(shape.top) < height * 0.9
        ]
        if title_text_shapes and divider_shapes:
            title_bottom = max(int(shape.top + shape.height) for shape in title_text_shapes)
            minimum_gap = int(height * 0.012)
            body_top = min(body_text_tops) if body_text_tops else int(height * 0.18)
            latest_divider_top = body_top - int(height * 0.006)
            desired_top = title_bottom + minimum_gap
            if header_art_bottom:
                desired_top = max(
                    desired_top,
                    int((float(header_art_bottom) + 0.015) * height),
                )
            if desired_top <= latest_divider_top:
                for divider in divider_shapes:
                    if int(divider.top) < desired_top:
                        divider.top = desired_top

    def _remove_duplicate_title_divider(self, slide, start_index: int, role: str) -> None:
        """Remove the renderer line when the selected template layout owns one."""
        if not self.template_path or not self._template_mapping:
            return
        role_map = self._template_mapping.get("roles", {}).get(role) or {}
        if not role_map.get("has_title_divider"):
            return
        width, height = int(self._slide_w), int(self._slide_h)
        for shape in list(slide.shapes)[start_index:]:
            try:
                is_generated_divider = (
                    not str(getattr(shape, "text", "") or "").strip()
                    and not getattr(shape, "has_chart", False)
                    and not getattr(shape, "has_table", False)
                    and int(shape.width) >= width * 0.45
                    and abs(int(shape.height)) <= height * 0.01
                    and height * 0.06 <= int(shape.top) <= height * 0.18
                )
            except (AttributeError, TypeError, ValueError):
                is_generated_divider = False
            if is_generated_divider:
                shape._element.getparent().remove(shape._element)

    def _add_cover(self, cover, dims):
        slide = self._blank_slide("cover")
        start = len(slide.shapes)
        build_cover(slide, cover, self.theme, dims)
        if self.template_path and len(slide.shapes) > start:
            # 模板母版/版式已经承担封面装饰，移除系统模板额外添加的顶部色带，
            # 避免两套视觉框架叠加；标题与客户信息继续映射到模板区域。
            first_new = slide.shapes[start]
            if first_new.width > self._slide_w * 0.8 and first_new.height < self._slide_h * 0.05:
                first_new._element.getparent().remove(first_new._element)
        self._apply_template_geometry(slide, start, "cover")

    def _add_toc(self, toc, dims):
        slide = self._blank_slide("toc")
        start = len(slide.shapes)
        build_toc(slide, toc, self.theme, dims)
        self._apply_template_geometry(slide, start, "toc")

    def _add_exec_summary(self, es, dims):
        slide = self._blank_slide("summary")
        start = len(slide.shapes)
        build_exec_summary(slide, es, self.theme, dims)
        self._apply_template_geometry(slide, start, "summary")

    def _add_chart_page(self, page, dims):
        role = "matrix" if isinstance(page, MultiGroupBarPageContent) else "chart"
        slide = self._blank_slide(role)
        start = len(slide.shapes)
        if isinstance(page, MultiGroupBarPageContent):
            build_multi_group_bar_page(slide, page, self.theme, dims)
        else:
            build_chart_page(slide, page, self.theme, dims)
        self._remove_duplicate_title_divider(slide, start, role)
        self._apply_template_geometry(slide, start, role)

    def _add_appendix(self, ap, dims):
        slide = self._blank_slide("appendix")
        start = len(slide.shapes)
        build_appendix(slide, ap, self.theme, dims)
        self._apply_template_geometry(slide, start, "appendix")

    # ------------------------- 模板占位符填充 -------------------------
    def fill_named_placeholders(self, slide, mapping: dict) -> None:
        """在 slide 上按占位符名称填充文本（模板预设占位符用法）。

        适用场景：模板里已放置名为 "title" / "client" / "date" 等占位符，
        你可以用它把内容填进去，而非用本包自绘形状。

        Args:
            slide: 目标幻灯片（通常是基于模板版式添加的）。
            mapping: ``{"占位符名称": "文本内容", ...}``。
        """
        for placeholder in slide.placeholders:
            name = getattr(placeholder, "name", "")
            if name in mapping:
                placeholder.text = str(mapping[name])
