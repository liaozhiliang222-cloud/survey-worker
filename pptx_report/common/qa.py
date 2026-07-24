"""Object-level QA for generated editable PowerPoint files."""

from __future__ import annotations

from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass
class ObjectQAResult:
    slide_count: int = 0
    checked_shapes: int = 0
    issues: list[dict[str, Any]] = field(default_factory=list)
    score: int = 100

    @property
    def ok(self) -> bool:
        return self.error_count == 0

    @property
    def error_count(self) -> int:
        return sum(issue.get("level") == "error" for issue in self.issues)

    @property
    def warning_count(self) -> int:
        return len(self.issues) - self.error_count

    @property
    def slides_with_issues(self) -> list[int]:
        return sorted({int(issue["slide"]) for issue in self.issues if issue.get("slide")})

    def to_dict(self) -> dict[str, Any]:
        return {
            **asdict(self),
            "ok": self.ok,
            "error_count": self.error_count,
            "warning_count": self.warning_count,
            "slides_with_issues": self.slides_with_issues,
        }


def _text_font_sizes(shape) -> list[float]:
    sizes: list[float] = []
    if not getattr(shape, "has_text_frame", False):
        return sizes
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip() and run.font.size is not None:
                sizes.append(float(run.font.size.pt))
    return sizes


def _bbox(shape) -> tuple[int, int, int, int]:
    left, top = int(shape.left), int(shape.top)
    return left, top, left + int(shape.width), top + int(shape.height)


def _overlap(box_a: tuple[int, int, int, int], box_b: tuple[int, int, int, int]) -> tuple[int, int]:
    return (
        max(0, min(box_a[2], box_b[2]) - max(box_a[0], box_b[0])),
        max(0, min(box_a[3], box_b[3]) - max(box_a[1], box_b[1])),
    )


def _is_horizontal_divider(shape, slide_width: int, slide_height: int) -> bool:
    shape_type = getattr(shape, "shape_type", None)
    is_line = shape_type == getattr(MSO_SHAPE_TYPE, "LINE", 9)
    return (
        int(shape.width) >= slide_width * 0.35
        and abs(int(shape.height)) <= slide_height * 0.012
        and (is_line or abs(int(shape.height)) <= slide_height * 0.004)
    )


def _chart_lengths(chart) -> tuple[int, list[int]]:
    try:
        category_count = len(list(chart.plots[0].categories)) if chart.plots else 0
    except (AttributeError, TypeError, ValueError):
        category_count = 0
    lengths: list[int] = []
    for series in chart.series:
        try:
            lengths.append(len(list(series.values)))
        except (AttributeError, TypeError, ValueError):
            lengths.append(0)
    return category_count, lengths


def inspect_presentation(
    path: str | Path,
    *,
    min_font_pt: float = 7.0,
    require_sources: bool = True,
) -> ObjectQAResult:
    prs = Presentation(str(path))
    width, height = int(prs.slide_width), int(prs.slide_height)
    result = ObjectQAResult(slide_count=len(prs.slides))
    for slide_index, slide in enumerate(prs.slides, 1):
        has_source = False
        nonempty_text = 0
        text_shapes = []
        header_text_shapes = []
        header_pictures = []
        divider_shapes = []
        for shape_index, shape in enumerate(slide.shapes, 1):
            result.checked_shapes += 1
            left, top, right, bottom = _bbox(shape)
            if left < 0 or top < 0 or right > width or bottom > height:
                result.issues.append({
                    "level": "error",
                    "code": "shape_out_of_bounds",
                    "slide": slide_index,
                    "shape": shape_index,
                })
            if int(shape.width) <= 0 or (int(shape.height) <= 0 and not _is_horizontal_divider(shape, width, height)):
                result.issues.append({
                    "level": "error",
                    "code": "invalid_shape_size",
                    "slide": slide_index,
                    "shape": shape_index,
                })
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                nonempty_text += 1
                text_shapes.append((shape_index, shape, text))
                if top < height * 0.20 and int(shape.width) > width * 0.25:
                    header_text_shapes.append((shape_index, shape, text))
                    if text.count("\n") >= 2:
                        result.issues.append({
                            "level": "warning",
                            "code": "title_over_two_lines",
                            "slide": slide_index,
                            "shape": shape_index,
                        })
                if "数据来源" in text or "Source" in text:
                    has_source = True
                if "nan" in text.lower():
                    result.issues.append({
                        "level": "error",
                        "code": "nan_text",
                        "slide": slide_index,
                        "shape": shape_index,
                    })
            sizes = _text_font_sizes(shape)
            if sizes and min(sizes) < min_font_pt:
                result.issues.append({
                    "level": "warning",
                    "code": "font_below_threshold",
                    "slide": slide_index,
                    "shape": shape_index,
                    "value": min(sizes),
                })
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and top < height * 0.20:
                header_pictures.append((shape_index, shape))
            if _is_horizontal_divider(shape, width, height):
                divider_shapes.append((shape_index, shape))
            if getattr(shape, "has_chart", False):
                chart = shape.chart
                category_count, series_lengths = _chart_lengths(chart)
                if not chart.series or not series_lengths or all(length == 0 for length in series_lengths):
                    result.issues.append({
                        "level": "error",
                        "code": "empty_chart",
                        "slide": slide_index,
                        "shape": shape_index,
                    })
                elif category_count and any(length != category_count for length in series_lengths):
                    result.issues.append({
                        "level": "error",
                        "code": "chart_data_length_mismatch",
                        "slide": slide_index,
                        "shape": shape_index,
                        "categories": category_count,
                        "series_lengths": series_lengths,
                    })

        for text_index, text_shape, _ in header_text_shapes:
            text_box = _bbox(text_shape)
            for picture_index, picture in header_pictures:
                picture_box = _bbox(picture)
                overlap_width, overlap_height = _overlap(text_box, picture_box)
                if overlap_width * overlap_height > 0:
                    result.issues.append({
                        "level": "error",
                        "code": "template_logo_overlaps_title",
                        "slide": slide_index,
                        "shape": text_index,
                        "related_shape": picture_index,
                    })
                    continue
                vertical_overlap = max(0, min(text_box[3], picture_box[3]) - max(text_box[1], picture_box[1]))
                horizontal_gap = max(text_box[0], picture_box[0]) - min(text_box[2], picture_box[2])
                if vertical_overlap > 0 and 0 <= horizontal_gap < width * 0.012:
                    result.issues.append({
                        "level": "warning",
                        "code": "title_logo_gap_too_small",
                        "slide": slide_index,
                        "shape": text_index,
                        "related_shape": picture_index,
                        "gap": horizontal_gap,
                    })

        for divider_pos, (divider_index, divider) in enumerate(divider_shapes):
            divider_box = _bbox(divider)
            divider_y = (divider_box[1] + divider_box[3]) // 2
            for other_index, other in divider_shapes[divider_pos + 1:]:
                other_box = _bbox(other)
                other_y = (other_box[1] + other_box[3]) // 2
                horizontal_overlap, _ = _overlap(
                    (divider_box[0], 0, divider_box[2], 1),
                    (other_box[0], 0, other_box[2], 1),
                )
                if abs(divider_y - other_y) <= height * 0.008 and horizontal_overlap >= width * 0.30:
                    result.issues.append({
                        "level": "warning",
                        "code": "duplicate_divider_lines",
                        "slide": slide_index,
                        "shape": divider_index,
                        "related_shape": other_index,
                    })
            for text_index, text_shape, text in text_shapes:
                if "数据来源" in text or "Source" in text:
                    continue
                text_box = _bbox(text_shape)
                horizontal_overlap, _ = _overlap(
                    (divider_box[0], 0, divider_box[2], 1),
                    (text_box[0], 0, text_box[2], 1),
                )
                if text_box[1] + height * 0.006 < divider_y < text_box[3] - height * 0.006 and horizontal_overlap >= min(int(text_shape.width), int(divider.width)) * 0.25:
                    result.issues.append({
                        "level": "warning",
                        "code": "divider_intersects_text_box",
                        "slide": slide_index,
                        "shape": divider_index,
                        "related_shape": text_index,
                    })
        if slide_index > 3 and require_sources and nonempty_text and not has_source:
            result.issues.append({
                "level": "warning",
                "code": "source_missing",
                "slide": slide_index,
            })
    result.score = max(0, 100 - result.error_count * 12 - result.warning_count * 2)
    return result
