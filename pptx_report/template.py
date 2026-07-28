"""PPTX 模板分析与主题提取。"""

from __future__ import annotations

import colorsys
import re
import zipfile
from collections import Counter
from statistics import median
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER

from .theme import Theme


DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _norm_rect(shape, width: int, height: int) -> dict:
    return {
        "x": round(max(0, int(shape.left)) / max(1, width), 4),
        "y": round(max(0, int(shape.top)) / max(1, height), 4),
        "w": round(max(0, int(shape.width)) / max(1, width), 4),
        "h": round(max(0, int(shape.height)) / max(1, height), 4),
    }


def _shape_text(shape) -> str:
    try:
        return str(shape.text or "").strip()
    except Exception:  # noqa: BLE001
        return ""


def _is_title_shape(shape, height: int) -> bool:
    try:
        if shape.is_placeholder:
            return shape.placeholder_format.type in {
                PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
            }
    except Exception:  # noqa: BLE001
        pass
    text = _shape_text(shape)
    if not text:
        return False
    top_ratio = int(shape.top) / max(1, height)
    height_ratio = int(shape.height) / max(1, height)
    font_size = _shape_font_size_pt(shape)
    # A loose "top fifth of slide" rule also catches the insight/body box that
    # many research templates place immediately below the title. That makes
    # the inferred title zone span both boxes and pushes the divider into the
    # body. Keep the geometry gate tight, with a small font-size escape hatch
    # for templates whose real title sits slightly lower.
    return (
        top_ratio < 0.12 and height_ratio < 0.15
    ) or (
        top_ratio < 0.16 and height_ratio < 0.13
        and font_size is not None and font_size >= 20
    )


def _shape_font_size_pt(shape) -> float | None:
    sizes = []
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size is not None:
                    sizes.append(float(run.font.size.pt))
    except Exception:  # noqa: BLE001
        return None
    return round(float(median(sizes)), 1) if sizes else None


def _safe_title_zone(shapes, width: int, height: int) -> dict | None:
    zone = _union_rect(shapes, width, height)
    if not zone:
        return None
    # A title area should never consume the body band. This is a final guard
    # for unusual grouped shapes or imported text boxes with misleading bounds.
    max_bottom = 0.17
    bottom = min(zone["y"] + zone["h"], max_bottom)
    zone["h"] = round(min(0.14, max(0.0, bottom - zone["y"])), 4)
    return zone


def _union_rect(shapes, width: int, height: int) -> dict | None:
    shapes = list(shapes)
    if not shapes:
        return None
    left = min(int(shape.left) for shape in shapes)
    top = min(int(shape.top) for shape in shapes)
    right = max(int(shape.left + shape.width) for shape in shapes)
    bottom = max(int(shape.top + shape.height) for shape in shapes)
    class Rect:
        pass
    rect = Rect()
    rect.left, rect.top, rect.width, rect.height = left, top, right - left, bottom - top
    return _norm_rect(rect, width, height)


def _has_title_divider(layout, width: int, height: int) -> bool:
    """Return whether a layout/master already draws a wide line below the title.

    Uploaded report templates often keep their title separator on the layout
    instead of the slide itself. Adding our own separator on top of that
    produces the visually obvious double-line defect.
    """
    for owner in (layout, layout.slide_master):
        for shape in owner.shapes:
            try:
                left = int(shape.left)
                top = int(shape.top)
                shape_width = int(shape.width)
                shape_height = abs(int(shape.height))
            except (AttributeError, TypeError, ValueError):
                continue
            if (
                shape_width >= width * 0.45
                and shape_height <= height * 0.035
                and height * 0.06 <= top <= height * 0.22
                and left <= width * 0.35
            ):
                return True
    return False


def _header_art_bottom(layout, width: int, height: int) -> float | None:
    """Return the bottom edge of visible header artwork such as a logo."""
    bottoms = []
    for owner in (layout, layout.slide_master):
        for shape in owner.shapes:
            try:
                if shape.is_placeholder:
                    continue
            except Exception:  # noqa: BLE001
                pass
            try:
                top = int(shape.top)
                bottom = int(shape.top + shape.height)
                shape_height = abs(int(shape.height))
            except (AttributeError, TypeError, ValueError):
                continue
            if (
                top < height * 0.18
                and 0 < shape_height < height * 0.22
                and bottom < height * 0.24
            ):
                bottoms.append(bottom / max(1, height))
    return round(max(bottoms), 4) if bottoms else None


def _header_art_left(layout, width: int, height: int) -> float | None:
    """Return the left edge of visible right-side header artwork."""
    lefts = []
    for owner in (layout, layout.slide_master):
        for shape in owner.shapes:
            try:
                if shape.is_placeholder:
                    continue
            except Exception:  # noqa: BLE001
                pass
            try:
                left = int(shape.left)
                top = int(shape.top)
                shape_height = abs(int(shape.height))
            except (AttributeError, TypeError, ValueError):
                continue
            if left > width * 0.55 and top < height * 0.18 and 0 < shape_height < height * 0.22:
                lefts.append(left / max(1, width))
    return round(min(lefts), 4) if lefts else None


def _classify_slide(slide, index: int) -> tuple[str, float]:
    texts = " ".join(_shape_text(shape) for shape in slide.shapes if _shape_text(shape)).lower()
    charts = sum(1 for shape in slide.shapes if getattr(shape, "has_chart", False))
    tables = sum(1 for shape in slide.shapes if getattr(shape, "has_table", False))
    pictures = sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
    if re.search(r"目录|contents?", texts):
        return "toc", 0.98
    if re.search(r"执行摘要|核心摘要|executive\s+summary", texts):
        return "summary", 0.98
    if re.search(r"附录|appendix", texts):
        return "appendix", 0.96
    if tables or charts >= 4:
        return "matrix", 0.9
    if charts:
        return "chart", 0.92
    if index == 0 and len(texts) < 260:
        return "cover", 0.86
    if re.search(r"chapter|part\s*\d|章节|篇章", texts) and len(texts) < 180:
        return "section", 0.82
    if pictures and len(texts) < 120:
        return "section", 0.72
    return "content", 0.62


def build_template_mapping(prs: Presentation, role_overrides: dict | None = None) -> dict:
    """识别示例页角色和版面区域，供渲染器按语义选版式与落位。"""
    width, height = int(prs.slide_width), int(prs.slide_height)
    candidates: dict[str, list[dict]] = {}
    all_candidates: list[dict] = []
    layout_indexes = {id(layout): idx for idx, layout in enumerate(prs.slide_layouts)}
    for index, slide in enumerate(prs.slides):
        role, confidence = _classify_slide(slide, index)
        title_shapes = [shape for shape in slide.shapes if _is_title_shape(shape, height)]
        footer_shapes = [
            shape for shape in slide.shapes
            if shape.top > height * 0.88 and (_shape_text(shape) or getattr(shape, "is_placeholder", False))
        ]
        content_shapes = [
            shape for shape in slide.shapes
            if shape not in title_shapes and shape not in footer_shapes
            and (
                bool(_shape_text(shape))
                or getattr(shape, "has_chart", False)
                or getattr(shape, "has_table", False)
                or shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                or getattr(shape, "is_placeholder", False)
            )
            and shape.top + shape.height > height * 0.18
            and shape.top < height * 0.9
        ]
        layout_index = layout_indexes.get(id(slide.slide_layout), 0)
        candidate = {
            "role": role,
            "slide_index": index,
            "layout_index": layout_index,
            "layout_name": slide.slide_layout.name or "未命名版式",
            "confidence": round(confidence, 2),
            "has_title_divider": _has_title_divider(slide.slide_layout, width, height),
            "header_art_bottom": _header_art_bottom(slide.slide_layout, width, height),
            "header_art_left": _header_art_left(slide.slide_layout, width, height),
            "title_font_size_pt": _shape_font_size_pt(title_shapes[0]) if title_shapes else None,
            "zones": {
                "title": _safe_title_zone(title_shapes, width, height),
                "content": _union_rect(content_shapes, width, height),
                "footer": _union_rect(footer_shapes, width, height),
            },
        }
        candidates.setdefault(role, []).append(candidate)
        all_candidates.append(candidate)
    roles = {}
    for role, items in candidates.items():
        items.sort(key=lambda item: (item["confidence"], bool(item["zones"]["content"])), reverse=True)
        roles[role] = items[0]
    # 通用正文映射可由图表页或普通内容页回退，确保所有报告页面都有可用落位。
    if role_overrides:
        by_slide = {item["slide_index"]: item for item in all_candidates}
        for role, requested in role_overrides.items():
            try:
                slide_index = int(requested) - 1
            except (TypeError, ValueError):
                continue
            selected = by_slide.get(slide_index)
            if selected is not None:
                roles[str(role)] = {
                    **selected,
                    "role": str(role),
                    "confidence": 1.0,
                    "user_confirmed": True,
                }
    fallback = roles.get("chart") or roles.get("content") or next(iter(roles.values()), None)
    for role in ("cover", "toc", "summary", "chart", "matrix", "appendix"):
        role_item = roles.get(role)
        zones = (role_item or {}).get("zones") or {}
        title_zone = zones.get("title") or {}
        content_zone = zones.get("content") or {}
        data_layout_is_unsafe = role in {"chart", "matrix"} and (
            not title_zone
            or content_zone.get("w", 0) < 0.60
            or content_zone.get("h", 0) < 0.45
        )
        if (
            role_item is None
            or (data_layout_is_unsafe and not role_item.get("user_confirmed"))
        ) and fallback:
            roles[role] = {**fallback, "role": role, "fallback": True,
                           "confidence": round(max(0.35, fallback["confidence"] - 0.2), 2)}
    confidences = [item["confidence"] for item in roles.values()]
    return {
        "version": 2,
        "mode": "semantic-layout-zones",
        "roles": roles,
        "coverage": round(min(1.0, len(candidates) / 6), 2),
        "confidence": round(sum(confidences) / max(1, len(confidences)), 2),
    }


def _structure_text_items(slide, width: int, height: int) -> list[dict]:
    """Return positioned text fragments used by the report-structure detector."""
    items = []
    for shape in slide.shapes:
        text = _shape_text(shape)
        if not text:
            continue
        try:
            top = int(shape.top) / max(1, height)
            left = int(shape.left) / max(1, width)
        except (AttributeError, TypeError, ValueError):
            top, left = 0.5, 0.5
        size = _shape_font_size_pt(shape) or 0.0
        for line in (part.strip() for part in re.split(r"[\r\n]+", text)):
            if line:
                items.append({"text": line, "top": top, "left": left, "font_size": size})
    return items


def _structure_title(items: list[dict], number_text: str) -> str:
    ignored = {"目录", "contents", "谢谢", "thank you", "感谢聆听"}
    candidates = []
    for item in items:
        text = item["text"].strip()
        if (
            text == number_text
            or text.lower() in ignored
            or "|" in text
            or "｜" in text
            or re.fullmatch(r"\d+(?:\.\d+)?", text)
            or item["top"] > 0.84
            or len(text) > 70
        ):
            continue
        score = float(item["font_size"]) * 10 - item["top"] * 3 - item["left"]
        candidates.append((score, text))
    return max(candidates, default=(0, ""))[1][:80]


def _structure_topics(items: list[dict]) -> list[str]:
    topic_lines = [item["text"] for item in items if "|" in item["text"] or "｜" in item["text"]]
    if not topic_lines:
        return []
    longest = max(topic_lines, key=len)
    return [part.strip()[:80] for part in re.split(r"[|｜]", longest) if part.strip()][:12]


def detect_report_structure(prs: Presentation) -> dict:
    """Detect numbered section and subsection divider slides."""
    width, height = int(prs.slide_width), int(prs.slide_height)
    top_sections = []
    subsections = []
    for slide_index, slide in enumerate(prs.slides):
        items = _structure_text_items(slide, width, height)
        if not items:
            continue
        top_numbers = [item for item in items if re.fullmatch(r"0?[1-9]|[1-9]\d", item["text"])
                       and item["top"] < 0.78
                       and (item["text"].startswith("0") or item["font_size"] >= 24 or item["left"] < 0.2)]
        sub_numbers = [item for item in items if re.fullmatch(r"[1-9]\d*\.[1-9]\d*", item["text"])
                       and item["font_size"] >= 10 and item["top"] < 0.82]
        if top_numbers:
            number = max(top_numbers, key=lambda item: item["font_size"])["text"]
            title = _structure_title(items, number)
            if title:
                top_sections.append({
                    "section_id": f"section_{int(number):02d}", "number": number.zfill(2),
                    "title": title, "topics": _structure_topics(items),
                    "start_slide_index": slide_index, "end_slide_index": slide_index,
                    "subsections": [],
                })
                continue
        if sub_numbers:
            number = max(sub_numbers, key=lambda item: item["font_size"])["text"]
            title = _structure_title(items, number)
            if title:
                subsections.append({
                    "subsection_id": f"subsection_{number.replace('.', '_')}", "number": number,
                    "title": title, "topics": _structure_topics(items),
                    "start_slide_index": slide_index, "end_slide_index": slide_index,
                })

    top_sections.sort(key=lambda item: item["start_slide_index"])
    if len(top_sections) < 2:
        return {"version": 1, "source": "template", "confidence": 0.0,
                "sections": [], "detected_slide_indices": []}
    for index, section in enumerate(top_sections):
        next_start = top_sections[index + 1]["start_slide_index"] if index + 1 < len(top_sections) else len(prs.slides)
        section["end_slide_index"] = next_start - 1
        children = [item for item in subsections
                    if section["start_slide_index"] < item["start_slide_index"] < next_start]
        for child_index, child in enumerate(children):
            child["end_slide_index"] = (children[child_index + 1]["start_slide_index"] - 1
                                        if child_index + 1 < len(children) else next_start - 1)
        section["subsections"] = children

    numbered_sequence = all(int(item["number"]) == int(top_sections[0]["number"]) + index
                            for index, item in enumerate(top_sections))
    confidence = 0.72 + len(top_sections) * 0.05 + len(subsections) * 0.025
    if numbered_sequence:
        confidence += 0.05
    detected_indices = [item["start_slide_index"] for item in top_sections]
    detected_indices.extend(item["start_slide_index"] for item in subsections)
    return {
        "version": 1, "source": "template", "confidence": round(min(0.99, confidence), 2),
        "sections": top_sections[:8], "detected_slide_indices": sorted(set(detected_indices)),
    }


def _theme_xml(path: str) -> bytes | None:
    with zipfile.ZipFile(path) as archive:
        members = archive.infolist()
        if len(members) > 5000 or sum(item.file_size for item in members) > 200 * 1024 * 1024:
            raise ValueError("模板解压后体积过大或内部文件过多")
        names = [name for name in archive.namelist() if re.fullmatch(r"ppt/theme/theme\d+\.xml", name)]
        return archive.read(names[0]) if names else None


def _extract_theme_tokens(path: str) -> tuple[list[str], list[str]]:
    raw = _theme_xml(path)
    if not raw:
        return [], []
    root = ET.fromstring(raw)
    ns = {"a": DRAWING_NS}
    colors = []
    scheme = root.find(".//a:themeElements/a:clrScheme", ns)
    if scheme is not None:
        for child in list(scheme):
            color = child.find("a:srgbClr", ns)
            if color is not None and color.get("val"):
                value = color.get("val").upper()
                if value not in {"FFFFFF", "000000"} and value not in colors:
                    colors.append(value)
    fonts = []
    for path_expr in (
        ".//a:themeElements/a:fontScheme/a:majorFont/a:ea",
        ".//a:themeElements/a:fontScheme/a:minorFont/a:ea",
        ".//a:themeElements/a:fontScheme/a:majorFont/a:latin",
        ".//a:themeElements/a:fontScheme/a:minorFont/a:latin",
    ):
        node = root.find(path_expr, ns)
        value = (node.get("typeface") if node is not None else "") or ""
        if value and value not in fonts:
            fonts.append(value)
    return colors, fonts


def analyze_template(path: str, filename: str = "template.pptx") -> dict:
    """读取模板基础信息，供前端确认和生成器换肤。"""
    prs = Presentation(path)
    colors, fonts = _extract_theme_tokens(path)
    width = int(prs.slide_width)
    height = int(prs.slide_height)
    ratio = width / height if height else 0
    layout_usage = Counter(slide.slide_layout.name or "未命名版式" for slide in prs.slides)
    layouts = []
    for layout in prs.slide_layouts:
        placeholders = sum(1 for shape in layout.shapes if getattr(shape, "is_placeholder", False))
        layouts.append({
            "name": layout.name or "未命名版式",
            "placeholders": placeholders,
            "used_by_slides": layout_usage.get(layout.name or "未命名版式", 0),
        })
    warnings = []
    mapping = build_template_mapping(prs)
    report_structure = detect_report_structure(prs)
    if not (1.70 <= ratio <= 1.82):
        warnings.append("模板不是常见的 16:9 宽屏比例，生成内容会按模板原始比例缩放布局。")
    if not any((layout.name or "").lower() in {"blank", "空白"} for layout in prs.slide_layouts):
        warnings.append("未识别到明确的空白版式，将使用最接近的可用版式并清除占位符。")
    if len(prs.slides) > 50:
        warnings.append("模板示例页较多；生成时会移除示例页，仅保留母版、版式和主题资源。")
    if mapping["confidence"] < 0.65:
        warnings.append("模板中的示例页语义不够明确，部分页面将回退到最接近的正文版式。")
    missing_roles = [role for role, item in mapping["roles"].items() if item.get("fallback")]
    if missing_roles:
        warnings.append("未找到全部页面类型，已为缺失类型配置安全回退版式。")
    return {
        "file_name": filename,
        "slide_count": len(prs.slides),
        "layout_count": len(prs.slide_layouts),
        "aspect_ratio": round(ratio, 3),
        "is_widescreen": 1.70 <= ratio <= 1.82,
        "theme_colors": colors[:8],
        "fonts": fonts[:6],
        "layouts": layouts[:20],
        "mapping": mapping,
        "report_structure": report_structure,
        "warnings": warnings,
    }


def theme_from_template(path: str, fallback: Theme) -> Theme:
    """在保留报告可读性规范的前提下，套用模板主色与字体。"""
    colors, fonts = _extract_theme_tokens(path)
    def usable_accent(value: str) -> bool:
        red, green, blue = (int(value[i:i + 2], 16) / 255 for i in (0, 2, 4))
        _, lightness, saturation = colorsys.rgb_to_hls(red, green, blue)
        return lightness < 0.82 and saturation > 0.10

    usable = [color for color in colors if usable_accent(color)]
    palette = usable[:4] if usable else list(fallback.palette)
    while len(palette) < 4:
        palette.append(fallback.palette[len(palette) % len(fallback.palette)])
    segment_palette = list(usable)
    segment_palette.extend(color for color in fallback.segment_palette if color not in segment_palette)
    return Theme(
        name="上传模板",
        font_name=fonts[0] if fonts else fallback.font_name,
        palette=palette[:4],
        segment_palette=segment_palette,
        background=fallback.background,
        text_dark=fallback.text_dark,
        text_light=fallback.text_light,
        data_label_color=fallback.data_label_color,
        data_label_size=fallback.data_label_size,
        pct_format=fallback.pct_format,
    )
