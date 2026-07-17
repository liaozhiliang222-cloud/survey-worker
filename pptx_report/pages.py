"""页面构建模块。

每个函数负责「一页」的绘制，输入是页面内容数据 + 主题 + 画布尺寸，
输出是把形状画到 slide 上。所有页面共用 :mod:`pptx_report.utils` 的字体 / 填充助手，
保证风格统一、中文字体正确。
"""

from __future__ import annotations

import os
import re
from typing import Tuple

import pandas as pd
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from lxml import etree

from . import charts, layouts
from .model import (
    AppendixContent,
    ChartPageContent,
    CoverContent,
    ExecutiveSummaryContent,
    MultiGroupBarPageContent,
    TocContent,
)
from .theme import Theme
from .utils import (
    remove_shape_outline,
    set_shape_fill,
    set_slide_background,
    style_font,
    style_textframe,
)

# 复用布局模块的几何常量，保证标题 / 内容区对齐一致
from .layouts import (
    BOTTOM_MARGIN,
    CAPTION_H,
    CONTENT_TOP,
    PAGE_MARGIN,
    TITLE_HEIGHT,
    TITLE_TOP,
)

Dims = Tuple[Emu, Emu]

# DrawingML 主命名空间（表格单元格底层 XML 填充/边框用）
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _text_width_units(text: str) -> float:
    """估算文本显示宽度（单位：以「一个中文字」为 1.0）。

    中文/全角字符按 1.0，ASCII/半角字符按 0.55，用于按列宽做等宽截断。
    """
    w = 0.0
    for ch in str(text):
        w += 1.0 if ord(ch) > 0x2E7F else 0.55
    return w


def _truncate_to_width(text: str, max_units: float) -> str:
    """把文本按显示宽度截断到 ``max_units``（以中文字为单位），超出加省略号。"""
    text = str(text)
    if _text_width_units(text) <= max_units:
        return text
    budget = max_units - 1.0  # 给省略号留位置
    out, acc = [], 0.0
    for ch in text:
        cw = 1.0 if ord(ch) > 0x2E7F else 0.55
        if acc + cw > budget:
            break
        out.append(ch)
        acc += cw
    return "".join(out).rstrip() + "…"


# ------------------------- 通用助手 -------------------------
def _add_textbox(slide, text, x, y, cx, cy, theme: Theme, size=18, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 line_spacing=None):
    """在 slide 上加一个文本框并写入（单行）文本。"""
    tb = slide.shapes.add_textbox(x, y, cx, cy)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    style_font(run.font, theme, size=size, bold=bold, color=color or theme.text_dark)
    return tb


def _add_bullets(slide, items, x, y, cx, cy, theme: Theme, size=13, color=None):
    """在 slide 上添加一个项目符号列表（支持多段）。"""
    tb = slide.shapes.add_textbox(x, y, cx, cy)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "• " + str(item)
        style_font(run.font, theme, size=size, color=color or theme.text_dark)
    return tb


def _add_numbered_insights(slide, items, x, y, cx, cy, theme: Theme,
                           size=12.5, color=None):
    """把洞察拆成清晰的编号段落，避免长句挤成一整行。"""
    tb = slide.shapes.add_textbox(x, y, cx, cy)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    for i, item in enumerate(items, 1):
        clean = re.sub(r"^\s*(?:[•·●]|\d+[\.、）)])\s*", "", str(item)).strip()
        if not clean:
            continue
        p = tf.paragraphs[0] if len(tf.paragraphs) == 1 and not tf.paragraphs[0].text else tf.add_paragraph()
        p.space_after = Pt(3)
        p.line_spacing = 1.05
        run = p.add_run()
        run.text = f"{i}. {clean}"
        style_font(run.font, theme, size=size, color=color or theme.text_dark)
    return tb


# ------------------------- 封面 -------------------------
def build_cover(slide, cover: CoverContent, theme: Theme, dims: Dims) -> None:
    slide_w, _ = dims
    # 顶部主色带
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, Inches(0.18))
    set_shape_fill(band, theme.primary)
    remove_shape_outline(band)

    _add_textbox(slide, cover.title, Inches(PAGE_MARGIN), Inches(2.1),
                  slide_w - Inches(2 * PAGE_MARGIN), Inches(1.6), theme,
                  size=40, bold=True, color=theme.primary)
    if cover.subtitle:
        _add_textbox(slide, cover.subtitle, Inches(PAGE_MARGIN), Inches(3.6),
                      slide_w - Inches(2 * PAGE_MARGIN), Inches(0.6), theme,
                      size=18, color=theme.text_dark)
    _add_textbox(slide, f"客户 / 品牌：{cover.client}", Inches(PAGE_MARGIN), Inches(4.7),
                  slide_w - Inches(2 * PAGE_MARGIN), Inches(0.5), theme,
                  size=16, bold=True, color=theme.secondary)
    _add_textbox(slide, f"报告日期：{cover.date}", Inches(PAGE_MARGIN), Inches(5.2),
                  slide_w - Inches(2 * PAGE_MARGIN), Inches(0.5), theme,
                  size=16, color=theme.secondary)

    # Logo：仅当文件存在时绘制（预留位置）
    if cover.logo_path and os.path.exists(cover.logo_path):
        try:
            slide.shapes.add_picture(cover.logo_path,
                                     slide_w - Inches(2.4), Inches(0.5),
                                     height=Inches(1.2))
        except Exception:
            pass


# ------------------------- 目录 -------------------------
def build_toc(slide, toc: TocContent, theme: Theme, dims: Dims) -> None:
    slide_w, slide_h = dims
    _add_textbox(slide, "目录 CONTENTS", Inches(PAGE_MARGIN), Inches(TITLE_TOP),
                  slide_w - Inches(2 * PAGE_MARGIN), Inches(TITLE_HEIGHT), theme,
                  size=32, bold=True, color=theme.primary)
    # 章节较多时压缩纵向节奏，确保最后一项完整落在安全区内。
    section_count = max(1, len(toc.sections))
    y = 1.55 if section_count >= 7 else 1.8
    step = min(0.85, (float(slide_h) / 914400.0 - y - BOTTOM_MARGIN - 0.18) / section_count)
    row_h = min(0.6, max(0.45, step - 0.12))
    for i, sec in enumerate(toc.sections):
        num = f"{i + 1:02d}"
        _add_textbox(slide, num, Inches(PAGE_MARGIN), Inches(y),
                      Inches(1.0), Inches(row_h), theme, size=22 if section_count >= 7 else 24, bold=True,
                      color=theme.accent)
        _add_textbox(slide, sec, Inches(PAGE_MARGIN + 1.1), Inches(y),
                      slide_w - Inches(2 * PAGE_MARGIN + 1.1), Inches(row_h), theme,
                      size=19 if section_count >= 7 else 20, color=theme.text_dark)
        # 底部分隔线
        line_y = y + row_h + 0.02
        if line_y < slide_h - BOTTOM_MARGIN:
            divider = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(PAGE_MARGIN), Inches(line_y),
                slide_w - Inches(2 * PAGE_MARGIN), Inches(0.012))
            set_shape_fill(divider, "D9D9D9")
            remove_shape_outline(divider)
        y += step


# ------------------------- 执行摘要 -------------------------
def build_exec_summary(slide, es, theme: Theme, dims: Dims) -> None:
    slide_w, slide_h = dims
    _add_textbox(slide, "执行摘要 EXECUTIVE SUMMARY", Inches(PAGE_MARGIN),
                  Inches(TITLE_TOP), slide_w - Inches(2 * PAGE_MARGIN),
                  Inches(TITLE_HEIGHT), theme, size=32, bold=True, color=theme.primary)

    n = len(es.kpis)
    gap = 0.3
    card_w = (float(slide_w) / 914400.0 - 2 * PAGE_MARGIN - gap * (n - 1)) / n
    card_h = 2.2
    y = 1.8
    for i, kpi in enumerate(es.kpis):
        x = Inches(PAGE_MARGIN + i * (card_w + gap))
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(y), Inches(card_w), Inches(card_h))
        card_color = theme.primary if i == 0 else theme.color(i)
        set_shape_fill(card, card_color)
        remove_shape_outline(card)

        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = kpi.value
        style_font(r.font, theme, size=40, bold=True, color=theme.text_light)

        if kpi.delta:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = kpi.delta
            style_font(r2.font, theme, size=12, color=theme.text_light)

        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = kpi.label
        style_font(r3.font, theme, size=14, color=theme.text_light)

    # 结论
    _add_textbox(slide, es.conclusion, Inches(PAGE_MARGIN),
                  Inches(y + card_h + 0.35),
                  slide_w - Inches(2 * PAGE_MARGIN),
                  slide_h - Inches(y + card_h + 0.35 + BOTTOM_MARGIN),
                  theme, size=16, color=theme.text_dark, line_spacing=1.3)


# ------------------------- 图表分析页 -------------------------
def build_chart_page(slide, page: ChartPageContent, theme: Theme, dims: Dims) -> None:
    slide_w, slide_h = dims
    title_units = _text_width_units(page.title)
    title_size = 28 if title_units <= 48 else 24 if title_units <= 72 else 21
    _add_textbox(slide, page.title, Inches(PAGE_MARGIN), Inches(TITLE_TOP),
                  slide_w - Inches(2 * PAGE_MARGIN), Inches(TITLE_HEIGHT), theme,
                  size=title_size, bold=True, color=theme.primary)
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(PAGE_MARGIN), Inches(0.98),
        slide_w - Inches(2 * PAGE_MARGIN), Inches(0.018),
    )
    set_shape_fill(divider, theme.primary)
    divider.line.fill.background()
    if page.subtitle:
        _add_textbox(slide, page.subtitle, Inches(PAGE_MARGIN), Inches(TITLE_TOP + 0.75),
                      slide_w - Inches(2 * PAGE_MARGIN), Inches(0.45), theme,
                      size=14, color=theme.text_dark)

    layout = layouts.resolve_layout(page, slide_w, slide_h)

    # 图文混排：右侧洞察栏
    if layout.sidebar_rect is not None:
        sr = layout.sidebar_rect
        _add_textbox(slide, "洞察 INSIGHT", sr.x, sr.y, sr.cx, Inches(0.5), theme,
                      size=16, bold=True, color=theme.primary)
        _add_bullets(slide, page.side_insights, sr.x, sr.y + Inches(0.6),
                      sr.cx, sr.cy - Inches(0.6), theme, size=13)

    # 标题下方直接呈现洞察正文，不再增加“核心洞察/核心结论”标签。
    insight_texts = list(dict.fromkeys(
        text for text in (
            list(page.insights or [])
            or [c.insight for c in page.charts]
        )
        if text and text not in page.title and page.title not in text
    ))
    if insight_texts:
        max_units = max(_text_width_units(text) for text in insight_texts)
        insight_h = min(0.92, max(0.34, 0.25 * len(insight_texts) + (0.18 if max_units > 70 else 0)))
        insight_y = Inches(1.05)
        _add_numbered_insights(
            slide, insight_texts, Inches(PAGE_MARGIN), insight_y,
            slide_w - Inches(2 * PAGE_MARGIN), Inches(insight_h), theme,
            size=12.5, color=theme.text_dark,
        )
    else:
        insight_h = 0

    # 数据来源统一固定在页面最底部。
    if page.data_source:
        slide_h_in = float(slide_h) / 914400.0
        src_y_emu = Inches(slide_h_in - 0.38)
        _add_textbox(slide, page.data_source, Inches(PAGE_MARGIN), src_y_emu,
                      slide_w - Inches(2 * PAGE_MARGIN), Inches(0.2), theme,
                      size=8, color="758D99")

    # 图表（如有洞察则下移起始位置）
    for slot, spec in zip(layout.slots, page.charts):
        chart_h = slot.cy - Inches(CAPTION_H)
        if insight_h > 0:
            # 把 slot 整体下移洞察高度
            slot_top = slot.y + Emu(int(insight_h * 914400))
            slot_cy = chart_h - Emu(int(insight_h * 914400))
            if slot_cy > Emu(0):
                charts.add_chart(slide, spec, slot.x, slot_top, slot.cx, slot_cy, theme)
            else:
                charts.add_chart(slide, spec, slot.x, slot.y, slot.cx, chart_h, theme)
        else:
            charts.add_chart(slide, spec, slot.x, slot.y, slot.cx, chart_h, theme)


# ------------------------- 多组多列条形图页（表格+图表叠加） -------------------------
def build_multi_group_bar_page(slide, page: MultiGroupBarPageContent,
                                theme: Theme, dims: Dims) -> None:
    """渲染 multi_group_bar 布局：底层 PPT 表格 + 上层每人群一个 BAR_CLUSTERED 图表。

    布局结构（参考 crosstab-to-ppt skill 的 build_multi_group_bar_slide）：
      - 底层：全宽表格（组 | 选项 | 各人群列），人群列留空
      - 上层：每个人群一个独立横向条形图，覆盖在对应列位置
      - 隐藏图表坐标轴，数据标签显示百分比
      - 组标题合并单元格 + 组间分隔线
    """
    slide_w, slide_h = dims

    # ── 标题栏 ──
    _add_textbox(slide, page.title, Inches(0.4), Inches(0.10),
                  Inches(12.5), Inches(0.58), theme,
                  size=20, bold=True, color=theme.primary)
    # 标题下方分隔线
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.72),
                                     Inches(12.5), Inches(0.02))
    set_shape_fill(divider, theme.primary)
    remove_shape_outline(divider)

    # 标题下方直接写洞察正文，不显示“分析结论”标签或占位话术。
    insights = [
        text for text in (page.insights or [])
        if text and text not in page.title and page.title not in text
    ]
    insight_y = 0.80
    if insights:
        insights = list(dict.fromkeys(insights))
        max_units = max(_text_width_units(text) for text in insights)
        region_h = min(0.92, max(0.34, 0.25 * len(insights) + (0.18 if max_units > 70 else 0)))
        _add_numbered_insights(
            slide, insights, Inches(0.4), Inches(insight_y),
            Inches(12.5), Inches(region_h), theme,
            size=12.5, color=theme.text_dark,
        )
    else:
        region_h = 0
    table_y = insight_y + region_h + 0.08

    segments = page.segments or []
    n_segs = len(segments)
    if n_segs == 0:
        return

    groups_data = page.groups_data or []
    colors = [theme.seg_color(i) for i in range(n_segs)]

    # ── 合并所有组的选项行（按组内排序） ──
    merged_rows = []          # list[dict]: {"选项": str, seg1: float, ...}
    group_boundaries = []     # [(start_idx, end_idx, title), ...]

    sort_col = 'Total' if 'Total' in segments else (segments[0] if segments else None)

    # v10: "其他"类选项关键词——始终置底，不参与降序排序
    _OTHER_PATTERNS = ["其他", "其它", "请说明", "other", "T2B", "B2B",
                       "无", "不适用", "跳过", "拒答"]

    def _is_other_option(opt: str) -> bool:
        opt_s = str(opt).strip().lower()
        return any(p.lower() in opt_s for p in _OTHER_PATTERNS)

    for g in groups_data:
        g_title = g.get("title", "")       # 完整题干（不再用短标签）
        g_df = g.get("data")
        if g_df is None or g_df.empty:
            continue
        g_df = g_df.copy()
        start_idx = len(merged_rows)

        cats = [str(v) for v in g_df["选项"].tolist()]
        should_sort = not _has_natural_order(cats)

        # v10: 拆分为"常规选项"和"其他/特殊选项"
        regular_mask = g_df["选项"].apply(lambda x: not _is_other_option(str(x)))
        regular_df = g_df[regular_mask].copy() if regular_mask.any() else pd.DataFrame(columns=g_df.columns)
        other_df = g_df[~regular_mask].copy() if (~regular_mask).any() else pd.DataFrame(columns=g_df.columns)

        # 仅对常规选项排序；"其他"项保持原始顺序追加到底部
        if should_sort and sort_col and sort_col in regular_df.columns and not regular_df.empty:
            regular_df = regular_df.sort_values(by=sort_col, ascending=False).reset_index(drop=True)

        # 合并：常规选项在前 + 其他/过滤项在后
        sorted_df = pd.concat([regular_df, other_df], ignore_index=True)

        # v10: 过滤掉 T2B / B2B 内部代码行（不渲染到报告中）
        filter_mask = sorted_df["选项"].apply(
            lambda x: not any(p.upper() in str(x).strip().upper() for p in ("T2B", "B2B"))
        )
        sorted_df = sorted_df[filter_mask].reset_index(drop=True)

        for _, row in sorted_df.iterrows():
            r = {"选项": str(row.get("选项", ""))}
            for seg in segments:
                v = row.get(seg, 0)
                if pd.isna(v):
                    v = 0.0
                if isinstance(v, (int, float)) and abs(v) > 1:
                    v = v / 100.0
                r[seg] = float(v)
            merged_rows.append(r)

        end_idx = len(merged_rows) - 1
        group_boundaries.append((start_idx, end_idx, g_title))

    n_total_rows = len(merged_rows)
    if n_total_rows == 0:
        return

    # ── 行数上限（防溢出） ──
    MAX_OPTS = 20
    if n_total_rows > MAX_OPTS:
        merged_rows = merged_rows[:MAX_OPTS]
        new_bounds = []
        for gs, ge, gt in group_boundaries:
            if gs >= MAX_OPTS:
                break
            new_ge = min(ge, MAX_OPTS - 1)
            new_bounds.append((gs, new_ge, gt))
            if new_ge < ge:
                break
        group_boundaries = new_bounds
        n_total_rows = len(merged_rows)

    # ── 列宽：簇状图(横向条形图)模式下组列+选项列固定 4.5cm（用户指定 v11） ──
    margin_left, margin_right = 0.25, 0.25
    bottom_limit = 7.5 - 0.35   # 给底部数据来源留空间
    table_w = 13.333 - margin_left - margin_right
    CM_TO_IN = 1 / 2.54
    LABEL_COL_CM = 4.5
    label_total_w = LABEL_COL_CM * 2 * CM_TO_IN   # 组+选项共 9cm
    group_col_w = option_col_w = LABEL_COL_CM * CM_TO_IN  # 各 4.5cm

    dark_gray = RGBColor(0x40, 0x40, 0x40)
    header_fill = RGBColor(0xE8, 0xE8, 0xE8)

    # ── 选项/组文本截断参数（供簇状表使用） ──
    n_table_rows_est = n_total_rows + 1
    table_h_est = bottom_limit - table_y
    table_row_h_est = table_h_est / n_table_rows_est
    if table_row_h_est < 0.22:
        opt_font = 6
    elif table_row_h_est < 0.28:
        opt_font = 7
    else:
        opt_font = 8
    line_h_in = opt_font / 72.0 * 1.25
    avail_h = table_row_h_est - 0.06
    max_lines = max(1, min(2, int(avail_h / line_h_in)))
    usable_w = max(0.8, option_col_w - 0.18)
    chars_per_line = usable_w / (opt_font / 72.0)
    opt_max_units = max_lines * chars_per_line * 0.98
    grp_max_lines = max(2, min(3, int(avail_h / line_h_in)))
    grp_usable_w = max(1.0, group_col_w - 0.18)
    grp_chars_per_line = grp_usable_w / (opt_font / 72.0)
    grp_max_units = grp_max_lines * grp_chars_per_line * 0.96

    # ═══════════════════ 判断图表模式 ═══════════════════
    all_options = list(dict.fromkeys(r["选项"] for r in merged_rows))  # 去重保序
    n_options = len(all_options)
    USE_STACKED_THRESHOLD = 7   # 选项数阈值：≤此值用堆积图

    if n_options <= USE_STACKED_THRESHOLD:
        # ── 100% 堆积图：控制在页面中部，不再铺满整页 ──
        available_h = bottom_limit - table_y
        stack_w = min(table_w, 10.2)
        stack_h = min(4.15, max(2.8, available_h - 0.35))
        stack_x = margin_left + (table_w - stack_w) / 2
        chart_y = table_y + max(0.12, (available_h - stack_h) * 0.35)
        _render_stacked_bars(
            slide, merged_rows, segments, all_options,
            stack_x, stack_w, chart_y, stack_h,
            colors, theme,
        )
    else:
        # ── 簇状条形图：保留底层表格（组列4.5cm + 选项列4.5cm） ──
        n_table_rows = n_total_rows + 1
        table_h = bottom_limit - table_y
        table_row_h = table_h / n_table_rows
        seg_col_w = (table_w - label_total_w) / n_segs

        n_cols = 2 + n_segs
        tbl_shape = slide.shapes.add_table(
            n_table_rows, n_cols,
            Inches(margin_left), Inches(table_y),
            Inches(table_w), Inches(table_h)
        )
        tbl = tbl_shape.table
        try:
            tbl.style = "None"
        except Exception:
            pass
        tbl.columns[0].width = Inches(group_col_w)
        tbl.columns[1].width = Inches(option_col_w)
        for j in range(n_segs):
            tbl.columns[2 + j].width = Inches(seg_col_w)
        for r in range(n_table_rows):
            tbl.rows[r].height = Inches(table_row_h)

        # 表头行
        cell_0 = tbl.cell(0, 0); cell_0.text = "组"
        _style_table_cell(cell_0, font_size=8, bold=True, color=dark_gray,
                          align=PP_ALIGN.CENTER, fill="E8E8E8")
        cell_1 = tbl.cell(0, 1); cell_1.text = "选项"
        _style_table_cell(cell_1, font_size=8, bold=True, color=dark_gray,
                          align=PP_ALIGN.CENTER, fill="E8E8E8")
        for j, seg in enumerate(segments):
            cs = tbl.cell(0, 2 + j); cs.text = "总体" if str(seg).strip().lower() == "total" else seg
            _style_table_cell(cs, font_size=8, bold=True, color=colors[j],
                              align=PP_ALIGN.CENTER, fill="E8E8E8")

        # 数据行
        for mi, mrow in enumerate(merged_rows):
            ri = mi + 1
            g_title, g_start = "", 0
            for gs, ge, gt in group_boundaries:
                if gs <= mi <= ge:
                    g_title, g_start = gt, gs
                    break
            cg = tbl.cell(ri, 0)
            cg.text = _truncate_to_width(g_title, grp_max_units) if mi == g_start else ""
            co = tbl.cell(ri, 1)
            co.text = _truncate_to_width(mrow["选项"], opt_max_units)
            _style_table_cell(cg, font_size=opt_font, bold=True, color=dark_gray, align=PP_ALIGN.CENTER)
            _style_table_cell(co, font_size=opt_font, color=dark_gray, align=PP_ALIGN.LEFT)
            for j in range(n_segs):
                cs = tbl.cell(ri, 2 + j); cs.text = ""
                _style_table_cell(cs, font_size=1, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

        # 合并组标题单元格
        for gs, ge, gt in group_boundaries:
            tbl_gs, tbl_ge = gs + 1, ge + 1
            if tbl_ge - tbl_gs > 0:
                c1, c2 = tbl.cell(tbl_gs, 0), tbl.cell(tbl_ge, 0)
                c1.merge(c2)
                _style_table_cell(c1, font_size=8, bold=True, color=dark_gray, align=PP_ALIGN.CENTER)

        # 组间分隔线
        sep_color = RGBColor(0xAA, 0xAA, 0xAA)
        for gs, ge, gt in group_boundaries:
            if ge >= n_total_rows - 1:
                continue
            sep_row = ge + 1
            for ci in range(n_cols):
                _apply_bottom_border(tbl.cell(sep_row, ci), sep_color)

        # ── 簇状条形图（每分群一列小图） ──
        chart_y = table_y + table_row_h
        chart_h = table_h - table_row_h
        _render_clustered_bars(
            slide, merged_rows, segments,
            margin_left, label_total_w, seg_col_w, n_segs,
            chart_y, chart_h, colors, theme,
        )

    # ── 数据来源 ──
    ds_text = page.data_source or ""
    if ds_text:
        _add_textbox(slide, ds_text, Inches(0.4), Inches(7.12),
                     Inches(12.5), Inches(0.18), theme,
                     size=8, color="758D99")


# ---- 多组多列条形图辅助函数 ----

_NATURAL_ORDER_PATTERNS = [
    r"\d+\s*[-–—～]\s*\d+",           # "18-25", "3-5年"
    r"^\d+$",                        # 纯数字
    r"(?:非常|比较|一般|不(?:太|)?)?(?:满意|同意|喜欢|认可|愿意|重要)",
    r"(?:非常|比较|一般|不(?:太|)?)?(?:不满意|不同意|不喜欢|不认可|不愿意|不重要)",
]
_NATURAL_ORDER_KW = [
    "年龄", "岁", "收入", "月收入", "时长", "小时", "分",
    "频率", "等级", "满意度", "非常", "满意", "一般",
    "不同意", "同意", "NPS", "推荐", " Likely",
]


def _has_natural_order(categories: list) -> bool:
    """判断类目是否有自然有序（年龄/收入/Likert/NPS 等），不应重新排序。"""
    import re
    joined = "|".join(str(c).strip() for c in categories)
    for pat in _NATURAL_ORDER_PATTERNS:
        if re.search(pat, joined, re.I):
            return True
    return any(kw in joined for kw in _NATURAL_ORDER_KW)


def _style_table_cell(cell, font_size=8, bold=False, color=None, align=PP_ALIGN.LEFT,
                      fill=None, border_color=None):
    """设置 PPT 表格单元格字体/填充/边框。"""
    from pptx.dml.color import RGBColor as _RGB
    if border_color is None:
        border_color = _RGB(0xC5, 0xC5, 0xC5)
    if color is None:
        color = _RGB(0x40, 0x40, 0x40)

    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = "微软雅黑"
    p.font.color.rgb = color
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    bc = f"{border_color[0]:02X}{border_color[1]:02X}{border_color[2]:02X}"
    if fill is None:
        fill_hex = None
    elif isinstance(fill, str):
        fill_hex = fill
    elif isinstance(fill, _RGB):
        fill_hex = f"{fill[0]:02X}{fill[1]:02X}{fill[2]:02X}"
    else:
        fill_hex = None
    _set_cell_tcpr(cell, bc, fill_hex)


# ---- 堆积条形图 / 簇状条形图 渲染辅助 ----

def _render_stacked_bars(slide, merged_rows, segments, all_options,
                         margin_left, chart_w, chart_y, chart_h,
                         colors, theme):
    """渲染一张 100% 堆积横向条形图（v11：全宽居中，无左侧标签区）。

    结构：X 轴=分群名（总体/都市中产/…），每根条形按选项堆积到 100%。
    图例在底部显示（避免右侧截断）。适合选项数 ≤7 的题目。
    """
    dark_gray = RGBColor(0x40, 0x40, 0x40)
    inset = 0.15   # 左右边距（给 Y 轴选项标签留空间）
    actual_cx = margin_left + inset
    actual_cw = chart_w - 2 * inset

    # 维度不多且标签较短时用纵向堆积柱状图，否则使用横向堆积条形图。
    use_column = (
        len(segments) <= 6
        and len(all_options) <= 5
        and max((_text_width_units(s) for s in segments), default=0) <= 9
    )
    stacked_type = (
        XL_CHART_TYPE.COLUMN_STACKED_100 if use_column
        else XL_CHART_TYPE.BAR_STACKED_100
    )
    # 构建数据：categories = 分群名, series = 各选项。
    # 横向条形图首类目显示在底部，因此反转写入，确保总体视觉上位于最上方。
    render_segments = list(segments) if use_column else list(reversed(segments))
    chart_data = ChartData()
    chart_data.categories = [
        "总体" if str(seg).strip().lower() == "total" else seg
        for seg in render_segments
    ]
    for opt in all_options:
        vals = []
        for seg in render_segments:
            v = 0.0
            for row in merged_rows:
                if row["选项"] == opt:
                    v = row.get(seg, 0.0)
                    break
            vals.append(v)
        chart_data.add_series(opt, tuple(vals))
    cht_shape = slide.shapes.add_chart(
        stacked_type,
        Inches(actual_cx), Inches(chart_y),
        Inches(actual_cw), Inches(chart_h),
        chart_data
    )
    cht = cht_shape.chart
    cht.has_title = False

    # 每个系列(选项)分配一个颜色
    n_opts = len(all_options)
    option_colors = []
    for oi in range(n_opts):
        if oi < len(colors):
            option_colors.append(colors[oi])
        else:
            r_val = theme.seg_color(oi % 8) if hasattr(theme, 'seg_color') else RGBColor(0x60 + oi * 25, 0x80, 0xA0)
            option_colors.append(r_val)

    for si, series in enumerate(cht.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = option_colors[si]

    # 数据标签
    try:
        for si, series in enumerate(cht.series):
            series.has_data_labels = True
            series.data_labels.show_value = True
            series.data_labels.number_format = "0.0%"
            series.data_labels.font.size = Pt(7)
            series.data_labels.font.name = "微软雅黑"
            series.data_labels.font.color.rgb = dark_gray
            _set_no_wrap_dl(series)
    except Exception:
        pass

    # 类目轴（Y 轴显示分群名 → 对于横向条形图实际是纵轴）
    cat_axis = cht.category_axis
    try:
        cat_axis.tick_labels.font.size = Pt(8)
        cat_axis.tick_labels.font.name = "微软雅黑"
        cat_axis.tick_labels.font.color.rgb = dark_gray
        # 堆积柱/条形图保留浅灰零基线；多列条形图由另一渲染器生成，仍保持无轴线。
        cat_axis.visible = True
        baseline = cat_axis.format.line
        baseline.fill.solid()
        baseline.fill.fore_color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
        baseline.width = Pt(0.8)
    except Exception:
        pass

    # 数值轴（隐藏）
    val_axis = cht.value_axis
    val_axis.maximum_scale = 1.0
    val_axis.number_format = "0%"
    try:
        val_axis.tick_labels.font.size = Pt(1)
        val_axis.has_major_gridlines = False
        val_axis.format.line.fill.background()
    except Exception:
        pass

    # ── 图例：放底部（避免右侧截断） ──
    cht.has_legend = True
    try:
        legend = cht.legend
        legend.position = XL_LEGEND_POSITION.BOTTOM
        legend.include_in_layout = False
        legend.font.size = Pt(7)
        legend.font.name = "微软雅黑"
        legend.font.color.rgb = dark_gray
        legend.include_in_layout = False
    except Exception:
        pass

    # plot area 收缩到图表上部（h=0.88），底部留出图例空间（不被覆盖）
    try:
        ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"
        chart_el = cht._element
        plot_area = chart_el.find(f".//{{{ns}}}plotArea")
        layout = plot_area.find(f"{{{ns}}}layout")
        if layout is None:
            layout = etree.SubElement(plot_area, f"{{{ns}}}layout")
        ml = layout.find(f"{{{ns}}}manualLayout")
        if ml is None:
            ml = etree.SubElement(layout, f"{{{ns}}}manualLayout")
        for tag, val in [("x", "0.0"), ("y", "0.0"), ("w", "1.0"), ("h", "0.86"),
                         ("xMode", "factor"), ("yMode", "factor"),
                         ("wMode", "factor"), ("hMode", "factor")]:
            el = ml.find(f"{{{ns}}}{tag}")
            if el is None:
                el = etree.SubElement(ml, f"{{{ns}}}{tag}")
            el.set("val", val)
    except Exception:
        pass


def _render_clustered_bars(slide, merged_rows, segments,
                            margin_left, label_total_w, seg_col_w, n_segs,
                            chart_y, chart_h, colors, theme):
    """渲染簇状条形图模式（每分群一列独立 BAR_CLUSTERED，对齐表格分群列）。

    图表起始位置 = margin_left + label_total_w（跳过组+选项列）。
    """
    dark_gray = RGBColor(0x40, 0x40, 0x40)

    for si, seg in enumerate(segments):
        cx = margin_left + label_total_w + si * seg_col_w

        if seg_col_w < 1.2:
            inset = 0.04
        elif seg_col_w < 1.8:
            inset = 0.06
        else:
            inset = 0.10
        actual_cx = cx + inset
        actual_cw = seg_col_w - 2 * inset
        right_edge = actual_cx + actual_cw
        if right_edge > 13.083:
            actual_cw = 13.083 - actual_cx
            if actual_cw < 0.5:
                continue

        color = colors[si]

        chart_data = ChartData()
        chart_data.categories = [r["选项"] for r in reversed(merged_rows)]
        vals = [r.get(seg, 0.0) for r in reversed(merged_rows)]
        chart_data.add_series(
            "总体" if str(seg).strip().lower() == "total" else seg,
            tuple(vals),
        )

        cht_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(actual_cx), Inches(chart_y),
            Inches(actual_cw), Inches(chart_h),
            chart_data
        )
        cht = cht_shape.chart
        cht.has_title = False

        series = cht.series[0]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color

        if actual_cw < 0.5:
            lbl_fmt, lbl_sz = "0%", 4
        elif actual_cw < 0.7:
            lbl_fmt, lbl_sz = "0%", 5
        elif actual_cw < 1.0:
            lbl_fmt, lbl_sz = "0.0%", 6
        else:
            lbl_fmt, lbl_sz = "0.0%", 7

        series.has_data_labels = True
        series.data_labels.show_value = True
        series.data_labels.number_format = lbl_fmt
        series.data_labels.font.size = Pt(lbl_sz)
        series.data_labels.font.name = "微软雅黑"
        series.data_labels.font.color.rgb = dark_gray
        _set_no_wrap_dl(series)

        cat_axis = cht.category_axis
        try:
            cat_axis.visible = False
        except Exception:
            pass
        cat_axis.tick_labels.font.size = Pt(1)
        cat_axis.tick_labels.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cat_axis.format.line.fill.background()

        val_axis = cht.value_axis
        val_axis.maximum_scale = 1.0
        val_axis.number_format = "0.0%"
        val_axis.tick_labels.font.size = Pt(1)
        val_axis.tick_labels.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        val_axis.has_major_gridlines = False
        val_axis.format.line.fill.background()

        cht.has_legend = False
        _set_plot_area_full_height(cht)


def _set_no_wrap_dl(series):
    """设置数据标签 wrap=none 防文字换行。"""
    try:
        dl = series.data_labels
        dl_txPr = dl._element.get_or_add_txPr()
        dl_bodyPr = dl_txPr.find(qn('a:bodyPr'))
        if dl_bodyPr is None:
            dl_bodyPr = parse_xml(
                '<a:bodyPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" wrap="none"/>')
            dl_txPr.insert(0, dl_bodyPr)
        else:
            dl_bodyPr.set('wrap', 'none')
    except Exception:
        pass


# DrawingML CT_TableCellProperties 子元素严格顺序：lnL, lnR, lnT, lnB, ..., 填充组
def _set_cell_tcpr(cell, border_hex, fill_hex, *, line_w=6350):
    """在现有 tcPr 上添加边框和填充（不删除重建），避免文本消失。

    之前用 ``parse_xml`` 删除旧 tcPr 再插入新元素会导致 PowerPoint 不渲染文字。
    现在改为：找到/创建 tcPr → 清理旧边框填充 → 用 SubElement 在原树上下文中添加新子元素。
    """
    tc = cell._tc
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn("a:tcPr"))

    # 清理旧的边框和填充子元素（保留其他属性如 anchor）
    for old_tag in ("lnL", "lnR", "lnT", "lnB",
                     "noFill", "solidFill", "gradFill",
                     "blipFill", "pattFill", "grpFill"):
        old_el = tcPr.find(qn(f"a:{old_tag}"))
        if old_el is not None:
            tcPr.remove(old_el)

    # 用原生 API 设置对齐（安全，不破坏 XML 树结构）
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 在原有 tcPr 下用 SubElement 创建边框（继承父级命名空间，避免 parse_xml 隔离问题）
    def _add_ln(tag):
        el = etree.SubElement(tcPr, qn(f"a:{tag}"))
        el.set("w", str(line_w))
        el.set("cap", "flat")
        el.set("cmpd", "sng")
        el.set("algn", "ctr")
        sf = etree.SubElement(el, qn("a:solidFill"))
        sc = etree.SubElement(sf, qn("a:srgbClr"))
        sc.set("val", border_hex)
        pd = etree.SubElement(el, qn("a:prstDash"))
        pd.set("val", "solid")
        return el

    # 四条边框按 schema 顺序插入到填充之前
    _add_ln("lnL")
    _add_ln("lnR")
    _add_ln("lnT")
    _add_ln("lnB")

    # 填充
    if fill_hex is None:
        etree.SubElement(tcPr, qn("a:noFill"))
    else:
        sf = etree.SubElement(tcPr, qn("a:solidFill"))
        sc = etree.SubElement(sf, qn("a:srgbClr"))
        sc.set("val", fill_hex)


def _apply_bottom_border(cell, border_color):
    """为单元格底部加较深分隔线（替换 a:lnB 的线宽/颜色）。"""
    tc = cell._tc
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is None:
        return
    bc = f"{border_color[0]:02X}{border_color[1]:02X}{border_color[2]:02X}"
    old_lnb = tcPr.find(qn("a:lnB"))
    if old_lnb is not None:
        tcPr.remove(old_lnb)
    # 用 SubElement 在原树上下文中创建（不用 parse_xml）
    lnB = etree.SubElement(tcPr, qn("a:lnB"))
    lnB.set("w", "12700")
    lnB.set("cap", "flat")
    lnB.set("cmpd", "sng")
    lnB.set("algn", "ctr")
    sf = etree.SubElement(lnB, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", bc)
    pd = etree.SubElement(lnB, qn("a:prstDash"))
    pd.set("val", "solid")
    # 确保 lnB 在填充元素之前
    fill_ref = None
    for ch in tcPr:
        tag = ch.tag.split("}")[-1] if "}" in ch.tag else ch.tag
        if tag in ("noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill"):
            fill_ref = ch
            break
    if fill_ref is not None:
        fill_ref.addprevious(lnB)


def _set_plot_area_full_height(chart):
    """通过 XML 设置 plot area 占满整个 chart 内部区域，使条形与表格行精确对齐。"""
    try:
        ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"
        chart_el = chart._element
        plot_area = chart_el.find(f".//{{{ns}}}plotArea")
        if plot_area is None:
            return
        layout = plot_area.find(f"{{{ns}}}layout")
        if layout is None:
            layout = etree.SubElement(plot_area, f"{{{ns}}}layout")
        else:
            old_ml = layout.find(f"{{{ns}}}manualLayout")
            if old_ml is not None:
                layout.remove(old_ml)
        ml = etree.SubElement(layout, f"{{{ns}}}manualLayout")
        for tag, val in [("x", "0.0"), ("y", "0.0"), ("w", "1.0"), ("h", "1.0"),
                         ("xMode", "factor"), ("yMode", "factor"),
                         ("wMode", "factor"), ("hMode", "factor")]:
            el = etree.SubElement(ml, f"{{{ns}}}{tag}")
            el.set("val", val)
    except Exception:
        pass


# ------------------------- 数据附录 -------------------------
def build_appendix(slide, ap: AppendixContent, theme: Theme, dims: Dims) -> None:
    slide_w, slide_h = dims
    _add_textbox(slide, ap.title, Inches(PAGE_MARGIN), Inches(TITLE_TOP),
                  slide_w - Inches(2 * PAGE_MARGIN), Inches(TITLE_HEIGHT), theme,
                  size=28, bold=True, color=theme.primary)

    tbl = ap.table
    if not tbl.headers:
        return

    rows = len(tbl.rows) + 1
    cols = len(tbl.headers)
    tx, ty = Inches(PAGE_MARGIN), Inches(1.4)
    tw = slide_w - Inches(2 * PAGE_MARGIN)
    th = slide_h - Inches(1.4) - Inches(0.9)  # 底部留出来源注
    gf = slide.shapes.add_table(rows, cols, tx, ty, tw, th)
    table = gf.table

    col_w = int(tw / cols)
    for c in range(cols):
        table.columns[c].width = col_w

    for c, h in enumerate(tbl.headers):
        cell = table.cell(0, c)
        cell.text = str(h)
        _style_cell(cell, theme, header=True, row_index=0)

    for r, row in enumerate(tbl.rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            _style_cell(cell, theme, header=False, row_index=r)


def _style_cell(cell, theme: Theme, header: bool, row_index: int) -> None:
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    if header:
        set_shape_fill(cell, theme.primary)
    else:
        zebra = "F2F6FB" if row_index % 2 == 0 else "FFFFFF"
        set_shape_fill(cell, zebra)
    tf = cell.text_frame
    for p in tf.paragraphs:
        for run in p.runs:
            style_font(run.font, theme, size=11, bold=header,
                       color=theme.text_light if header else theme.text_dark)
