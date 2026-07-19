"""Editable research proposal deck renderer driven by validated Deck JSON.

The AI layer supplies semantics only.  All geometry, typography and colors are
owned by this deterministic renderer so the HTML preview and PPTX export can
share one content contract without accepting model-generated coordinates.
"""
from __future__ import annotations

import copy
import re
from dataclasses import dataclass
from datetime import date
from typing import Any

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_MARKER_STYLE, XL_TICK_MARK
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx_report.illustrative_dataset import DISCLAIMER, EXAMPLE_LABEL, audit_illustrative_dataset, normalize_illustrative_dataset
from pptx_report.proposal_templates import plan_deck_templates


VISUAL_TYPES = {
    "cover_product_hero",
    "context_tension_map",
    "decision_tree",
    "evidence_threshold_matrix",
    "dual_track_research_flow",
    "qualitative_design_canvas",
    "quantitative_sample_architecture",
    "questionnaire_decision_journey",
    "concept_funnel_maxdiff_example",
    "pricing_segment_example",
    "decision_output_map",
    "timeline_gantt_risk",
    "deliverable_map",
    "risk_matrix",
    "plan_comparison",
    "pricing_table",
}

VISUAL_FALLBACKS = {
    "flow": "dual_track_research_flow",
    "timeline": "timeline_gantt_risk",
    "matrix": "evidence_threshold_matrix",
    "tree": "decision_tree",
    "map": "decision_output_map",
    "sample": "quantitative_sample_architecture",
    "funnel": "concept_funnel_maxdiff_example",
    "cards": "deliverable_map",
}

TITLE_LIMIT = 38
KEY_MESSAGE_LIMIT = 65
ITEM_LABEL_LIMIT = 12
ITEM_DESCRIPTION_LIMIT = 110
MAX_ITEMS = 7

NAVY = "17324D"
BLUE = "1E5AA8"
CYAN = "23A6B5"
ORANGE = "E88735"
INK = "1B2733"
MUTED = "607080"
LINE = "D8E2EA"
PAPER = "F6F8FA"
WHITE = "FFFFFF"
LIGHT_BLUE = "EAF2FA"
LIGHT_CYAN = "E8F7F8"
LIGHT_ORANGE = "FFF2E7"


class DeckValidationError(ValueError):
    """Raised when Deck JSON cannot be repaired safely."""


def _text(value: Any, limit: int = 0) -> str:
    result = re.sub(r"\s+", " ", str(value or "")).strip()
    return result[:limit] if limit else result


def _items(slide: dict) -> list[dict]:
    raw = slide.get("content")
    if not isinstance(raw, list):
        return []
    normalized = []
    for index, item in enumerate(raw[:MAX_ITEMS]):
        if isinstance(item, str):
            item = {"headline": item}
        if not isinstance(item, dict):
            continue
        normalized.append({
            "id": _text(item.get("id")) or f"item_{index + 1:02d}",
            "label": _text(item.get("label"), ITEM_LABEL_LIMIT),
            "headline": _text(item.get("headline") or item.get("title"), 32),
            "description": _text(item.get("description"), ITEM_DESCRIPTION_LIMIT),
            "priority": int(item.get("priority") or index + 1),
            "source_type": _text(item.get("source_type")) or "project_brief",
        })
    return normalized


def _fallback_visual(value: str) -> str:
    value = _text(value).lower()
    aliases = {
        "project_background": "context_tension_map",
        "opportunity_context": "context_tension_map",
        "key_business_decisions": "decision_tree",
        "decision_framework": "decision_tree",
        "research_path": "dual_track_research_flow",
        "methodology_flow": "dual_track_research_flow",
        "execution_design": "qualitative_design_canvas",
        "qualitative_design": "qualitative_design_canvas",
        "sample_design": "quantitative_sample_architecture",
        "quantitative_design": "quantitative_sample_architecture",
        "report_example": "concept_funnel_maxdiff_example",
        "concept_result_example": "concept_funnel_maxdiff_example",
        "pricing_result_example": "pricing_segment_example",
        "gantt": "timeline_gantt_risk",
        "timeline": "timeline_gantt_risk",
        "delivery_plan": "timeline_gantt_risk",
    }
    value = aliases.get(value, value)
    if value in VISUAL_TYPES:
        return value
    for token, fallback in VISUAL_FALLBACKS.items():
        if token in value:
            return fallback
    return "decision_output_map"


def normalize_deck(deck: dict) -> tuple[dict, list[dict]]:
    """Repair safe schema defects and return (normalized_deck, issues)."""
    if not isinstance(deck, dict):
        raise DeckValidationError("Deck JSON 必须是对象。")
    result = copy.deepcopy(deck)
    issues: list[dict] = []
    slides = result.get("slides")
    if not isinstance(slides, list) or not slides:
        raise DeckValidationError("Deck JSON 缺少 slides。")
    normalized_slides = []
    seen_ids: set[str] = set()
    seen_titles: set[str] = set()
    for index, raw in enumerate(slides[:20]):
        if not isinstance(raw, dict):
            issues.append({"level": "error", "code": "invalid_slide", "slide": index + 1})
            continue
        slide_id = _text(raw.get("id")) or f"slide_{index + 1:02d}"
        if slide_id in seen_ids:
            slide_id = f"slide_{index + 1:02d}"
            issues.append({"level": "warning", "code": "duplicate_id_repaired", "slide": index + 1})
        seen_ids.add(slide_id)
        title = _text(raw.get("title"), TITLE_LIMIT)
        if not title:
            title = "待完善页面"
            issues.append({"level": "error", "code": "empty_title", "slide": index + 1})
        if title in seen_titles:
            issues.append({"level": "warning", "code": "duplicate_title", "slide": index + 1})
        seen_titles.add(title)
        visual = _fallback_visual(raw.get("visual_type") or raw.get("slide_type"))
        if visual != raw.get("visual_type"):
            issues.append({"level": "warning", "code": "visual_type_repaired", "slide": index + 1, "value": visual})
        content = _items(raw)
        if index and not content:
            issues.append({"level": "error", "code": "empty_content", "slide": index + 1})
        refs = raw.get("source_references") if isinstance(raw.get("source_references"), list) else []
        key_message = _text(raw.get("key_message"), KEY_MESSAGE_LIMIT)
        numeric_claim = bool(re.search(r"(?<!slide_)\b\d+(?:\.\d+)?%?\b", " ".join(
            [title, key_message] + [f"{item['headline']} {item['description']}" for item in content]
        )))
        if numeric_claim and not refs and index != 0:
            issues.append({"level": "warning", "code": "unreferenced_number", "slide": index + 1})
        normalized_slides.append({
            "id": slide_id,
            "slide_id": _text(raw.get("slide_id")) or slide_id,
            "order": index + 1,
            "slide_type": _text(raw.get("slide_type")) or ("cover" if index == 0 else visual),
            "title": title,
            "subtitle": _text(raw.get("subtitle"), 80),
            "key_message": key_message,
            "slide_question": _text(raw.get("slide_question") or raw.get("question"), 100),
            "unique_purpose": _text(raw.get("unique_purpose"), 140),
            "previous_slide_relation": _text(raw.get("previous_slide_relation"), 140),
            "next_slide_relation": _text(raw.get("next_slide_relation"), 140),
            "visual_type": "cover_product_hero" if index == 0 else visual,
            "layout_variant": _text(raw.get("layout_variant")) or "default",
            "relation_type": _text(raw.get("relation_type")) or "sequence",
            "content_density": _text(raw.get("content_density")) or "professional",
            "node_count": int(raw.get("node_count") or len(content)),
            "has_stage_output": bool(raw.get("has_stage_output")),
            "has_parallel_tracks": bool(raw.get("has_parallel_tracks") or visual == "dual_track_research_flow"),
            "target_canvas_occupancy": min(.85, max(.5, float(raw.get("target_canvas_occupancy") or .72))),
            "content": content,
            "charts": raw.get("charts") if isinstance(raw.get("charts"), list) else [],
            "source_references": refs,
            "data_status": raw.get("data_status") if raw.get("data_status") in {"verified", "illustrative", "framework_only"} else "framework_only",
            "example_output": bool(raw.get("example_output")),
            "timeline_tasks": raw.get("timeline_tasks") if isinstance(raw.get("timeline_tasks"), list) else [],
            "notes": _text(raw.get("notes"), 240),
            "locked": bool(raw.get("locked")),
        })
    if not normalized_slides:
        raise DeckValidationError("Deck JSON 没有可用页面。")
    page_size = int(result.get("page_size") or len(normalized_slides))
    if page_size != len(normalized_slides):
        issues.append({"level": "warning", "code": "page_size_repaired", "value": len(normalized_slides)})
    normalized = {
        "deck_id": _text(result.get("deck_id")) or "deck_local",
        "project_id": _text(result.get("project_id")) or "project_local",
        "title": _text(result.get("title"), 80) or normalized_slides[0]["title"],
        "subtitle": _text(result.get("subtitle"), 120),
        "language": "zh-CN",
        "aspect_ratio": "16:9",
        "theme": "modern_insight_v2",
        "purpose": _text(result.get("purpose")) or "client_proposal",
        "example_output_mode": _text(result.get("example_output_mode")) or "illustrative",
        "illustrative_dataset_id": _text(result.get("illustrative_dataset_id")),
        "illustrative_dataset": normalize_illustrative_dataset(result.get("illustrative_dataset")),
        "content_density": _text(result.get("content_density")) or "professional",
        "target_canvas_occupancy": min(.82, max(.5, float(result.get("target_canvas_occupancy") or .72))),
        "min_body_characters": int(result.get("min_body_characters") or 160),
        "max_body_characters": int(result.get("max_body_characters") or 300),
        "page_size": len(normalized_slides),
        "source_summary": result.get("source_summary") if isinstance(result.get("source_summary"), dict) else {},
        "slides": normalized_slides,
    }
    planned_slides, template_issues = plan_deck_templates(normalized_slides)
    normalized["slides"] = planned_slides
    issues.extend(template_issues)
    return normalized, issues


def audit_deck(deck: dict) -> dict:
    normalized, issues = normalize_deck(deck)
    required = {"cover", "project_context", "business_decisions", "decision_evidence", "research_path", "sample_design", "questionnaire_design", "decision_outputs", "timeline"}
    present = {slide["slide_type"] for slide in normalized["slides"]}
    for slide_type in sorted(required - present):
        issues.append({"level": "warning", "code": "missing_structure", "value": slide_type})
    scores = []
    for slide in normalized["slides"]:
        score = 10
        if not slide["key_message"] and slide["visual_type"] != "cover_product_hero":
            score -= 1
        if not slide["content"] and slide["visual_type"] != "cover_product_hero":
            score -= 2
        if len(slide["content"]) > 5:
            score -= 1
        if len(slide["title"]) > 32:
            score -= 1
        body_chars = sum(len(item["headline"] + item["description"]) for item in slide["content"])
        if slide["visual_type"] != "cover_product_hero" and not slide["example_output"] and body_chars < 80:
            score -= 2; issues.append({"level": "warning", "code": "low_content_density", "slide": slide["order"]})
        components = {
            "content_professionalism": 9 if body_chars >= 120 or slide["example_output"] else 7,
            "title_body_alignment": 9 if slide["key_message"] and slide["content"] else 7,
            "information_density": 9 if body_chars >= 100 or slide["example_output"] else 6,
            "visual_completion": 9,
            "layout_fit": 9 if slide["target_canvas_occupancy"] >= .6 else 6,
            "project_specificity": 9 if any(token in (slide["title"] + " ".join(i["description"] for i in slide["content"])) for token in ["宠物", "饮水", "养宠", "概念", "样本"]) else 7,
            "data_status": 10 if slide["data_status"] in {"verified", "illustrative", "framework_only"} else 4,
            "proposal_value": 9 if slide["slide_question"] else 6,
            "narrative_link": 9 if slide["previous_slide_relation"] or slide["order"] == 1 else 6,
            "editability": 10,
        }
        if min(components.values()) < 7:
            issues.append({"level": "warning", "code": "slide_score_below_threshold", "slide": slide["order"]})
        scores.append({"slide_id": slide["id"], "score": round(sum(components.values()) / len(components), 1), "components": components, "body_characters": body_chars})
    for index in range(1, len(normalized["slides"])):
        if normalized["slides"][index]["visual_type"] == normalized["slides"][index - 1]["visual_type"]:
            issues.append({"level": "warning", "code": "consecutive_layout", "slide": index + 1})
    for left in range(len(normalized["slides"])):
        left_text = set(re.findall(r"[\u4e00-\u9fff]{2}", normalized["slides"][left]["slide_question"] + normalized["slides"][left]["title"]))
        for right in range(left + 1, len(normalized["slides"])):
            right_text = set(re.findall(r"[\u4e00-\u9fff]{2}", normalized["slides"][right]["slide_question"] + normalized["slides"][right]["title"]))
            union = left_text | right_text
            similarity = len(left_text & right_text) / len(union) if union else 0
            if similarity > .7:
                issues.append({"level": "warning", "code": "story_similarity", "slides": [left + 1, right + 1], "value": round(similarity, 2)})
    example_slides = [slide for slide in normalized["slides"] if slide["example_output"]]
    if normalized["example_output_mode"] == "illustrative" and example_slides:
        issues.extend(audit_illustrative_dataset(normalized["illustrative_dataset"]))
        for slide in example_slides:
            if slide["data_status"] != "illustrative" or not slide["charts"]:
                issues.append({"level": "error", "code": "invalid_example_slide", "slide": slide["order"]})
    for slide in normalized["slides"]:
        match = re.search(r"([一二三四五六七八九十]|\d+)(?:个?模块|项)", slide["title"])
        if match:
            values = {"一": 1, "二": 2, "三": 3, "四": 4, "五": 5, "六": 6, "七": 7, "八": 8, "九": 9, "十": 10}
            expected = int(match.group(1)) if match.group(1).isdigit() else values[match.group(1)]
            if expected != len(slide["content"]):
                issues.append({"level": "error", "code": "module_count_mismatch", "slide": slide["order"]})
    timeline = next((slide for slide in normalized["slides"] if slide["visual_type"] == "timeline_gantt_risk"), None)
    if timeline:
        tasks = timeline["timeline_tasks"]
        if not tasks:
            issues.append({"level": "error", "code": "gantt_tasks_missing", "slide": timeline["order"]})
        else:
            invalid_bounds = [task for task in tasks if not isinstance(task.get("start_day"), (int, float)) or not isinstance(task.get("end_day"), (int, float))]
            if invalid_bounds:
                issues.append({"level": "error", "code": "gantt_bounds_missing", "slide": timeline["order"]})
            else:
                if len({int(task["start_day"]) for task in tasks}) == 1:
                    issues.append({"level": "error", "code": "gantt_same_start", "slide": timeline["order"]})
                if len({int(task["end_day"]) - int(task["start_day"]) for task in tasks}) == 1:
                    issues.append({"level": "error", "code": "gantt_same_duration", "slide": timeline["order"]})
                if any(int(task["start_day"]) < 1 or int(task["end_day"]) > 14 or int(task["end_day"]) < int(task["start_day"]) for task in tasks):
                    issues.append({"level": "error", "code": "gantt_bounds_invalid", "slide": timeline["order"]})
    return {"deck": normalized, "issues": issues, "slide_scores": scores, "ok": not any(i["level"] == "error" for i in issues)}


@dataclass(frozen=True)
class Theme:
    font_cn: str = "Microsoft YaHei"
    font_en: str = "Arial"


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _set_text(shape, text: str, size: float, color: str = INK, bold: bool = False,
              align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin: float = 0.05) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = Theme.font_cn
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _textbox(slide, text, x, y, w, h, size=14, color=INK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.04):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set_text(shape, text, size, color, bold, align, valign, margin)
    return shape


def _rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(line)
    shape.line.width = Pt(0.8)
    return shape


def _line(slide, x, y, w, h, color=BLUE, width=1.6, end_arrow=False):
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + w), Inches(y + h))
    shape.line.color.rgb = _rgb(color); shape.line.width = Pt(width)
    if end_arrow:
        shape.line.end_arrowhead = True
    return shape


def _base_slide(prs, slide_data, page_no: int, project_short: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _rgb(PAPER)
    if slide_data["visual_type"] != "cover_product_hero":
        _textbox(slide, slide_data["title"], .62, .34, 11.85, .58, 27, NAVY, True)
        if slide_data["key_message"]:
            _textbox(slide, slide_data["key_message"], .64, .98, 11.7, .38, 11.5, MUTED)
        _textbox(slide, project_short, .62, 7.12, 4.8, .2, 8, MUTED)
        _textbox(slide, str(page_no), 12.05, 7.08, .55, .25, 8, MUTED, align=PP_ALIGN.RIGHT)
    return slide


def _node_text(slide, item, x, y, w, h, fill=WHITE, accent=BLUE):
    shape = _rect(slide, x, y, w, h, fill, accent, True)
    headline = item.get("headline") or item.get("label") or "待完善"
    description = item.get("description") or ""
    _textbox(slide, headline, x + .16, y + .14, w - .32, .38, 14, NAVY, True)
    if description:
        _textbox(slide, description, x + .16, y + .56, w - .32, h - .68, 10.5, MUTED)
    return shape


def _render_cover(slide, deck, data):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = _rgb(NAVY)
    _rect(slide, .78, .76, .17, 5.7, CYAN, CYAN)
    _textbox(slide, data["title"], 1.28, 1.55, 10.7, 1.25, 31, WHITE, True)
    _textbox(slide, data.get("subtitle") or deck.get("subtitle") or "研究方案", 1.32, 3.0, 9.9, .72, 17, "C9DAEA")
    _textbox(slide, date.today().isoformat(), 1.32, 5.95, 2.8, .32, 10.5, "B7CADB")
    _textbox(slide, "SurveyKit · AI 研究方案", 9.05, 5.95, 3.1, .32, 10.5, WHITE, True, PP_ALIGN.RIGHT)


def _render_horizontal(slide, items, y=2.55):
    items = items[:5] or [{"headline": "待完善"}]
    gap = .28; total_w = 12.0; w = (total_w - gap * (len(items) - 1)) / len(items); x = .66
    for index, item in enumerate(items):
        _node_text(slide, item, x, y, w, 2.15, LIGHT_BLUE if index % 2 == 0 else WHITE, BLUE if index % 2 == 0 else CYAN)
        if index < len(items) - 1:
            _line(slide, x + w + .04, y + 1.06, gap - .08, 0, CYAN, 1.8, True)
        x += w + gap


def _render_challenge_chain(slide, items):
    items = items[:4] or [{"headline": "待完善"}]
    widths = [4.1, 5.35, 6.6, 7.85]
    fills = [LIGHT_BLUE, LIGHT_CYAN, WHITE, LIGHT_ORANGE]
    for index, item in enumerate(items):
        w = widths[index]
        x = (13.333 - w) / 2
        y = 1.62 + index * 1.28
        _node_text(slide, item, x, y, w, .96, fills[index], [BLUE, CYAN, NAVY, ORANGE][index])
        if index < len(items) - 1:
            _line(slide, 6.66, y + .96, 0, .28, CYAN, 1.5, True)


def _render_decision_tree(slide, items):
    items = items[:5] or [{"headline": "待完善", "description": "待补充决策信息"}]
    _textbox(slide, "核心业务决策", .72, 1.58, 1.45, .3, 10, CYAN, True)
    _textbox(slide, "从研究问题到可执行判断", .72, 1.95, 4.9, .55, 22, NAVY, True)
    _line(slide, 2.05, 3.55, 9.95, 0, "AABCC9", 1.4, True)
    centers = [2.05 + i * 2.48 for i in range(len(items))]
    accents = [BLUE, CYAN, NAVY, CYAN, ORANGE]
    for index, (item, cx) in enumerate(zip(items, centers)):
        accent = accents[index]
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - .18), Inches(3.37), Inches(.36), Inches(.36))
        node.fill.solid(); node.fill.fore_color.rgb = _rgb(accent); node.line.fill.background()
        _textbox(slide, f"0{index + 1}", cx - .23, 2.93, .46, .24, 9, accent, True, PP_ALIGN.CENTER)
        _textbox(slide, item.get("headline") or item.get("label") or "待完善", cx - 1.0, 3.95, 2.0, .5, 13, NAVY, True, PP_ALIGN.CENTER)
        _textbox(slide, _text(item.get("description") or "形成对应决策证据", 48), cx - 1.0, 4.58, 2.0, .82, 9.5, MUTED, False, PP_ALIGN.CENTER)
        _rect(slide, cx - .65, 5.65, 1.3, .06, accent, accent)
    _textbox(slide, "最终输出", .72, 6.18, .75, .26, 9.5, ORANGE, True)
    _textbox(slide, "明确继续推进、优化迭代或停止投入的决策条件", 1.58, 6.13, 7.6, .34, 12.5, NAVY, True)


def _render_matrix(slide, items):
    headers = ["业务决策", "研究问题", "关键指标", "决策输出"]
    rows = []
    for item in (items[:5] or [{"headline": "待完善"}]):
        headline = item.get("headline") or "待完善"
        rows.append([item.get("label") or headline, headline, item.get("description"), f"{headline}阶段结论"])
    table_shape = slide.shapes.add_table(len(rows) + 1, 4, Inches(.68), Inches(1.65), Inches(12.0), Inches(4.85))
    table = table_shape.table
    for col in range(4): table.columns[col].width = Inches(3.0)
    for col, header in enumerate(headers):
        cell = table.cell(0, col); cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(NAVY)
        _set_text(cell, header, 11, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, .08)
    for row_index, row in enumerate(rows, 1):
        for col, value in enumerate(row):
            cell = table.cell(row_index, col); cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(WHITE if row_index % 2 else LIGHT_BLUE)
            _set_text(cell, _text(value, 42), 9.5, INK, col == 0, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, .08)


def _render_relationship(slide, items):
    titles = ["研究问题", "研究方法", "计划产出"]
    x_values = [.7, 4.75, 8.8]
    for col, (x, title) in enumerate(zip(x_values, titles)):
        _textbox(slide, title, x, 1.55, 3.3, .35, 13, [BLUE, CYAN, NAVY][col], True, PP_ALIGN.CENTER)
    values = items[:3] or [{"headline": "待完善", "description": "待补充研究信息"}]
    for row, item in enumerate(values):
        y = 2.05 + row * 1.35
        headline = item.get("headline") or item.get("label") or "待完善"
        description = item.get("description") or "待补充研究方法"
        cells = [
            {"headline": headline, "description": "需要验证的关键场景与障碍"},
            {"headline": "组合验证", "description": description},
            {"headline": "决策输入", "description": f"形成{headline}的优先行动"},
        ]
        for col, (x, cell) in enumerate(zip(x_values, cells)):
            _node_text(slide, cell, x, y, 3.3, 1.02, WHITE, [BLUE, CYAN, NAVY][col])
        _line(slide, 4.05, y + .51, .62, 0, LINE, 1.4, True)
        _line(slide, 8.1, y + .51, .62, 0, LINE, 1.4, True)


def _render_hierarchy(slide, items):
    items = items[:4] or [{"headline": "待完善"}]
    widths = [11.4, 9.4, 7.4, 5.4]
    fills = [LIGHT_BLUE, LIGHT_CYAN, WHITE, LIGHT_ORANGE]
    for index, item in enumerate(items):
        w = widths[index]; x = (13.333 - w) / 2; y = 1.62 + index * 1.2
        shape = _rect(slide, x, y, w, .92, fills[index], [BLUE, CYAN, LINE, ORANGE][index], True)
        _set_text(shape, f"{item.get('headline') or item.get('label') or '待完善'}  {item.get('description') or ''}", 12, NAVY, index == 0, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, .08)


def _render_blueprint(slide, items):
    _line(slide, 1.35, 5.95, 10.7, 0, MUTED, 1.2, True)
    _line(slide, 1.35, 5.95, 0, -3.85, MUTED, 1.2, True)
    _textbox(slide, "执行优先级  低 → 高", 8.7, 6.08, 3.35, .3, 10, MUTED, True, PP_ALIGN.RIGHT)
    _textbox(slide, "证据强度\n高\n↑\n低", .55, 2.35, .65, 2.45, 9.5, MUTED, True, PP_ALIGN.CENTER)
    _line(slide, 6.7, 2.1, 0, 3.85, LINE, .75)
    _line(slide, 1.35, 4.03, 10.7, 0, LINE, .75)
    _textbox(slide, "优先验证", 1.65, 2.35, 1.3, .3, 9.5, MUTED, True)
    _textbox(slide, "优先落地", 10.15, 2.35, 1.3, .3, 9.5, CYAN, True)
    _textbox(slide, "计划输出 · 不代表真实数据", 8.8, 1.6, 3.4, .35, 10.5, ORANGE, True, PP_ALIGN.RIGHT)
    positions = [(2.1, 4.7), (4.2, 3.8), (6.3, 4.35), (8.4, 2.65), (10.3, 3.25)]
    for item, (x, y) in zip(items[:5] or [{"headline": "待完善"}], positions):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(.72), Inches(.72))
        shape.fill.solid(); shape.fill.fore_color.rgb = _rgb(LIGHT_CYAN); shape.line.color.rgb = _rgb(CYAN)
        _textbox(slide, item.get("headline") or item.get("label") or "待完善", x - .48, y + .78, 1.7, .48, 9.5, INK, True, PP_ALIGN.CENTER)


def _render_timeline(slide, items):
    items = items[:6] or [{"headline": "待完善"}]
    start_y = 1.62; row_h = .78
    _textbox(slide, "阶段", .72, 1.36, 2.4, .25, 10.5, MUTED, True)
    _textbox(slide, "项目推进节奏", 3.35, 1.36, 8.75, .25, 10.5, MUTED, True)
    for index, item in enumerate(items):
        y = start_y + index * row_h
        _textbox(slide, item.get("headline") or item.get("label") or f"阶段{index + 1}", .72, y + .12, 2.35, .36, 10.5, NAVY, True)
        _rect(slide, 3.3, y + .18, 8.85, .34, "E6EBF0", "E6EBF0", True)
        bar_w = min(8.4, 1.4 + index * 1.25)
        _rect(slide, 3.34, y + .21, bar_w, .28, CYAN if index % 2 else BLUE, CYAN if index % 2 else BLUE, True)
        if item.get("description"):
            _textbox(slide, item["description"], 3.5 + min(bar_w, 6.7), y + .06, 1.85, .42, 8.5, MUTED)


def _render_deliverables(slide, items):
    center = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), Inches(2.35), Inches(2.95), Inches(2.0))
    center.fill.solid(); center.fill.fore_color.rgb = _rgb(NAVY); center.line.color.rgb = _rgb(NAVY)
    _set_text(center, "研究决策\n输出", 17, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    positions = [(.75, 1.7), (9.7, 1.7), (.75, 4.55), (9.7, 4.55)]
    for item, (x, y) in zip(items[:4] or [{"headline": "待完善"}], positions):
        _node_text(slide, item, x, y, 2.9, 1.35, WHITE, ORANGE if y > 4 else CYAN)
        _line(slide, 6.68, 3.34, x + 1.45 - 6.68, y + .68 - 3.34, LINE, 1.2)


def _render_risk(slide, items):
    _rect(slide, 1.25, 1.65, 10.8, 4.95, WHITE, LINE)
    _line(slide, 6.65, 1.65, 0, 4.95, LINE, 1.0); _line(slide, 1.25, 4.12, 10.8, 0, LINE, 1.0)
    labels = ["高影响 / 低概率", "高影响 / 高概率", "低影响 / 低概率", "低影响 / 高概率"]
    positions = [(1.55, 1.9), (6.95, 1.9), (1.55, 4.35), (6.95, 4.35)]
    for index, ((x, y), label) in enumerate(zip(positions, labels)):
        _textbox(slide, label, x, y, 2.2, .3, 10, ORANGE if index == 1 else MUTED, True)
        if index < len(items): _node_text(slide, items[index], x, y + .48, 4.5, 1.25, LIGHT_ORANGE if index == 1 else PAPER, ORANGE if index == 1 else CYAN)


def _render_product_cover(slide, deck, data):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = _rgb(NAVY)
    # Editorial cover grid: one dominant title block and one large editable product visual.
    _textbox(slide, "USER RESEARCH PROPOSAL", .82, .72, 4.7, .28, 10, "7ED0D8", True)
    _textbox(slide, data["title"], .82, 1.42, 6.9, 1.55, 30.5, WHITE, True)
    subtitle = data.get("subtitle") or deck.get("subtitle") or "从需求洞察到产品与上市决策"
    _textbox(slide, subtitle, .86, 3.18, 5.95, .68, 15.5, "C9DAEA")
    _textbox(slide, "研究范围", .86, 4.28, .82, .26, 9.5, "7ED0D8", True)
    _textbox(slide, "概念验证 · 卖点优先级 · 价格空间 · 目标人群", 1.72, 4.25, 5.2, .32, 11.5, WHITE, True)
    _textbox(slide, date.today().isoformat(), .86, 6.34, 1.8, .26, 9.5, "B7CADB")
    _textbox(slide, "SurveyKit · AI 研究方案", 4.45, 6.34, 2.7, .26, 9.5, "B7CADB", True, PP_ALIGN.RIGHT)

    halo = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.05), Inches(.65), Inches(4.45), Inches(4.45))
    halo.fill.solid(); halo.fill.fore_color.rgb = _rgb("1E4663"); halo.line.fill.background()
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.62), Inches(1.22), Inches(3.3), Inches(3.3))
    ring.fill.background(); ring.line.color.rgb = _rgb("3E7087"); ring.line.width = Pt(1)
    # Editable, intentionally generic pet-water-dispenser silhouette.
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.22), Inches(1.25), Inches(2.08), Inches(3.55))
    body.fill.solid(); body.fill.fore_color.rgb = _rgb("173A55"); body.line.color.rgb = _rgb("71C8D2"); body.line.width = Pt(1.4)
    window = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.58), Inches(1.62), Inches(1.36), Inches(1.95))
    window.fill.solid(); window.fill.fore_color.rgb = _rgb("24516B"); window.line.fill.background()
    _line(slide, 10.26, 1.72, 0, 1.62, "6BC4CF", 1.0)
    for y in [2.02, 2.42, 2.82]: _line(slide, 9.82, y, .88, 0, "4A7089", .75)
    bowl = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(8.72), Inches(4.04), Inches(3.15), Inches(1.28))
    bowl.fill.solid(); bowl.fill.fore_color.rgb = _rgb("173A55"); bowl.line.color.rgb = _rgb("71C8D2"); bowl.line.width = Pt(1.2)
    drop = slide.shapes.add_shape(MSO_SHAPE.TEAR, Inches(11.55), Inches(3.22), Inches(.48), Inches(.65))
    drop.fill.solid(); drop.fill.fore_color.rgb = _rgb(CYAN); drop.line.fill.background(); drop.rotation = 20
    _textbox(slide, "智能饮水管理", 8.15, 5.62, 4.25, .4, 15, WHITE, True, PP_ALIGN.CENTER)
    _textbox(slide, "产品示意 · 非真实外观", 8.15, 6.08, 4.25, .24, 9, "9CB6C8", False, PP_ALIGN.CENTER)


def _render_context(slide, items):
    items = items[:4]
    labels = ["MARKET SIGNAL", "USER TENSION", "PRODUCT OPPORTUNITY", "DECISION RISK"]
    accents = [BLUE, "718091", CYAN, ORANGE]
    centers = [1.65, 4.92, 8.18, 11.45]
    _line(slide, 1.65, 3.44, 9.8, 0, "B9C8D4", 1.5, True)
    for index, (item, cx) in enumerate(zip(items, centers)):
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - .18), Inches(3.26), Inches(.36), Inches(.36))
        node.fill.solid(); node.fill.fore_color.rgb = _rgb(accents[index]); node.line.fill.background()
        _textbox(slide, f"0{index + 1}", cx - .16, 2.78, .32, .24, 9, accents[index], True, PP_ALIGN.CENTER)
        _textbox(slide, labels[index], cx - 1.25, 1.55, 2.5, .25, 9.5, accents[index], True, PP_ALIGN.CENTER)
        _textbox(slide, item["headline"], cx - 1.28, 1.94, 2.56, .76, 15.5, NAVY, True, PP_ALIGN.CENTER)
        parts = [part.strip() for part in item["description"].split("｜")]
        signal = parts[0].split("：", 1)[-1] if parts else "待验证"
        response = parts[-1].split("：", 1)[-1] if parts else "形成决策依据"
        _line(slide, cx, 3.62, 0, .45, "C5D1DA", .9)
        _textbox(slide, "关键信号", cx - 1.2, 4.13, .74, .24, 9, accents[index], True)
        _textbox(slide, signal, cx - .42, 4.08, 1.62, .62, 10.2, MUTED)
        _textbox(slide, "研究响应", cx - 1.2, 4.96, .74, .24, 9, accents[index], True)
        _textbox(slide, response, cx - .42, 4.91, 1.62, .72, 10.2, NAVY, True)
    _rect(slide, .72, 6.08, .08, .58, ORANGE, ORANGE)
    _textbox(slide, "研究判断", .95, 6.11, .88, .26, 10, ORANGE, True)
    _textbox(slide, "验证真实需求 × 持续体验 × 商业边界，形成推进、优化或停止的证据。", 1.86, 6.06, 10.4, .38, 12.5, NAVY, True)


def _render_evidence_matrix(slide, items):
    headers = ["业务决策", "核心证据", "判断方式", "对应行动"]
    rows = []
    for item in items[:5]:
        parts = item["description"].split("｜", 1)
        rows.append([item.get("label") or item["headline"], item["headline"], parts[0], parts[1] if len(parts) > 1 else "形成对应行动"])
    shape = slide.shapes.add_table(len(rows) + 1, 4, Inches(.64), Inches(1.58), Inches(12.05), Inches(4.55))
    table = shape.table
    widths = [2.15, 3.2, 3.45, 3.25]
    for col, width in enumerate(widths): table.columns[col].width = Inches(width)
    for col, header in enumerate(headers):
        cell = table.cell(0, col); cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(NAVY)
        _set_text(cell, header, 11, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, .08)
    for row_index, row in enumerate(rows, 1):
        for col, value in enumerate(row):
            cell = table.cell(row_index, col); cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(WHITE if row_index % 2 else LIGHT_BLUE)
            _set_text(cell, value, 10, NAVY if col == 0 else INK, col == 0, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, .08)
    # Keep the matrix compact and use the lower band for the decision rule, rather
    # than stretching sparse rows to fill the slide.
    _rect(slide, .64, 6.28, 12.05, .5, LIGHT_CYAN, "B7E4E8", True)
    _textbox(slide, "判断原则", .83, 6.39, .78, .24, 9.5, CYAN, True)
    _textbox(slide, "优先使用相对比较；有历史基准时校准阈值；无统一通过线时由客户内部标准完成最终裁决。", 1.65, 6.35, 10.75, .3, 10.5, NAVY, True)


def _render_dual_flow(slide, items):
    items = items[:5]
    stage_actions = [
        "确认业务假设、刺激物与核心决策边界",
        "开展用户深访与场景追问，识别行为动因与核心张力",
        "统一概念刺激物信息量，完成可测化改写",
        "执行概念测试、MaxDiff与价格测量",
        "整合产品、人群、价格与传播证据",
    ]
    stage_meta = [
        ("经确认的研究简报与刺激物清单", "用于锁定招募、访谈和测量边界"),
        ("定性小结、用户语言与假设池", "用于优化概念并形成问卷输入"),
        ("标准化概念刺激物与评价维度", "作为定量刺激物和测量对象"),
        ("数据底表、优先级与价格模型", "用于形成产品、人群和价格建议"),
        ("决策结论与上市行动路线图", "进入客户评审与执行规划"),
    ]
    _textbox(slide, "研究动作｜做什么", .66, 1.53, 1.6, .25, 9.5, CYAN, True)
    _textbox(slide, "阶段产出｜交付什么", .66, 3.78, 1.8, .25, 9.5, ORANGE, True)
    _textbox(slide, "下一阶段输入｜如何被使用", .66, 5.25, 2.3, .25, 9.5, BLUE, True)
    _line(slide, .68, 3.43, 11.94, 0, "AABCC9", 1.2, True)
    w = 2.28
    for index, item in enumerate(items):
        x = .66 + index * 2.43
        accent = ORANGE if index == len(items) - 1 else (CYAN if index % 2 else BLUE)
        output, usage = stage_meta[index]
        _textbox(slide, f"0{index + 1}", x, 1.92, .45, .28, 12, accent, True)
        _textbox(slide, item["headline"], x, 2.32, w, .48, 14.5, NAVY, True)
        _textbox(slide, stage_actions[index], x, 2.84, w, .48, 9.6, MUTED)
        _line(slide, x + .12, 3.43, 0, .32, accent, 1.5)
        _textbox(slide, output, x, 4.18, w, .58, 10.3, NAVY, True)
        _textbox(slide, usage, x, 5.65, w, .52, 9.5, MUTED)
        _rect(slide, x, 6.28, w, .07, accent, accent)


def _render_qualitative(slide, items):
    groups = [("人群 / WHO", items[:2], .68, BLUE), ("主题 / WHAT", items[2:4], 4.72, CYAN), ("输出 / SO WHAT", items[4:6], 8.76, ORANGE)]
    for col, (title, values, x, accent) in enumerate(groups):
        _textbox(slide, f"0{col + 1}", x, 1.52, .45, .3, 12, accent, True)
        _textbox(slide, title, x + .52, 1.5, 2.8, .34, 15, NAVY, True)
        _line(slide, x, 2.0, 3.35, 0, accent, 2.0)
        for row, item in enumerate(values or [{"headline": "待完善", "description": ""}]):
            y = 2.38 + row * 1.35
            _textbox(slide, item["headline"], x, y, 3.25, .38, 13, NAVY, True)
            _textbox(slide, _text(item["description"], 70), x, y + .48, 3.25, .62, 9.8, MUTED)
            if row == 0: _line(slide, x, y + 1.17, 3.05, 0, "D8E2EA", .8)
        if col < 2: _textbox(slide, "→", x + 3.55, 3.05, .35, .32, 15, "90A4B3", True, PP_ALIGN.CENTER)
    _rect(slide, .68, 5.45, 11.43, .08, NAVY, NAVY)
    _textbox(slide, "设计守则", .68, 5.78, .82, .28, 10, BLUE, True)
    rules = [("场景真实", "在真实补水与清洁任务中观察行为"), ("刺激可比", "概念、功能与价格材料保持同等信息量"), ("输出可测", "每条发现都转译为问卷选项或待验证假设")]
    for index, (label, value) in enumerate(rules):
        x = 1.82 + index * 3.46
        _textbox(slide, label, x, 5.73, .78, .28, 10, NAVY, True)
        _textbox(slide, value, x + .8, 5.69, 2.45, .48, 9.5, MUTED)


def _render_sample_architecture(slide, items):
    if not items: return
    _textbox(slide, items[0]["headline"], .72, 1.56, 3.6, .42, 17, NAVY, True)
    _textbox(slide, items[0]["description"].split("；")[0], 4.4, 1.6, 3.0, .3, 11, CYAN, True)
    _line(slide, .72, 2.15, 11.9, 0, NAVY, 1.4)
    branches = items[1:3]
    widths = [6.15, 5.75]
    starts = [.72, 6.87]
    for index, (item, x, width) in enumerate(zip(branches, starts, widths)):
        accent = BLUE if index == 0 else CYAN
        _rect(slide, x, 2.55, width, .22, accent, accent)
        _textbox(slide, item["headline"], x, 3.02, 1.7, .38, 14.5, NAVY, True)
        _textbox(slide, _text(item["description"], 92), x + 1.72, 2.96, width - 1.85, .72, 10, MUTED)
        _textbox(slide, "主样本" if index == 0 else "机会样本", x, 3.82, 1.1, .26, 9.5, accent, True)
    _line(slide, 6.55, 2.44, 0, 1.72, "CCD7DF", .8)
    for index, item in enumerate(items[3:6]):
        x = .72 + index * 4.0
        accent = [BLUE, CYAN, ORANGE][index]
        _textbox(slide, f"0{index + 1}", x, 4.55, .42, .26, 10, accent, True)
        _textbox(slide, item["headline"], x + .52, 4.5, 2.95, .38, 13, NAVY, True)
        _textbox(slide, _text(item["description"], 90), x + .52, 5.0, 3.0, .72, 9.8, MUTED)
        _rect(slide, x, 5.92, 3.42, .07, accent, accent)
    _textbox(slide, "样本量来自用户输入或方案计算；概念分配与配额为建议方案，正式执行前待客户确认。", .72, 6.34, 11.85, .3, 9.5, ORANGE, True)


def _render_questionnaire_journey(slide, items):
    items = items[:7]
    centers = [1.08 + index * 1.86 for index in range(7)]
    _rect(slide, .72, 3.18, 11.9, .34, "E9EFF4", "E9EFF4")
    _line(slide, .9, 3.35, 11.45, 0, "9DB2C2", 1.25, True)
    for index, (item, cx) in enumerate(zip(items, centers)):
        parts = [part.strip() for part in item["description"].split("｜")]
        values = {}
        for part in parts:
            label, _, value = part.partition("：")
            values[label] = value or part
        above = index % 2 == 0
        node_color = ORANGE if index in (4, 5) else (CYAN if index in (2, 3, 6) else BLUE)
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - .16), Inches(3.19), Inches(.32), Inches(.32))
        node.fill.solid(); node.fill.fore_color.rgb = _rgb(node_color); node.line.fill.background()
        _textbox(slide, str(index + 1).zfill(2), cx - .22, 3.62, .44, .24, 9, node_color, True, PP_ALIGN.CENTER)
        title_y = 1.58 if above else 4.2
        detail_y = 2.06 if above else 4.72
        line_y = 2.72 if above else 3.51
        _line(slide, cx, line_y, 0, (3.19 - line_y) if above else .62, "C2CFD8", .9)
        _textbox(slide, item["headline"], cx - .75, title_y, 1.5, .44, 13, NAVY, True, PP_ALIGN.CENTER)
        _textbox(slide, "指标", cx - .75, detail_y, .36, .22, 8.5, node_color, True)
        _textbox(slide, values.get("指标", "关键指标"), cx - .37, detail_y - .03, 1.12, .48, 9.3, MUTED)
        _textbox(slide, "输出", cx - .75, detail_y + .55, .36, .22, 8.5, ORANGE, True)
        _textbox(slide, values.get("输出", "阶段结论"), cx - .37, detail_y + .52, 1.12, .48, 9.3, NAVY, True)
    _textbox(slide, "行为基线", .72, 6.24, 1.1, .25, 9.5, MUTED, True)
    _line(slide, 1.66, 6.36, 9.8, 0, "C5D1DA", 1.0, True)
    _textbox(slide, "上市策略", 11.55, 6.24, 1.0, .25, 9.5, CYAN, True, PP_ALIGN.RIGHT)


def _example_header(slide):
    _textbox(slide, EXAMPLE_LABEL, 9.25, 1.46, 3.45, .28, 10, ORANGE, True, PP_ALIGN.RIGHT)
    _textbox(slide, DISCLAIMER, .65, 6.78, 12.0, .25, 8.5, MUTED)


def _native_bar_chart(slide, categories, values, x, y, w, h, colors, maximum=None, value_suffix=""):
    """Add a PowerPoint-native bar chart with an embedded editable workbook."""
    data = ChartData()
    # PowerPoint renders the first horizontal-bar category at the bottom. Reverse
    # the data so the semantic order supplied by the planner remains top-down.
    render_categories = list(reversed(categories))
    render_values = list(reversed(values))
    render_colors = list(reversed(colors))
    data.categories = render_categories
    data.add_series("示例值", render_values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data
    ).chart
    chart.has_title = False
    chart.has_legend = False
    chart.chart_style = 2
    chart.font.name = "微软雅黑"
    chart.font.size = Pt(9.5)
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = _rgb(CYAN)
    try: chart.plots[0].gap_width = 68
    except Exception: pass
    for point, color in zip(series.points, render_colors):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = _rgb(color)
        point.format.line.fill.background()
    series.has_data_labels = True
    labels = series.data_labels
    labels.show_value = True
    labels.position = XL_LABEL_POSITION.OUTSIDE_END
    labels.number_format = f'0"{value_suffix}"' if value_suffix else "0"
    labels.number_format_is_linked = False
    labels.font.name = "微软雅黑"
    labels.font.size = Pt(9.5)
    labels.font.bold = True
    labels.font.color.rgb = _rgb(MUTED)
    cat_axis = chart.category_axis
    cat_axis.visible = True
    cat_axis.major_tick_mark = XL_TICK_MARK.NONE
    cat_axis.minor_tick_mark = XL_TICK_MARK.NONE
    cat_axis.tick_labels.font.name = "微软雅黑"
    cat_axis.tick_labels.font.size = Pt(10)
    cat_axis.tick_labels.font.color.rgb = _rgb(NAVY)
    cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.minimum_scale = 0
    if maximum is not None:
        val_axis.maximum_scale = maximum
    val_axis.visible = False
    val_axis.has_major_gridlines = False
    val_axis.format.line.fill.background()
    return chart


def _native_funnel_chart(slide, categories, values, x, y, w, h, colors, maximum=100):
    """Editable centered-bar funnel built as a native stacked bar chart."""
    render_categories = list(reversed(categories))
    render_values = list(reversed(values))
    render_colors = list(reversed(colors))
    data = ChartData()
    data.categories = render_categories
    data.add_series("居中留白", [(maximum - value) / 2 for value in render_values])
    data.add_series("指标值", render_values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED, Inches(x), Inches(y), Inches(w), Inches(h), data
    ).chart
    chart.has_title = False
    chart.has_legend = False
    chart.chart_style = 2
    chart.font.name = "微软雅黑"
    chart.font.size = Pt(9.5)
    try: chart.plots[0].gap_width = 44
    except Exception: pass
    offset, visible = chart.series
    offset.format.fill.background(); offset.format.line.fill.background()
    visible.format.fill.solid(); visible.format.fill.fore_color.rgb = _rgb(CYAN)
    for point, color in zip(visible.points, render_colors):
        point.format.fill.solid(); point.format.fill.fore_color.rgb = _rgb(color)
        point.format.line.fill.background()
    visible.has_data_labels = True
    labels = visible.data_labels
    labels.show_value = True
    labels.position = XL_LABEL_POSITION.CENTER
    labels.number_format = '0"%"'
    labels.number_format_is_linked = False
    labels.font.name = "微软雅黑"; labels.font.size = Pt(9.5); labels.font.bold = True
    labels.font.color.rgb = _rgb(WHITE)
    cat_axis = chart.category_axis
    cat_axis.visible = True
    cat_axis.major_tick_mark = XL_TICK_MARK.NONE; cat_axis.minor_tick_mark = XL_TICK_MARK.NONE
    cat_axis.tick_labels.font.name = "微软雅黑"; cat_axis.tick_labels.font.size = Pt(9.5)
    cat_axis.tick_labels.font.color.rgb = _rgb(NAVY); cat_axis.format.line.fill.background()
    val_axis = chart.value_axis
    val_axis.minimum_scale = 0; val_axis.maximum_scale = maximum; val_axis.visible = False
    val_axis.has_major_gridlines = False; val_axis.format.line.fill.background()
    return chart


def _native_line_chart(slide, categories, values, x, y, w, h, maximum=100, highlight_index=None):
    """Add a PowerPoint-native line chart with editable categories and values."""
    data = ChartData()
    data.categories = categories
    data.add_series("购买意愿", values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(x), Inches(y), Inches(w), Inches(h), data
    ).chart
    chart.has_title = False
    chart.has_legend = False
    chart.chart_style = 2
    chart.font.name = "微软雅黑"
    chart.font.size = Pt(9)
    series = chart.series[0]
    series.format.line.color.rgb = _rgb(CYAN)
    series.format.line.width = Pt(2)
    series.marker.style = XL_MARKER_STYLE.CIRCLE
    series.marker.size = 7
    series.marker.format.fill.solid()
    series.marker.format.fill.fore_color.rgb = _rgb(CYAN)
    series.marker.format.line.color.rgb = _rgb(CYAN)
    if highlight_index is not None and 0 <= highlight_index < len(series.points):
        point = series.points[highlight_index]
        point.format.fill.solid(); point.format.fill.fore_color.rgb = _rgb(ORANGE)
        point.format.line.color.rgb = _rgb(ORANGE)
    series.has_data_labels = True
    labels = series.data_labels
    labels.show_value = True
    labels.position = XL_LABEL_POSITION.ABOVE
    labels.number_format = '0"%"'
    labels.number_format_is_linked = False
    labels.font.name = "微软雅黑"
    labels.font.size = Pt(8.5)
    labels.font.bold = True
    labels.font.color.rgb = _rgb(NAVY)
    cat_axis = chart.category_axis
    cat_axis.visible = True
    cat_axis.major_tick_mark = XL_TICK_MARK.NONE
    cat_axis.minor_tick_mark = XL_TICK_MARK.NONE
    cat_axis.tick_labels.font.name = "微软雅黑"
    cat_axis.tick_labels.font.size = Pt(8.5)
    cat_axis.tick_labels.font.color.rgb = _rgb(MUTED)
    cat_axis.format.line.color.rgb = _rgb(LINE)
    val_axis = chart.value_axis
    val_axis.minimum_scale = 0
    val_axis.maximum_scale = maximum
    val_axis.visible = False
    val_axis.has_major_gridlines = False
    val_axis.format.line.fill.background()
    return chart


def _render_concept_example(slide, dataset, framework_only=False):
    _example_header(slide)
    _textbox(slide, "概念表现漏斗", .62, 1.48, 5.9, .35, 14.5, NAVY, True)
    _textbox(slide, "卖点优先级（MaxDiff）", 6.94, 1.48, 5.7, .35, 14.5, NAVY, True)
    if framework_only or not dataset:
        _textbox(slide, "正确理解 → 相关性 → 独特性 → 可信度 → 购买意愿", .8, 2.35, 5.3, 2.6, 16, MUTED, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        _textbox(slide, "未来将展示卖点相对重要性及优先级", 7.15, 2.35, 5.0, 2.6, 15, MUTED, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        return
    # Only canonical report metrics are placed in native charts. Dataset keys,
    # source tables and debugging metadata never become visible slide content.
    metric_specs = [
        ("concept_understanding", "正确理解"),
        ("concept_relevance", "相关性"),
        ("concept_uniqueness", "独特性"),
        ("concept_credibility", "技术可信度"),
        ("purchase_intention_t2b", "购买意愿"),
    ]
    metrics = dataset.get("metrics") or {}
    metric_labels = [label for _, label in metric_specs]
    metric_values = [float(metrics.get(key) or 0) for key, _ in metric_specs]
    _native_funnel_chart(slide, metric_labels, metric_values, .6, 1.84, 5.96, 4.12,
                         [BLUE, "2E76B8", CYAN, "52B9C4", ORANGE], 100)
    ranked_points = sorted((dataset.get("selling_point_scores") or {}).items(), key=lambda pair: pair[1], reverse=True)[:5]
    selling_labels = [_text(label, 16) for label, _ in ranked_points]
    selling_values = [float(value) for _, value in ranked_points]
    _native_bar_chart(slide, selling_labels, selling_values, 6.88, 1.84, 5.75, 4.12,
                      [CYAN, CYAN, BLUE, BLUE, BLUE], max(selling_values) * 1.28)
    gap = metric_values[0] - metric_values[-1]
    top_two = sum(selling_values[:2])
    _rect(slide, .68, 6.05, .08, .48, ORANGE, ORANGE)
    _textbox(slide, "示例解读", .92, 6.08, .88, .25, 9.5, ORANGE, True)
    _textbox(slide, f"概念漏斗存在 {gap:g}pt 转化落差；排名前两位卖点合计 {top_two:g} 分，应优先验证其组合表达与购买驱动。", 1.8, 6.02, 10.5, .46, 11.2, NAVY, True)


def _render_pricing_example(slide, dataset, framework_only=False):
    _example_header(slide)
    _textbox(slide, "价格—购买意愿曲线", .62, 1.48, 5.9, .35, 14.5, NAVY, True)
    _textbox(slide, "机会人群购买意愿", 6.94, 1.48, 5.7, .35, 14.5, NAVY, True)
    if framework_only or not dataset:
        _textbox(slide, "未来将展示不同价格点的购买概率", 1.35, 3.25, 4.7, .6, 14, MUTED, True, PP_ALIGN.CENTER)
        _textbox(slide, "未来将比较机会人群的购买意愿", 7.4, 3.25, 4.7, .6, 14, MUTED, True, PP_ALIGN.CENTER)
        return
    pricing = dataset.get("pricing") or {}
    curve = sorted(pricing.get("purchase_curve") or [], key=lambda point: point["price"])
    _native_line_chart(slide, [f"¥{point['price']}" for point in curve],
                       [point["intention"] for point in curve], .6, 1.84, 5.96, 3.92, 70, 1)
    ranked_segments = sorted((dataset.get("segment_purchase_intention") or {}).items(), key=lambda pair: pair[1], reverse=True)[:5]
    segment_labels = [_text(label, 18) for label, _ in ranked_segments]
    segment_values = [float(value) for _, value in ranked_segments]
    _native_bar_chart(slide, segment_labels, segment_values, 6.88, 1.84, 5.75, 3.92,
                      [CYAN, CYAN, BLUE, BLUE], 100, "%")
    key_points = [("可接受下限", pricing["acceptable_low"], NAVY), ("最优价位", pricing["optimal_price"], ORANGE), ("可接受上限", pricing["acceptable_high"], NAVY)]
    for index, (label, value, color) in enumerate(key_points):
        x = .72 + index * 1.78
        _textbox(slide, label, x, 5.91, 1.28, .22, 8.3, ORANGE if index == 1 else MUTED, True, PP_ALIGN.CENTER)
        _textbox(slide, f"¥{value}", x, 6.15, 1.28, .36, 17.5 if index != 1 else 20, color, True, PP_ALIGN.CENTER)
        if index < 2: _line(slide, x + 1.32, 6.3, .42, 0, "AABCC9", 1.0)
    _rect(slide, 6.98, 5.94, .07, .55, CYAN, CYAN)
    _textbox(slide, "示例解读", 7.24, 5.96, .86, .24, 9.2, CYAN, True)
    _textbox(slide, f"机会人群购买意愿领先；围绕 ¥{pricing['optimal_price']} 锚点验证溢价，并以 ¥{pricing['acceptable_low']}—¥{pricing['acceptable_high']} 作为可接受区间。", 8.1, 5.9, 4.18, .58, 10.2, NAVY, True)


def _render_decision_output(slide, items):
    stages = [("研究指标", "需求、概念、卖点、价格、人群"), ("分析模型", "漏斗、MaxDiff、PSM、驱动、细分"), ("关键发现", "总体规律、差异与关键障碍"), ("业务行动", "产品、价格、人群与营销动作")]
    for index, (label, description) in enumerate(stages):
        x = .68 + index * 3.05
        _textbox(slide, f"0{index + 1}", x, 1.55, .42, .38, 17, [BLUE, CYAN, NAVY, ORANGE][index], True)
        _textbox(slide, label, x + .55, 1.57, 1.62, .34, 13.5, NAVY, True)
        _textbox(slide, description, x, 2.08, 2.62, .55, 10.2, MUTED)
        if index < 3: _line(slide, x + 2.48, 1.76, .42, 0, "B8C7D2", 1.0, True)
    band = _rect(slide, .65, 2.88, 12.02, .62, NAVY, NAVY)
    band.line.fill.background()
    _textbox(slide, "研究证据最终汇聚为 7 类可执行上市结论", .92, 3.03, 5.6, .28, 13, WHITE, True)
    _textbox(slide, "从判断到行动", 10.45, 3.04, 1.9, .24, 9.5, "9FD4DA", True, PP_ALIGN.RIGHT)
    for index, item in enumerate(items[:7]):
        x = .65 + index * 1.72
        accent = ORANGE if index in (4, 6) else (CYAN if index in (2, 5) else BLUE)
        _rect(slide, x, 3.82, 1.48, .06, accent, accent)
        _textbox(slide, str(index + 1).zfill(2), x, 4.08, .35, .28, 9.5, accent, True)
        _textbox(slide, item["headline"], x, 4.48, 1.48, .62, 12.5, NAVY, True)
        _textbox(slide, item["description"], x, 5.18, 1.48, 1.12, 9.3, MUTED)
        if index < 6: _line(slide, x + 1.58, 4.12, 0, 1.92, "D7E0E7", .75)


def _render_real_gantt(slide, tasks, items):
    clean_tasks = []
    for task in tasks[:8]:
        try:
            start = max(1, min(14, int(task.get("start_day") or 1)))
            end = max(start, min(14, int(task.get("end_day") or start)))
        except (TypeError, ValueError):
            continue
        clean_tasks.append({**task, "start_day": start, "end_day": end})
    tasks = clean_tasks
    x0, day_w = 2.92, .54; y0, row_h = 2.05, .44
    timeline_w = day_w * 14
    _textbox(slide, "工作流", .62, 1.55, 1.85, .26, 9.5, MUTED, True)
    _textbox(slide, "D1", x0, 1.55, .5, .24, 8.5, MUTED, True, PP_ALIGN.CENTER)
    _textbox(slide, "D5", x0 + day_w * 4, 1.55, .5, .24, 8.5, MUTED, True, PP_ALIGN.CENTER)
    _textbox(slide, "D10", x0 + day_w * 9, 1.55, .55, .24, 8.5, MUTED, True, PP_ALIGN.CENTER)
    _textbox(slide, "D14", x0 + day_w * 13, 1.55, .55, .24, 8.5, MUTED, True, PP_ALIGN.CENTER)
    _textbox(slide, "关键交付", 10.88, 1.55, 1.55, .26, 9.5, MUTED, True)
    for marker_day in [1, 5, 10, 14]:
        x = x0 + (marker_day - 1) * day_w + day_w / 2
        _line(slide, x, 1.86, 0, 3.72, "D8E2EA", .75)
    # Client gates are semantic milestones, separate from the sparse day guides.
    for day, _ in [(2, "方案确认"), (5, "工具冻结"), (10, "阶段数据")]:
        x = x0 + (day - 1) * day_w + day_w / 2
        _line(slide, x, 1.86, 0, 3.72, "E9A25F", .9)
    for index, task in enumerate(tasks):
        y = y0 + index * row_h
        if index % 2 == 0:
            lane = _rect(slide, .58, y - .03, 11.92, .38, "F0F4F7", "F0F4F7")
            lane.line.fill.background()
        _textbox(slide, task.get("task"), .66, y + .03, 1.95, .24, 9.3, NAVY, True)
        start, end = task["start_day"], task["end_day"]
        bar = _rect(slide, x0 + (start - 1) * day_w + .04, y + .04, (end - start + 1) * day_w - .08, .23,
                    CYAN if index % 2 else BLUE, CYAN if index % 2 else BLUE)
        bar.line.fill.background()
        if task.get("deliverable"):
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.9), Inches(y + .1), Inches(.1), Inches(.1))
            dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(ORANGE if index in (0, 2, 7) else CYAN); dot.line.fill.background()
            _textbox(slide, task["deliverable"], 11.08, y + .02, 1.32, .24, 8.2, MUTED)
    for day, label in [(2, "方案确认"), (5, "工具冻结"), (10, "阶段数据")]:
        x = x0 + (day - 1) * day_w + day_w / 2
        diamond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x - .09), Inches(1.78), Inches(.18), Inches(.18))
        diamond.fill.solid(); diamond.fill.fore_color.rgb = _rgb(ORANGE); diamond.line.fill.background()
        _textbox(slide, label, x - .48, 1.39, .96, .2, 7.8, ORANGE, True, PP_ALIGN.CENTER)
    _line(slide, .66, 5.86, 11.72, 0, "C7D3DC", .8)
    bottom = [("客户参与", "D2方案确认 · D5工具冻结 · D10阶段数据", BLUE),
              ("关键交付", "研究方案 · 定性小结 · 数据底表 · 最终报告", CYAN),
              ("风险控制", items[-1]["description"] if items else "素材变更与确认延迟需前置管理", ORANGE)]
    for index, (label, value, color) in enumerate(bottom):
        x = .68 + index * 4.05
        if index: _line(slide, x - .22, 6.02, 0, .62, "D4DEE6", .8)
        _textbox(slide, label, x, 6.03, .78, .24, 9, color, True)
        _textbox(slide, value, x + .82, 5.99, 3.0, .55, 8.8, NAVY if index < 2 else MUTED, index < 2)



def render_project_background(slide, deck, data) -> dict:
    items = data["content"]
    if data.get("primary_structure") == "contrast":
        _render_relationship(slide, items)
    else:
        _render_context(slide, items)
    return {"renderer": "render_project_background", "template_id": data.get("template_id")}


def render_key_decisions(slide, deck, data) -> dict:
    _render_decision_tree(slide, data["content"])
    return {"renderer": "render_key_decisions", "template_id": data.get("template_id")}


def render_research_path(slide, deck, data) -> dict:
    if data.get("has_parallel_tracks") or data.get("primary_structure") == "dual_track":
        _render_dual_flow(slide, data["content"])
    else:
        _render_timeline(slide, data["content"])
    return {"renderer": "render_research_path", "template_id": data.get("template_id")}


def render_sample_design(slide, deck, data) -> dict:
    if data.get("page_family") == "sample_design" and data.get("primary_structure") == "execution_canvas":
        _render_qualitative(slide, data["content"])
    else:
        _render_sample_architecture(slide, data["content"])
    return {"renderer": "render_sample_design", "template_id": data.get("template_id")}


def render_report_example(slide, deck, data) -> dict:
    visual = data["visual_type"]
    if visual == "pricing_segment_example":
        _render_pricing_example(slide, deck.get("illustrative_dataset"), data["data_status"] != "illustrative")
    else:
        _render_concept_example(slide, deck.get("illustrative_dataset"), data["data_status"] != "illustrative")
    return {"renderer": "render_report_example", "template_id": data.get("template_id")}


def render_gantt(slide, deck, data) -> dict:
    _render_real_gantt(slide, data.get("timeline_tasks") or [], data["content"])
    return {"renderer": "render_gantt", "template_id": data.get("template_id")}


def _render_by_template(slide, deck, data) -> dict:
    renderer = data.get("renderer")
    if data["visual_type"] == "cover_product_hero":
        _render_product_cover(slide, deck, data)
        return {"renderer": "render_cover", "template_id": data.get("template_id")}
    dispatch = {
        "render_project_background": render_project_background,
        "render_key_decisions": render_key_decisions,
        "render_research_path": render_research_path,
        "render_sample_design": render_sample_design,
        "render_report_example": render_report_example,
        "render_gantt": render_gantt,
    }
    if renderer in dispatch:
        return dispatch[renderer](slide, deck, data)
    visual = data["visual_type"]
    items = data["content"]
    if visual == "evidence_threshold_matrix": _render_evidence_matrix(slide, items)
    elif visual == "questionnaire_decision_journey": _render_questionnaire_journey(slide, items)
    elif visual == "decision_output_map": _render_decision_output(slide, items)
    elif visual in {"plan_comparison", "pricing_table"}: _render_matrix(slide, items)
    elif visual == "risk_matrix": _render_risk(slide, items)
    else: _render_deliverables(slide, items)
    return {"renderer": "legacy_visual_renderer", "template_id": data.get("template_id")}


def render_proposal_deck(deck: dict, output_path: str) -> dict:
    audit = audit_deck(deck)
    normalized = audit["deck"]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = normalized["title"]
    prs.core_properties.subject = "SurveyKit AI 调研方案 PPT"
    prs.core_properties.author = "SurveyKit"
    project_short = normalized["title"][:18]
    for page_no, data in enumerate(normalized["slides"], 1):
        slide = _base_slide(prs, data, page_no, project_short)
        diagnostics = _render_by_template(slide, normalized, data)
        data.setdefault("render_diagnostics", diagnostics)
    prs.save(output_path)
    return audit
