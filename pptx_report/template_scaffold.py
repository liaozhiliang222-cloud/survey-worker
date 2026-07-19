"""Build placeholder skeleton PPTX files for proposal page families."""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from pptx_report.proposal_templates import REGISTRY, THEME, registry_as_dict


def _rgb(hex_value: str):
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(hex_value.lstrip("#"))


def _placeholder(slide, name: str, x: float, y: float, w: float, h: float, size: float = 9.5):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = name
    run.font.name = THEME["fonts"]["cn"]
    run.font.size = Pt(size)
    run.font.color.rgb = _rgb(THEME["colors"]["muted"])
    return shape


def _line(slide, x: float, y: float, w: float, h: float, color: str = "D9E1EA"):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    return shape


def _base_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _decorate(slide, template_id: str, family: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(THEME["colors"]["bg"])
    _line(slide, 0.55, 0.28, 0.12, 6.72, THEME["colors"]["cyan"])
    _placeholder(slide, "SK_KICKER", 0.82, 0.25, 2.5, 0.28, 8.5)
    _placeholder(slide, "SK_TITLE", 0.82, 0.62, 8.9, 0.55, 13)
    _placeholder(slide, "SK_SUBTITLE", 0.84, 1.16, 8.8, 0.34, 8.5)
    _placeholder(slide, "SK_BRAND", 0.82, 7.08, 2.2, 0.22, 7.5)
    _placeholder(slide, "SK_PAGE_NO", 12.08, 7.08, 0.45, 0.22, 7.5)
    _placeholder(slide, "SK_FOOTNOTE", 3.2, 7.08, 6.8, 0.22, 7.5)
    keep = _placeholder(slide, "SK_KEEP", 10.55, 0.38, 1.85, 0.28, 7.5)
    keep.text = f"{family} / {template_id}"


def _place_required(slide, placeholders: list[str]) -> None:
    body = [name for name in placeholders if name.startswith("SK_BODY") or name.startswith("SK_NOTE")]
    nodes = [name for name in placeholders if name.startswith("SK_NODE")]
    stages = [name for name in placeholders if name.startswith("SK_STAGE")]
    charts = [name for name in placeholders if name.startswith("SK_CHART") or name.startswith("SK_INSIGHT") or name in {"SK_EXAMPLE_TAG", "SK_DISCLAIMER"}]
    gantt = [name for name in placeholders if name.startswith("SK_TASK") or name.startswith("SK_DELIVERABLE") or name.startswith("SK_RISK") or name == "SK_TIMELINE"]
    sample = [name for name in placeholders if name.startswith("SK_TOTAL") or name.startswith("SK_BRANCH") or name.startswith("SK_ADDON")]

    for index, name in enumerate(body):
        _placeholder(slide, name, 0.9 + (index % 3) * 3.7, 1.8 + (index // 3) * 0.7, 3.2, 0.45)
    for index, name in enumerate(nodes):
        _placeholder(slide, name, 0.9 + (index % 5) * 2.35, 2.05 + (index // 5) * 0.72, 1.8, 0.42)
    for index, name in enumerate(stages):
        _placeholder(slide, name, 0.95 + (index % 4) * 3.0, 2.0 + (index // 4) * 0.82, 2.3, 0.42)
    for index, name in enumerate(charts):
        w = 5.2 if name.startswith("SK_CHART_0") else 2.8
        _placeholder(slide, name, 0.9 + (index % 2) * 6.0, 1.85 + (index // 2) * 0.82, w, 0.52)
    for index, name in enumerate(gantt):
        _placeholder(slide, name, 0.9 + (index % 4) * 3.0, 1.85 + (index // 4) * 0.72, 2.3, 0.42)
    for index, name in enumerate(sample):
        _placeholder(slide, name, 1.0 + (index % 3) * 3.65, 1.9 + (index // 3) * 0.78, 3.0, 0.42)


def create_template_library(output_dir: str | Path) -> dict:
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    for variant in REGISTRY:
        prs = _base_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _decorate(slide, variant.template_id, variant.family)
        _place_required(slide, list(variant.required_placeholders))
        prs.save(output_dir / (variant.template_file or f"{variant.template_id}.pptx"))

    registry_path = output_dir / "template_registry.json"
    registry_path.write_text(json.dumps(registry_as_dict(), ensure_ascii=False, indent=2), encoding="utf-8")
    return {"output_dir": str(output_dir), "template_count": len(REGISTRY), "registry": str(registry_path)}


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser()
    parser.add_argument("output_dir", nargs="?", default="templates/scheme_proposal")
    args = parser.parse_args()
    print(json.dumps(create_template_library(args.output_dir), ensure_ascii=False))

