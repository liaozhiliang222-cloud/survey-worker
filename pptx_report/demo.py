"""完整使用示例。

运行：python demo.py
产出（位于 ../outputs/）：
  - 调研报告示例.pptx         ：无模板版本，覆盖全部 8 种图表 + 4 种布局
  - 调研报告_模板版.pptx       ：加载 .pptx 模板作为设计底
  - 占位符示例.pptx            ：演示 fill_named_placeholders 模板占位符填充
  - sample_data.xlsx            ：演示用的 Excel 数据源

两种「数据」写法都给出：
  1) 字典（ReportSpec.from_dict）—— 最适合接口 / 配置驱动；
  2) dataclass + 便捷工厂（ChartSpec.combo 等）—— 代码里更直观。
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

import pandas as pd

from pptx_report import (
    ChartPageContent,
    LayoutType,
    ReportRenderer,
    ReportSpec,
)
from pptx_report import loaders
from pptx_report.model import ChartSpec

HERE = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.normpath(os.path.join(HERE, "..", "outputs"))
os.makedirs(OUT_DIR, exist_ok=True)


# ===================== 1) 用「字典」描述一份报告内容 =====================
REPORT_DICT = {
    "cover": {
        "title": "2026 年中国智能家居市场调研报告",
        "client": "云栖智能",
        "date": "2026-07-11",
        "subtitle": "定量研究 · 消费者使用与态度（U&A）",
        "logo_path": None,
    },
    "toc": {"sections": []},  # 留空 → 渲染器自动提取章节标题
    "executive_summary": {
        "kpis": [
            {"label": "目标市场规模(亿元)", "value": "1,280", "delta": "+18% 同比"},
            {"label": "品牌认知度", "value": "76%", "delta": "+9pct"},
            {"label": "净推荐值 NPS", "value": "42", "delta": "行业领先"},
            {"label": "复购率", "value": "58%", "delta": "+6pct"},
        ],
        "conclusion": (
            "智能家居市场进入高速成长期，消费者对「互联互通」与「本地化服务」诉求强烈；"
            "云栖在品牌认知上已建立优势，但性价比与售后仍是转化为复购的关键短板，"
            "建议优先补齐服务体系并强化跨设备联动的差异化叙事。"
        ),
    },
    "chart_pages": [
        # —— 单图大版面：组合图（柱状 + 折线双轴）——
        {
            "title": "各产品线销售额及增长率",
            "layout": "single",
            "charts": [{
                "title": "销售额(万元) 与 同比增长率(%)",
                "type": "combo",
                "categories": ["智能音箱", "扫地机器人", "智能门锁", "安防摄像头"],
                "series": [
                    {"name": "销售额(万元)", "values": [320, 280, 210, 150], "axis": "primary"},
                    {"name": "增长率(%)", "values": [12.5, 8.3, 15.2, 5.1], "axis": "secondary", "chart_type": "line"},
                ],
                "secondary_axis_title": "增长率(%)",
                "insight": "智能音箱规模领先，智能门锁增速最快。",
            }],
        },
        # —— 对比式双图：我司 vs 竞品 ——
        {
            "title": "市场份额与销售趋势对比",
            "layout": "dual",
            "charts": [
                {
                    "title": "各品牌市场份额(%)",
                    "type": "bar",
                    "categories": ["云栖", "竞品A", "竞品B", "竞品C", "其他"],
                    "series": [{"name": "份额", "values": [24, 19, 16, 13, 28]}],
                    "insight": "云栖以 24% 居首，但长尾「其他」仍占近三成。",
                },
                {
                    "title": "近 6 个月销售指数趋势",
                    "type": "line",
                    "categories": ["1月", "2月", "3月", "4月", "5月", "6月"],
                    "series": [
                        {"name": "云栖", "values": [100, 108, 115, 121, 130, 138]},
                        {"name": "行业均值", "values": [100, 103, 107, 110, 112, 116]},
                    ],
                    "insight": "云栖增速持续跑赢行业均值。",
                },
            ],
        },
        # —— 仪表盘网格 2x2：4 个构成 / 比较图 ——
        {
            "title": "市场结构与构成分析",
            "layout": "dashboard",
            "charts": [
                {
                    "title": "渠道销售占比",
                    "type": "pie",
                    "categories": ["线上自营", "平台电商", "线下门店", "运营商"],
                    "series": [{"name": "占比", "values": [38, 31, 21, 10]}],
                    "insight": "线上自营 + 平台电商合计近七成。",
                },
                {
                    "title": "营销预算分配",
                    "type": "doughnut",
                    "categories": ["内容种草", "效果广告", "品牌投放", "线下活动"],
                    "series": [{"name": "预算", "values": [34, 29, 22, 15]}],
                    "insight": "内容种草已成第一预算去向。",
                },
                {
                    "title": "各区域产品线销售额构成(万元)",
                    "type": "stacked_bar",
                    "categories": ["华东", "华南", "华北", "西部"],
                    "series": [
                        {"name": "硬件", "values": [180, 120, 140, 80]},
                        {"name": "软件", "values": [90, 70, 60, 40]},
                        {"name": "服务", "values": [60, 40, 55, 30]},
                    ],
                    "insight": "硬件仍是营收主力，软件在华东占比最高。",
                },
                {
                    "title": "竞品多维度对比（5 分制）",
                    "type": "radar",
                    "categories": ["品牌", "质量", "服务", "性价比", "创新"],
                    "series": [
                        {"name": "云栖", "values": [4.5, 4.2, 4.0, 3.8, 4.1]},
                        {"name": "竞品A", "values": [4.0, 4.5, 3.6, 4.2, 3.9]},
                    ],
                    "insight": "云栖品牌 / 服务领先，性价比为短板。",
                },
            ],
        },
        # —— 图文混排：散点 + 侧栏洞察 ——
        {
            "title": "广告投入与销售额相关性",
            "layout": "mixed",
            "charts": [{
                "title": "各城市广告投入(万元) vs 销售额(万元)",
                "type": "scatter",
                "x_values": [10, 20, 30, 40, 50, 60],
                "y_values": [35, 60, 90, 120, 150, 185],
                "series": [{"name": "城市", "values": [35, 60, 90, 120, 150, 185]}],
                "insight": "广告投入与销售额呈强正相关（R²≈0.99）。",
            }],
            "side_insights": [
                "投放 ROI 在 40 万元附近出现边际拐点。",
                "高线城市的单位投放效率显著高于下沉市场。",
                "建议将增量预算向华东 / 华南的高 ROI 城市倾斜。",
            ],
        },
        # —— 仪表盘网格 3x2：6 个图（验证最多 6 图路径）——
        {
            "title": "用户画像与多期趋势",
            "layout": "dashboard",
            "charts": [
                {"title": "用户月增长率(%)", "type": "line",
                 "categories": ["1月", "2月", "3月", "4月", "5月", "6月"],
                 "series": [{"name": "增长率", "values": [5.2, 6.1, 7.4, 8.0, 9.1, 10.3]}],
                 "insight": "用户增长逐月提速。"},
                {"title": "各地区销售额(万元)", "type": "bar",
                 "categories": ["华东", "华南", "华北", "西部"],
                 "series": [{"name": "销售额", "values": [330, 230, 255, 150]}],
                 "insight": "华东为绝对核心市场。"},
                {"title": "用户性别占比", "type": "pie",
                 "categories": ["男", "女"],
                 "series": [{"name": "占比", "values": [54, 46]}],
                 "insight": "性别分布较均衡。"},
                {"title": "年龄段分布", "type": "doughnut",
                 "categories": ["<25", "25-34", "35-44", "45+"],
                 "series": [{"name": "占比", "values": [22, 41, 26, 11]}],
                 "insight": "25-34 岁为消费主力。"},
                {"title": "各季度产品线构成(万元)", "type": "stacked_bar",
                 "categories": ["Q1", "Q2", "Q3", "Q4"],
                 "series": [
                     {"name": "硬件", "values": [200, 230, 260, 290]},
                     {"name": "软件", "values": [110, 130, 150, 170]},
                 ],
                 "insight": "硬件 / 软件同步扩张。"},
                {"title": "服务能力评估(5分制)", "type": "radar",
                 "categories": ["售前", "安装", "售后", "响应", "培训"],
                 "series": [{"name": "云栖", "values": [4.3, 4.0, 3.6, 3.9, 3.7]}],
                 "insight": "售后评分最低，需重点优化。"},
            ],
        },
    ],
    "appendix": {
        "title": "数据附录",
        "table": {
            "headers": ["区域", "销售额(万元)", "同比", "份额", "NPS"],
            "rows": [
                ["华东", "330", "+20%", "26%", "45"],
                ["华南", "230", "+15%", "18%", "41"],
                ["华北", "255", "+17%", "20%", "43"],
                ["西部", "150", "+22%", "12%", "38"],
            ],
        },
        "source": "云栖智能 2026 H1 消费者调研（N=2,400，线上 + 线下配额抽样）",
    },
}


def build_sample_template(path: str) -> None:
    """创建一个带品牌浅色背景的 .pptx 模板，用于演示「模板作为设计底」。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    master = prs.slide_masters[0]
    master.background.fill.solid()
    master.background.fill.fore_color.rgb = RGBColor.from_string("F4F6F9")
    prs.save(path)


def demo_placeholders(template_path: str) -> None:
    """演示 fill_named_placeholders：在模板预设占位符上按名称填充。"""
    prs = Presentation(template_path)
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title 版式含 Title/Subtitle 占位符
    renderer = ReportRenderer()
    renderer.fill_named_placeholders(slide, {
        "Title": "占位符填充示例",
        "Subtitle": "由 ReportRenderer.fill_named_placeholders 写入",
    })
    out = os.path.join(OUT_DIR, "占位符示例.pptx")
    prs.save(out)
    print("  ✓ 已生成:", out)


def main() -> None:
    print("=== pptx_report 示例开始 ===")

    # 1. 由字典构建报告规格
    spec = ReportSpec.from_dict(REPORT_DICT)
    print(f"  ✓ 已解析报告：{len(spec.chart_pages)} 个图表页，"
          f"目录章节 {len(spec.auto_toc())} 项")

    # 2. 追加一个「来自 Excel」的图表页，演示 pandas + openpyxl 数据源
    xlsx_path = os.path.join(OUT_DIR, "sample_data.xlsx")
    pd.DataFrame({
        "区域": ["华东", "华南", "华北", "西部"],
        "销售额": [330, 230, 255, 150],
        "增长率": [0.20, 0.15, 0.17, 0.22],
    }).to_excel(xlsx_path, index=False, engine="openpyxl")
    excel_chart: ChartSpec = loaders.crosstab_to_chart(
        xlsx_path, title="各区域销售额（Excel 数据源）", insight="华东贡献最高营收。")
    spec.chart_pages.append(ChartPageContent(
        title="区域销售明细（来自 Excel）", charts=[excel_chart],
        layout=LayoutType.SINGLE))
    print("  ✓ 已追加 Excel 数据源图表页")

    # 3. 无模板渲染
    renderer = ReportRenderer()  # 可传入 theme=Theme(...) 换肤
    out1 = os.path.join(OUT_DIR, "调研报告示例.pptx")
    renderer.render(spec, out1)
    print("  ✓ 已生成:", out1)

    # 4. 模板作为设计底渲染
    template_path = os.path.join(OUT_DIR, "template.pptx")
    build_sample_template(template_path)
    renderer_t = ReportRenderer(template_path=template_path)
    out2 = os.path.join(OUT_DIR, "调研报告_模板版.pptx")
    renderer_t.render(spec, out2)
    print("  ✓ 已生成(模板版):", out2)

    # 5. 占位符填充演示
    demo_placeholders(template_path)

    print("=== 全部完成 ===")


if __name__ == "__main__":
    main()
