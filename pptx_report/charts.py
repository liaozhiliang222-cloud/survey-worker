"""图表渲染模块。

把 :class:`~pptx_report.model.ChartSpec`（纯数据）渲染为 python-pptx 的原生图表。
统一规范：
  - 去掉网格线，保持简洁；
  - 图例置于底部、不占用绘图区；
  - 全图中文（拉丁 + 东亚字形）；
  - 按主题配色板为系列着色；
  - 标题 / 图例 / 坐标轴严格统一字体。
"""

from __future__ import annotations

from typing import Dict

from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    XL_LEGEND_POSITION,
)
from pptx.oxml.ns import qn
from pptx.util import Inches

from .exceptions import RenderingError, UnsupportedChartTypeError
from .model import ChartSpec, ChartType
from .theme import Theme
from .utils import style_font, style_textframe


# 图表类型 → python-pptx 原生 XL_CHART_TYPE
_CHART_TYPE_MAP: Dict[ChartType, XL_CHART_TYPE] = {
    ChartType.BAR: XL_CHART_TYPE.BAR_CLUSTERED,
    ChartType.LINE: XL_CHART_TYPE.LINE,
    ChartType.PIE: XL_CHART_TYPE.PIE,
    ChartType.DOUGHNUT: XL_CHART_TYPE.DOUGHNUT,
    ChartType.STACKED_BAR: XL_CHART_TYPE.COLUMN_STACKED,
    ChartType.SCATTER: XL_CHART_TYPE.XY_SCATTER,
    ChartType.RADAR: XL_CHART_TYPE.RADAR,
    ChartType.COMBO: XL_CHART_TYPE.COLUMN_CLUSTERED,  # 组合图以柱状为底座
}


def add_chart(slide, spec: ChartSpec, x, y, cx, cy, theme: Theme) -> None:
    """在 slide 的指定矩形 (x, y, cx, cy) 内添加一张图表并套用统一样式。

    Args:
        slide: 目标幻灯片。
        spec: 图表数据规格。
        x, y, cx, cy: 以 EMU 表示的矩形位置与尺寸。
        theme: 视觉主题。
    """
    try:
        ctype = spec.type if isinstance(spec.type, ChartType) else ChartType(spec.type)
        builder = _BUILDERS.get(ctype, _build_category_chart)
        chart_data = builder(spec)
        chart_type = _CHART_TYPE_MAP[ctype]
        gf = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
        chart = gf.chart

        _style_common(chart, spec, theme)
        _color_series(chart, theme)
        if ctype in (ChartType.PIE, ChartType.DOUGHNUT):
            _style_pie_doughnut(chart, theme)
        else:
            # 分类图（柱状 / 条形 / 折线 / 堆积 / 雷达 / 组合底座）统一加数据标签
            if ctype == ChartType.COMBO:
                _apply_combo(chart, spec, theme)
            _apply_data_labels(chart, spec, theme)
    except UnsupportedChartTypeError:
        raise
    except Exception as exc:  # noqa: BLE001
        raise RenderingError(f"图表渲染失败: {spec.title}", chart=spec.title) from exc


# ------------------------- 数据构造 -------------------------
def _build_category_chart(spec: ChartSpec) -> CategoryChartData:
    """通用分类图（柱状 / 条形 / 折线 / 堆积 / 雷达 / 组合底座）。"""
    cd = CategoryChartData()
    cd.categories = [str(c) for c in spec.categories]
    for s in spec.series:
        cd.add_series(s.name, [float(v) for v in s.values])
    return cd


def _build_scatter(spec: ChartSpec) -> XyChartData:
    """散点图使用 XY 数据构造器。"""
    cd = XyChartData()
    name = spec.series[0].name if spec.series else "散点"
    s = cd.add_series(name)
    for xv, yv in zip(spec.x_values, spec.y_values):
        s.add_data_point(float(xv), float(yv))
    return cd


_BUILDERS = {
    ChartType.SCATTER: _build_scatter,
}


# ------------------------- 统一样式 -------------------------
def _style_common(chart, spec: ChartSpec, theme: Theme) -> None:
    """标题、图例、全局字体、网格线等通用处理。"""
    # 标题
    chart.has_title = True
    tf = chart.chart_title.text_frame
    tf.text = spec.title
    style_textframe(tf, theme, size=14, bold=True, color=theme.text_dark)

    # 图例（饼 / 环形图用数据标签代替图例，关闭图例）
    if spec.type not in (ChartType.PIE, ChartType.DOUGHNUT):
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        style_font(chart.legend.font, theme, size=10)
    else:
        chart.has_legend = False

    # 全局字体（影响坐标轴 / 数据标签等）
    style_font(chart.font, theme, size=10)

    # 坐标轴：去网格线 + 统一刻度字体
    # 注意：饼 / 环形 / 散点等图表没有 category_axis，直接访问会抛 ValueError，
    # 因此用「惰性求值 + 异常兜底」逐个安全获取。
    for axis_getter in (lambda: chart.category_axis, lambda: chart.value_axis):
        try:
            axis = axis_getter()
        except Exception:
            continue
        if axis is None:
            continue
        try:
            axis.has_major_gridlines = False
            axis.has_minor_gridlines = False
        except Exception:
            pass
        try:
            style_font(axis.tick_labels.font, theme, size=9)
        except Exception:
            pass


def _color_series(chart, theme: Theme) -> None:
    """按主题人群对比色板为各系列着色（自动循环取色）。"""
    for i, s in enumerate(chart.series):
        try:
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = theme.seg_color(i)
        except Exception:
            pass


def _no_wrap_data_labels(dl) -> None:
    """关闭数据标签自动换行，避免「49.0%」被拆成两行（XML 级 wrap=none）。"""
    try:
        from pptx.oxml import parse_xml
        txPr = dl._element.get_or_add_txPr()
        bodyPr = txPr.find(qn("a:bodyPr"))
        if bodyPr is None:
            bodyPr = parse_xml(
                '<a:bodyPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                'wrap="none"/>'
            )
            txPr.insert(0, bodyPr)
        else:
            bodyPr.set("wrap", "none")
    except Exception:
        pass


def _apply_data_labels(chart, spec: ChartSpec, theme: Theme) -> None:
    """为分类图（柱 / 条 / 折线 / 堆积 / 雷达 / 组合）的每个系列添加数据标签。

    规则（对齐调研公司交付规范）：
      - 百分比图默认显示 ``49.0%``（数据已 ×100，使用 :attr:`Theme.pct_format`）；
      - 系列数 ≤ 4 时全部标注；> 4 时仅标注第一个系列（通常为 Total），
        其余系列靠图例区分，避免标签拥挤重叠；
      - 字号随系列数动态缩小（9 → 7 → 5pt）；
      - 关闭自动换行，保证数字与 % 同行。
    """
    ctype = spec.type if isinstance(spec.type, ChartType) else ChartType(spec.type)
    # 散点图坐标点、雷达图轮廓靠形状对比，标签意义不大且易拥挤，跳过
    if ctype in (ChartType.SCATTER, ChartType.RADAR):
        return
    n_series = len(chart.series)
    if n_series == 0:
        return
    label_all = n_series <= 4
    if n_series >= 5:
        font_size = max(5, theme.data_label_size - 3)
    elif n_series >= 3:
        font_size = max(7, theme.data_label_size - 1)
    else:
        font_size = theme.data_label_size

    for idx, s in enumerate(chart.series):
        if not label_all and idx != 0:
            continue
        spec_series = spec.series[idx] if idx < len(spec.series) else None
        fmt = getattr(spec_series, "value_format", None) or theme.pct_format
        try:
            s.has_data_labels = True
            dl = s.data_labels
            dl.show_value = True
            dl.show_percentage = False
            dl.number_format = fmt
            dl.number_format_is_linked = False
            style_font(dl.font, theme, size=font_size, color=theme.data_label_color)
            _no_wrap_data_labels(dl)
        except Exception:
            pass


def _style_pie_doughnut(chart, theme: Theme) -> None:
    """饼 / 环形图：显示百分比数据标签，每个扇区独立着色。"""
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    # 数据组装阶段已 ×100（33.8 表示 33.8%），直接显示数值 + 字面 %，避免比例重算误差
    dl.show_value = True
    dl.show_percentage = False
    dl.number_format = theme.pct_format
    dl.number_format_is_linked = False
    try:
        dl.position = XL_LABEL_POSITION.BEST_FIT
    except Exception:
        pass
    style_font(dl.font, theme, size=theme.data_label_size, color=theme.data_label_color)
    _no_wrap_data_labels(dl)

    # 扇区着色
    series = chart.series[0]
    for i, point in enumerate(series.points):
        try:
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = theme.color(i)
        except Exception:
            pass


def _apply_combo(chart, spec: ChartSpec, theme: Theme) -> None:
    """组合图：把标记为副轴的系列改为折线并挪到次坐标轴。

    python-pptx 1.x 尚未在高层 API 暴露 ``series.chart_type``，
    因此这里直接操作 chart 的底层 XML（plotArea）实现真正的双轴组合图：
      1) 在 barChart 之后插入一个 lineChart 绘图区；
      2) 把副轴系列从 barChart 移动到 lineChart；
      3) 新增一个次数值轴（valAx）并被 lineChart 引用。
    任何一步失败都会静默回退为「全部柱状」，保证不崩。
    """
    line_idx = [i for i, s in enumerate(spec.series)
                 if (s.axis or "primary") == "secondary"]
    if not line_idx:
        return
    try:
        from pptx.oxml.xmlchemy import OxmlElement
        cs = chart._chartSpace
        pa = cs.plotArea
        bar_chart = pa.find(qn("c:barChart"))
        if bar_chart is None:
            return
        # 主图已有的两个轴 id（分类轴、主轴）
        ax_ids = bar_chart.findall(qn("c:axId"))
        if len(ax_ids) < 2:
            return
        cat_ax = ax_ids[0].get("val")
        val_ax = ax_ids[1].get("val")
        sec_ax = str(int(float(cat_ax)) + 5000)

        bar_sers = bar_chart.findall(qn("c:ser"))
        line_chart = OxmlElement("c:lineChart")
        grouping = OxmlElement("c:grouping")
        grouping.set("val", "standard")
        line_chart.append(grouping)
        # 把副轴系列移入 lineChart
        for idx in sorted(line_idx, reverse=True):
            ser = bar_sers[idx]
            bar_chart.remove(ser)
            line_chart.append(ser)
        smooth = OxmlElement("c:smooth")
        smooth.set("val", "0")
        line_chart.append(smooth)
        ax1 = OxmlElement("c:axId")
        ax1.set("val", cat_ax)
        ax2 = OxmlElement("c:axId")
        ax2.set("val", sec_ax)
        line_chart.append(ax1)
        line_chart.append(ax2)
        # 插入到 barChart 之后（绘图区必须位于坐标轴之前）
        bar_chart.addnext(line_chart)

        # 新增次数值轴
        sec_val = OxmlElement("c:valAx")
        sec_val.set("axId", sec_ax)
        scaling = OxmlElement("c:scaling")
        orient = OxmlElement("c:orientation")
        orient.set("val", "minMax")
        scaling.append(orient)
        sec_val.append(scaling)
        delete = OxmlElement("c:delete")
        delete.set("val", "0")
        sec_val.append(delete)
        ax_pos = OxmlElement("c:axPos")
        ax_pos.set("val", "r")
        sec_val.append(ax_pos)
        crosses = OxmlElement("c:crosses")
        crosses.set("val", "max")
        sec_val.append(crosses)
        cross_ax = OxmlElement("c:crossAx")
        cross_ax.set("val", cat_ax)
        sec_val.append(cross_ax)
        if spec.secondary_axis_title:
            sec_val.append(_build_axis_title(spec.secondary_axis_title))
        primary_val = pa.find(qn("c:valAx"))
        if primary_val is not None:
            primary_val.addnext(sec_val)

        # 折线系列数据标签（副轴均值，如平均年龄），数字格式 0.0
        for ser in line_chart.findall(qn("c:ser")):
            try:
                dLbls = OxmlElement("c:dLbls")
                for tag, val in (
                    ("c:showVal", "1"), ("c:showSerName", "0"),
                    ("c:showCatName", "0"), ("c:showLegendKey", "0"),
                    ("c:showPercent", "0"),
                ):
                    e = OxmlElement(tag)
                    e.set("val", val)
                    dLbls.append(e)
                numFmt = OxmlElement("c:numFmt")
                numFmt.set("formatCode", "0.0")
                numFmt.set("sourceLinked", "0")
                dLbls.append(numFmt)
                txPr = OxmlElement("c:txPr")
                bodyPr = OxmlElement("a:bodyPr")
                bodyPr.set("wrap", "none")
                txPr.append(bodyPr)
                txPr.append(OxmlElement("a:lstStyle"))
                p = OxmlElement("a:p")
                endPr = OxmlElement("a:endParaRPr")
                endPr.set("lang", "zh-CN")
                latin = OxmlElement("a:latin")
                latin.set("typeface", "微软雅黑")
                ea = OxmlElement("a:ea")
                ea.set("typeface", "微软雅黑")
                endPr.append(latin)
                endPr.append(ea)
                p.append(endPr)
                txPr.append(p)
                dLbls.append(txPr)
                ser.append(dLbls)
            except Exception:
                pass
    except Exception:
        # 回退：保持全部柱状，不影响整体报告
        pass


def _build_axis_title(text: str):
    """构造 <c:title> 元素（含中文 latin + ea 字形）。"""
    from pptx.oxml.xmlchemy import OxmlElement
    title = OxmlElement("c:title")
    tx = OxmlElement("c:tx")
    rich = OxmlElement("c:rich")
    rich.append(OxmlElement("a:bodyPr"))
    rich.append(OxmlElement("a:lstStyle"))
    p = OxmlElement("a:p")
    r = OxmlElement("a:r")
    rpr = OxmlElement("a:rPr")
    rpr.set("lang", "zh-CN")
    latin = OxmlElement("a:latin")
    latin.set("typeface", "微软雅黑")
    ea = OxmlElement("a:ea")
    ea.set("typeface", "微软雅黑")
    rpr.append(latin)
    rpr.append(ea)
    r.append(rpr)
    t = OxmlElement("a:t")
    t.text = text
    r.append(t)
    p.append(r)
    rich.append(p)
    tx.append(rich)
    title.append(tx)
    overlay = OxmlElement("c:overlay")
    overlay.set("val", "0")
    title.append(overlay)
    return title

