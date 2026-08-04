"""Single-slide editable PowerPoint exports for research models.

The browser research tools keep their calculations in JavaScript. This module
accepts those calculated results and renders one PowerPoint slide containing a
native Office chart, editable text boxes, and editable KPI cards.
"""

from __future__ import annotations

from io import BytesIO
import math
from uuid import uuid4
from zipfile import ZIP_DEFLATED, ZipFile
from typing import Any, Iterable

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_DATA_LABEL_POSITION,
    XL_LEGEND_POSITION,
    XL_MARKER_STYLE,
    XL_TICK_LABEL_POSITION,
    XL_TICK_MARK,
)
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


MODEL_TYPES = {"psm", "kano", "maxdiff", "driver"}
FONT = "Microsoft YaHei"
COLORS = {
    "navy": "17365D",
    "blue": "2563EB",
    "sky": "60A5FA",
    "teal": "0F766E",
    "green": "18875B",
    "orange": "E87645",
    "amber": "B7791F",
    "red": "B42318",
    "purple": "7C3AED",
    "ink": "24364B",
    "muted": "64748B",
    "grid": "DCE5EE",
    "panel": "F4F7FA",
    "white": "FFFFFF",
}


class ModelChartExportError(ValueError):
    """Raised when a model chart export request is invalid."""


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _number(value: Any, *, default: float | None = None) -> float | None:
    try:
        result = float(value)
    except (TypeError, ValueError):
        return default
    return result if math.isfinite(result) else default


def _rows(value: Any, *, limit: int = 60) -> list[dict[str, Any]]:
    if not isinstance(value, list):
        raise ModelChartExportError("model data must be an array")
    rows = [dict(item) for item in value if isinstance(item, dict)]
    if not rows:
        raise ModelChartExportError("model data cannot be empty")
    return rows[:limit]


def _add_text(
    slide,
    text: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 12,
    bold: bool = False,
    color: str = COLORS["ink"],
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(0.04)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text or "")
    paragraph.alignment = align
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = _rgb(color)


def _add_header(slide, title: str, subtitle: str) -> None:
    _add_text(slide, title, 0.65, 0.30, 12.0, 0.55, size=25, bold=True)
    _add_text(slide, subtitle, 0.66, 0.86, 11.9, 0.34, size=10.5, color=COLORS["muted"])


def _add_panel(slide, x: float, y: float, w: float, h: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(COLORS["white"])
    shape.line.color.rgb = _rgb(COLORS["grid"])
    shape.line.width = Pt(0.8)


def _add_badge(slide, text: str, x: float, y: float, w: float, color: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.28)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    _add_text(slide, text, x + 0.03, y + 0.01, w - 0.06, 0.24, size=8.5, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)


def _set_plot_area_border(chart, color: str, width: float = 0.75) -> None:
    plot_area = chart._chartSpace.chart.plotArea
    existing = plot_area.find(qn("c:spPr"))
    if existing is not None:
        plot_area.remove(existing)
    sp_pr = OxmlElement("c:spPr")
    sp_pr.append(OxmlElement("a:noFill"))
    line = OxmlElement("a:ln")
    line.set("w", str(int(Pt(width))))
    solid_fill = OxmlElement("a:solidFill")
    rgb = OxmlElement("a:srgbClr")
    rgb.set("val", color)
    solid_fill.append(rgb)
    line.append(solid_fill)
    sp_pr.append(line)
    plot_area.append(sp_pr)


def _set_scatter_cell_range_labels(
    series,
    label_values: list[str],
    *,
    range_formula: str,
    font_size: float,
    color: str,
) -> None:
    """Bind native scatter labels to cells in the embedded workbook."""
    ser = series._ser
    existing = ser.find(qn("c:dLbls"))
    if existing is not None:
        ser.remove(existing)

    data_labels = OxmlElement("c:dLbls")
    c15_ns = "http://schemas.microsoft.com/office/drawing/2012/chart"
    for index in range(len(label_values)):
        data_label = OxmlElement("c:dLbl")
        idx = OxmlElement("c:idx")
        idx.set("val", str(index))
        data_label.append(idx)

        tx = OxmlElement("c:tx")
        rich = OxmlElement("c:rich")
        rich.append(OxmlElement("a:bodyPr"))
        rich.append(OxmlElement("a:lstStyle"))
        paragraph = OxmlElement("a:p")
        field = OxmlElement("a:fld")
        field.set("id", "{" + str(uuid4()).upper() + "}")
        field.set("type", "CELLRANGE")
        field.append(OxmlElement("a:rPr"))
        field_text = OxmlElement("a:t")
        field_text.text = "[CELLRANGE]"
        field.append(field_text)
        paragraph.append(field)
        paragraph.append(OxmlElement("a:endParaRPr"))
        rich.append(paragraph)
        tx.append(rich)
        data_label.append(tx)

        label_ext_list = OxmlElement("c:extLst")
        label_ext = OxmlElement("c:ext")
        label_ext.set("uri", "{CE6537A1-D6FC-4F65-9D91-7224C49458BB}")
        etree.SubElement(label_ext, "{http://schemas.microsoft.com/office/drawing/2012/chart}dlblFieldTable")
        show_range = etree.SubElement(label_ext, "{http://schemas.microsoft.com/office/drawing/2012/chart}showDataLabelsRange")
        show_range.set("val", "1")
        label_ext_list.append(label_ext)
        data_label.append(label_ext_list)
        data_labels.append(data_label)

    for tag, value in (
        ("c:dLblPos", "r"),
        ("c:showLegendKey", "0"),
        ("c:showVal", "0"),
        ("c:showCatName", "0"),
        ("c:showSerName", "0"),
        ("c:showPercent", "0"),
        ("c:showBubbleSize", "0"),
        ("c:showLeaderLines", "1"),
    ):
        element = OxmlElement(tag)
        element.set("val", value)
        data_labels.append(element)

    text_properties = OxmlElement("c:txPr")
    text_properties.append(OxmlElement("a:bodyPr"))
    text_properties.append(OxmlElement("a:lstStyle"))
    paragraph = OxmlElement("a:p")
    paragraph_properties = OxmlElement("a:pPr")
    run_properties = OxmlElement("a:defRPr")
    run_properties.set("sz", str(int(font_size * 100)))
    run_properties.set("b", "1")
    solid_fill = OxmlElement("a:solidFill")
    rgb = OxmlElement("a:srgbClr")
    rgb.set("val", color)
    solid_fill.append(rgb)
    run_properties.append(solid_fill)
    for tag in ("a:latin", "a:ea", "a:cs"):
        typeface = OxmlElement(tag)
        typeface.set("typeface", FONT)
        run_properties.append(typeface)
    paragraph_properties.append(run_properties)
    paragraph.append(paragraph_properties)
    paragraph.append(OxmlElement("a:endParaRPr"))
    text_properties.append(paragraph)
    data_labels.append(text_properties)

    labels_ext_list = OxmlElement("c:extLst")
    labels_ext = OxmlElement("c:ext")
    labels_ext.set("uri", "{CE6537A1-D6FC-4F65-9D91-7224C49458BB}")
    show_range = etree.SubElement(labels_ext, "{http://schemas.microsoft.com/office/drawing/2012/chart}showDataLabelsRange")
    show_range.set("val", "1")
    labels_ext_list.append(labels_ext)
    data_labels.append(labels_ext_list)

    x_values = ser.find(qn("c:xVal"))
    if x_values is not None:
        x_values.addprevious(data_labels)
    else:
        ser.append(data_labels)

    series_ext_list = ser.find(qn("c:extLst"))
    if series_ext_list is None:
        series_ext_list = OxmlElement("c:extLst")
        ser.append(series_ext_list)
    data_range_ext = OxmlElement("c:ext")
    data_range_ext.set("uri", "{02D57815-91ED-43CB-92C2-25804820EDAC}")
    data_range = etree.SubElement(data_range_ext, "{http://schemas.microsoft.com/office/drawing/2012/chart}datalabelsRange")
    formula = etree.SubElement(data_range, "{http://schemas.microsoft.com/office/drawing/2012/chart}f")
    formula.text = range_formula
    cache = etree.SubElement(data_range, "{http://schemas.microsoft.com/office/drawing/2012/chart}dlblRangeCache")
    point_count = OxmlElement("c:ptCount")
    point_count.set("val", str(len(label_values)))
    cache.append(point_count)
    for index, label in enumerate(label_values):
        point = OxmlElement("c:pt")
        point.set("idx", str(index))
        value = OxmlElement("c:v")
        value.text = label
        point.append(value)
        cache.append(point)
    series_ext_list.append(data_range_ext)


def _set_chart_workbook_label_column(chart, label_values: list[str]) -> None:
    """Add worse, better, and KANO label columns to the embedded workbook."""
    xlsx_part = chart.part.chart_workbook.xlsx_part
    source = BytesIO(xlsx_part.blob)
    output = BytesIO()
    spreadsheet_ns = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
    ns = {"s": spreadsheet_ns}

    with ZipFile(source, "r") as archive:
        shared_root = etree.fromstring(archive.read("xl/sharedStrings.xml"))
        strings = ["".join(item.itertext()) for item in shared_root.findall("s:si", ns)]

        def shared_index(value: str) -> int:
            if value in strings:
                return strings.index(value)
            strings.append(value)
            item = etree.SubElement(shared_root, f"{{{spreadsheet_ns}}}si")
            item_text = etree.SubElement(item, f"{{{spreadsheet_ns}}}t")
            item_text.text = value
            return len(strings) - 1

        worse_index = shared_index("worse")
        header_index = shared_index("KANO\u5c5e\u6027\u6807\u7b7e")
        label_indexes = [shared_index(label) for label in label_values]
        shared_root.set("count", str(len(label_values) + 3))
        shared_root.set("uniqueCount", str(len(strings)))

        sheet_root = etree.fromstring(archive.read("xl/worksheets/sheet1.xml"))
        dimension = sheet_root.find("s:dimension", ns)
        dimension.set("ref", f"A1:C{len(label_values) + 1}")
        sheet_data = sheet_root.find("s:sheetData", ns)
        rows = sheet_data.findall("s:row", ns)
        for row in rows:
            row.set("spans", "1:3")

        first_row = rows[0]
        first_cell = etree.Element(f"{{{spreadsheet_ns}}}c", r="A1", t="s")
        etree.SubElement(first_cell, f"{{{spreadsheet_ns}}}v").text = str(worse_index)
        first_row.insert(0, first_cell)
        header_cell = etree.SubElement(first_row, f"{{{spreadsheet_ns}}}c", r="C1", t="s")
        etree.SubElement(header_cell, f"{{{spreadsheet_ns}}}v").text = str(header_index)

        for row_number, (row, label_index) in enumerate(zip(rows[1:], label_indexes), start=2):
            label_cell = etree.SubElement(row, f"{{{spreadsheet_ns}}}c", r=f"C{row_number}", t="s")
            etree.SubElement(label_cell, f"{{{spreadsheet_ns}}}v").text = str(label_index)

        replacements = {
            "xl/sharedStrings.xml": etree.tostring(
                shared_root, xml_declaration=True, encoding="UTF-8", standalone=True
            ),
            "xl/worksheets/sheet1.xml": etree.tostring(
                sheet_root, xml_declaration=True, encoding="UTF-8", standalone=True
            ),
        }
        with ZipFile(output, "w", ZIP_DEFLATED) as rewritten:
            for item in archive.infolist():
                rewritten.writestr(
                    item,
                    replacements.get(item.filename, archive.read(item.filename)),
                )
    xlsx_part._blob = output.getvalue()


def _short_label(value: Any, limit: int = 16) -> str:
    text = str(value or "").strip()
    return text if len(text) <= limit else text[: limit - 1] + "…"

def _add_kpis(slide, title: str, entries: Iterable[tuple[str, Any]]) -> None:
    x, y, w = 10.25, 1.38, 2.43
    _add_text(slide, title, x, y, w, 0.35, size=13, bold=True, color=COLORS["navy"])
    cursor = y + 0.48
    for label, value in list(entries)[:5]:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(cursor), Inches(w), Inches(0.78),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(COLORS["panel"])
        card.line.fill.background()
        _add_text(slide, label, x + 0.13, cursor + 0.08, w - 0.26, 0.22, size=9, color=COLORS["muted"])
        value_text = str(value)
        value_size = 15 if len(value_text) <= 16 else 12.5 if len(value_text) <= 24 else 10.5
        _add_text(
            slide,
            value_text,
            x + 0.13,
            cursor + 0.31,
            w - 0.26,
            0.34,
            size=value_size,
            bold=True,
            color=COLORS["navy"],
        )
        cursor += 0.91


def _add_footer(slide, source: str) -> None:
    _add_text(
        slide,
        f"来源：SurveyKit 研究模型计算结果 · {source} · 图表与数据可在 PowerPoint 中编辑",
        0.67, 7.09, 12.0, 0.20,
        size=8.5,
        color=COLORS["muted"],
    )


def _style_chart(chart, *, legend: bool = False) -> None:
    chart.has_title = False
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = FONT
        chart.legend.font.size = Pt(9)
    chart.font.name = FONT
    chart.font.size = Pt(9)
    for getter in (lambda: chart.category_axis, lambda: chart.value_axis):
        try:
            axis = getter()
            axis.tick_labels.font.name = FONT
            axis.tick_labels.font.size = Pt(9)
            axis.tick_labels.font.color.rgb = _rgb(COLORS["muted"])
            axis.has_minor_gridlines = False
        except Exception:
            continue


def _new_presentation(title: str) -> tuple[Presentation, Any]:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.title = title
    presentation.core_properties.subject = "Editable SurveyKit research model chart"
    presentation.core_properties.author = "SurveyKit"
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb("F7F9FC")
    return presentation, slide


def _render_psm(slide, payload: dict[str, Any]) -> None:
    curve = sorted(
        _rows(payload.get("curve")),
        key=lambda item: _number(item.get("price"), default=0) or 0,
    )
    valid = [
        row for row in curve
        if _number(row.get("price")) is not None
        and all(_number(row.get(key)) is not None for key in (
            "tooCheap", "cheap", "expensive", "tooExpensive"
        ))
    ]
    if len(valid) < 2:
        raise ModelChartExportError("PSM curve requires at least two valid points")

    _add_header(slide, "PSM 价格敏感度分析", "四条价格认知曲线与关键交点 · 单页可编辑图表")
    _add_panel(slide, 0.62, 1.30, 9.28, 5.62)
    data = XyChartData()
    definitions = [
        ("太便宜", "tooCheap", COLORS["navy"]),
        ("比较便宜", "cheap", COLORS["blue"]),
        ("比较贵", "expensive", COLORS["green"]),
        ("太贵", "tooExpensive", COLORS["orange"]),
    ]
    for label, key, _ in definitions:
        series = data.add_series(label)
        for row in valid:
            series.add_data_point(float(row["price"]), float(row[key]) / 100.0)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS,
        Inches(0.88), Inches(1.62), Inches(8.72), Inches(4.88), data,
    ).chart
    _style_chart(chart, legend=True)
    chart.category_axis.minimum_scale = float(valid[0]["price"])
    chart.category_axis.maximum_scale = float(valid[-1]["price"])
    chart.category_axis.has_major_gridlines = False
    chart.value_axis.minimum_scale = 0.0
    chart.value_axis.maximum_scale = 1.0
    chart.value_axis.major_unit = 0.25
    chart.value_axis.number_format = "0%"
    chart.value_axis.number_format_is_linked = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = _rgb(COLORS["grid"])
    for series, (_, _, color) in zip(chart.series, definitions):
        series.format.line.color.rgb = _rgb(color)
        series.format.line.width = Pt(2.2)

    def price(name: str) -> str:
        value = _number(payload.get(name))
        return "—" if value is None else f"¥{value:,.1f}"

    _add_kpis(slide, "关键价格点", [
        ("可接受价格区间", payload.get("acceptable") or "—"),
        ("OPP 最优价格点", price("opp")),
        ("IPP 无差异价格点", price("ipp")),
        ("PMC 价格下限", price("pmc")),
        ("PME 价格上限", price("pme")),
    ])
    _add_footer(slide, f"PSM · N={int(_number(payload.get('sampleCount'), default=0) or 0)}")


def _render_kano(slide, payload: Any) -> None:
    rows = _rows(payload, limit=80)
    valid = []
    for row in rows:
        better = _number(row.get("better"))
        worse = _number(row.get("worse"))
        name = str(row.get("name") or "").strip()
        if name and better is not None and worse is not None:
            valid.append({**row, "name": name, "better": better, "worse": worse, "_code": f"A{len(valid) + 1:02d}"})
    if not valid:
        raise ModelChartExportError("KANO data requires name, better, and worse")

    _add_header(slide, "KANO Better–Worse 分析", "需求属性在满意提升与不满意风险上的位置 · 单页可编辑图表")
    _add_panel(slide, 0.62, 1.30, 9.28, 5.62)
    classifications = [
        ("魅力属性", COLORS["green"]),
        ("期望属性", COLORS["blue"]),
        ("必备属性", COLORS["amber"]),
        ("无差异属性", COLORS["muted"]),
    ]
    classification_colors = dict(classifications)
    compact_labels = len(valid) > 20
    label_values = [
        row["_code"] if compact_labels else _short_label(row["name"])
        for row in valid
    ]
    point_colors = [
        classification_colors.get(
            str(row.get("classification") or ""),
            COLORS["blue"],
        )
        for row in valid
    ]
    data = XyChartData()
    data_series = data.add_series("better")
    for row in valid:
        data_series.add_data_point(abs(float(row["worse"])), float(row["better"]))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        Inches(0.88), Inches(1.62), Inches(8.72), Inches(4.88), data,
    ).chart
    _style_chart(chart, legend=False)
    mean_worse = sum(abs(float(row["worse"])) for row in valid) / len(valid)
    mean_better = sum(float(row["better"]) for row in valid) / len(valid)
    for axis in (chart.category_axis, chart.value_axis):
        axis.minimum_scale = 0.0
        axis.maximum_scale = 1.0
        axis.major_unit = 0.25
        axis.number_format = "0.00"
        axis.number_format_is_linked = False
        axis.has_major_gridlines = False
        axis.has_minor_gridlines = False
        axis.major_tick_mark = XL_TICK_MARK.NONE
        axis.minor_tick_mark = XL_TICK_MARK.NONE
        axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        axis.format.line.color.rgb = _rgb("8FA3B8")
        axis.format.line.width = Pt(0.75)
    chart.category_axis.crosses_at = mean_worse
    chart.value_axis.crosses_at = mean_better
    _set_plot_area_border(chart, COLORS["grid"])
    series = chart.series[0]
    series.marker.style = XL_MARKER_STYLE.CIRCLE
    series.marker.size = 8
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = _rgb(COLORS["blue"])
    series.format.line.fill.background()
    for point, color in zip(series.points, point_colors):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = _rgb(color)
        point.format.line.fill.background()
    _set_chart_workbook_label_column(chart, label_values)
    _set_scatter_cell_range_labels(
        series,
        label_values,
        range_formula="Sheet1!$C$2:$C$" + str(len(label_values) + 1),
        font_size=5.8 if compact_labels else 8.2,
        color=COLORS["ink"],
    )

    plot_left, plot_top, plot_width, plot_height = 1.46, 1.91, 7.70, 3.92
    line_x = plot_left + mean_worse * plot_width
    line_y = plot_top + (1.0 - mean_better) * plot_height
    _add_text(slide, f"Worse 均值 {mean_worse:.2f}", line_x + 0.05, plot_top + plot_height - 0.22, 1.18, 0.18, size=7, bold=True, color=COLORS["muted"])
    _add_text(slide, f"Better 均值 {mean_better:.2f}", plot_left + 0.05, line_y - 0.22, 1.18, 0.18, size=7, bold=True, color=COLORS["muted"])
    _add_badge(slide, "魅力属性", plot_left + 0.05, plot_top + 0.05, 0.92, COLORS["green"])
    _add_badge(slide, "期望属性", plot_left + plot_width - 0.97, plot_top + 0.05, 0.92, COLORS["blue"])
    _add_badge(slide, "无差异属性", plot_left + 0.05, plot_top + plot_height - 0.33, 1.05, COLORS["muted"])
    _add_badge(slide, "必备属性", plot_left + plot_width - 0.97, plot_top + plot_height - 0.33, 0.92, COLORS["amber"])

    ranked = sorted(valid, key=lambda row: abs(float(row["worse"])) + float(row["better"]), reverse=True)
    entries = [
        (
            f"{row.get('_code')} · {_short_label(row.get('name'), 10)}",
            f"B {float(row['better']):.2f} · W {float(row['worse']):.2f}",
        )
        for row in ranked[:5]
    ]
    _add_kpis(slide, "重点属性坐标", entries)
    _add_footer(slide, f"KANO · {len(valid)} 个属性 · 象限交叉线为有效属性系数均值" + (" · A01 起按导入顺序编号" if compact_labels else ""))


def _add_bar_chart(
    slide,
    labels: list[str],
    values: list[float],
    *,
    number_format: str,
    positive_color: str,
    negative_color: str | None = None,
    maximum: float | None = None,
) -> None:
    data = CategoryChartData()
    data.categories = labels
    data.add_series("结果", values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.88), Inches(1.62), Inches(8.72), Inches(4.88), data,
    ).chart
    _style_chart(chart, legend=False)
    chart.category_axis.reverse_order = True
    chart.category_axis.has_major_gridlines = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = _rgb(COLORS["grid"])
    chart.value_axis.number_format = number_format
    chart.value_axis.number_format_is_linked = False
    if maximum is not None:
        chart.value_axis.minimum_scale = 0.0
        chart.value_axis.maximum_scale = maximum
    elif values:
        peak = max(abs(value) for value in values) or 1.0
        chart.value_axis.minimum_scale = -peak * 1.2
        chart.value_axis.maximum_scale = peak * 1.2
    plot = chart.plots[0]
    plot.gap_width = 58
    plot.has_data_labels = True
    labels_format = plot.data_labels
    labels_format.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    labels_format.show_value = True
    labels_format.number_format = number_format
    labels_format.number_format_is_linked = False
    labels_format.font.name = FONT
    labels_format.font.size = Pt(9)
    labels_format.font.bold = True
    labels_format.font.color.rgb = _rgb(COLORS["ink"])
    series = chart.series[0]
    series.invert_if_negative = False
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = _rgb(positive_color)
    series.format.line.fill.background()
    for index, value in enumerate(values):
        color = negative_color if negative_color and value < 0 else positive_color
        point = series.points[index]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = _rgb(color)
        point.format.line.color.rgb = _rgb(color)
        point.format.line.width = Pt(1.75)


def _render_maxdiff(slide, payload: Any) -> None:
    rows = _rows(payload, limit=30)
    valid = []
    for row in rows:
        score = _number(row.get("score"))
        item = str(row.get("item") or "").strip()
        if item and score is not None:
            valid.append({**row, "item": item, "score": score})
    if not valid:
        raise ModelChartExportError("MaxDiff data requires item and score")
    valid.sort(key=lambda row: float(row["score"]), reverse=True)
    _add_header(slide, "MaxDiff 相对偏好得分", "Best–Worst 快速计分结果与优先级排序 · 单页可编辑图表")
    _add_panel(slide, 0.62, 1.30, 9.28, 5.62)
    _add_bar_chart(
        slide,
        [str(row["item"]) for row in valid],
        [float(row["score"]) for row in valid],
        number_format="0.00",
        positive_color=COLORS["green"],
        negative_color=COLORS["red"],
    )
    _add_kpis(slide, "排序摘要", [
        (f"TOP {index + 1}", f"{row['item']} · {float(row['score']):.2f}")
        for index, row in enumerate(valid[:5])
    ])
    _add_footer(slide, f"MaxDiff · {len(valid)} 个项目")


def _render_driver(slide, payload: dict[str, Any]) -> None:
    if not isinstance(payload, dict):
        raise ModelChartExportError("driver data must be an object")
    rows = _rows(payload.get("results"), limit=30)
    valid = []
    for row in rows:
        importance = _number(row.get("importance"))
        name = str(row.get("name") or "").strip()
        if name and importance is not None:
            valid.append({**row, "name": name, "importance": importance})
    if not valid:
        raise ModelChartExportError("driver data requires name and importance")
    valid.sort(key=lambda row: float(row["importance"]), reverse=True)
    dv_name = str(payload.get("dvName") or "目标指标")
    _add_header(slide, "关键驱动分析", f"各维度对“{dv_name}”的相对贡献 · 单页可编辑图表")
    _add_panel(slide, 0.62, 1.30, 9.28, 5.62)
    maximum = max(float(row["importance"]) for row in valid)
    axis_maximum = min(1.0, max(0.1, math.ceil(maximum * 12) / 10))
    _add_bar_chart(
        slide,
        [str(row["name"]) for row in valid],
        [float(row["importance"]) for row in valid],
        number_format="0.0%",
        positive_color=COLORS["blue"],
        maximum=axis_maximum,
    )
    r2 = _number(payload.get("r2"), default=0) or 0
    sample = int(_number(payload.get("n"), default=0) or 0)
    top_entries = [
        ("R² / 样本量", f"{r2:.3f} / N={sample}"),
        ("因变量", dv_name),
    ] + [
        (f"驱动 {index + 1}", f"{row['name']} · {float(row['importance']):.1%}")
        for index, row in enumerate(valid[:3])
    ]
    _add_kpis(slide, "模型摘要", top_entries)
    _add_footer(slide, f"关键驱动 · R²={r2:.3f} · N={sample}")


def render_model_chart_pptx(model_type: str, data: Any) -> bytes:
    """Render one editable PPTX slide for a supported research model."""
    normalized = str(model_type or "").strip().lower()
    if normalized not in MODEL_TYPES:
        raise ModelChartExportError(f"unsupported model_type: {normalized or '(empty)'}")
    title = {
        "psm": "PSM 价格敏感度分析",
        "kano": "KANO Better–Worse 分析",
        "maxdiff": "MaxDiff 相对偏好得分",
        "driver": "关键驱动分析",
    }[normalized]
    presentation, slide = _new_presentation(title)
    if normalized == "psm":
        if not isinstance(data, dict):
            raise ModelChartExportError("PSM data must be an object")
        _render_psm(slide, data)
    elif normalized == "kano":
        _render_kano(slide, data)
    elif normalized == "maxdiff":
        _render_maxdiff(slide, data)
    else:
        _render_driver(slide, data)
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()
