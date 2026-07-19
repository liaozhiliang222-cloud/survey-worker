"""Placeholder-based validation for proposal PPTX templates."""

from __future__ import annotations

from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation

from pptx_report.proposal_templates import REGISTRY_BY_ID


def _shape_name(shape) -> str:
    try:
        return str(shape.name or "").strip()
    except Exception:  # noqa: BLE001
        return ""


def _shape_text(shape) -> str:
    try:
        return str(shape.text or "").strip()
    except Exception:  # noqa: BLE001
        return ""


def collect_placeholders(pptx_path: str | Path) -> dict[str, list[dict[str, Any]]]:
    prs = Presentation(str(pptx_path))
    found: dict[str, list[dict[str, Any]]] = {}
    for slide_index, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            name = _shape_name(shape)
            if not name.startswith("SK_"):
                continue
            found.setdefault(name, []).append({
                "slide": slide_index,
                "name": name,
                "text": _shape_text(shape),
                "left": int(shape.left),
                "top": int(shape.top),
                "width": int(shape.width),
                "height": int(shape.height),
            })
    return found


def validate_template_file(pptx_path: str | Path, template_id: str | None = None) -> dict[str, Any]:
    """Validate one placeholder template file.

    When template_id is provided, required placeholders are read from the
    proposal template registry.  Otherwise the validator performs generic
    structural checks.
    """
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    issues: list[dict[str, Any]] = []
    ratio = int(prs.slide_width) / max(1, int(prs.slide_height))
    if not 1.70 <= ratio <= 1.82:
        issues.append({"level": "error", "code": "invalid_page_size", "value": round(ratio, 3)})
    if len(prs.slides) < 1:
        issues.append({"level": "error", "code": "empty_template"})

    placeholders = collect_placeholders(pptx_path)
    counts = Counter({name: len(entries) for name, entries in placeholders.items()})
    for name, count in counts.items():
        if count > 1 and name != "SK_KEEP":
            issues.append({"level": "warning", "code": "duplicate_placeholder", "placeholder": name, "count": count})
    for name, entries in placeholders.items():
        if name == "SK_KEEP":
            continue
        if any(not entry["text"] for entry in entries):
            issues.append({"level": "warning", "code": "empty_placeholder_text", "placeholder": name})

    if template_id:
        variant = REGISTRY_BY_ID.get(template_id)
        if not variant:
            issues.append({"level": "error", "code": "unknown_template_id", "template_id": template_id})
        else:
            missing = sorted(set(variant.required_placeholders) - set(placeholders))
            for name in missing:
                issues.append({"level": "error", "code": "missing_placeholder", "template_id": template_id, "placeholder": name})

    return {
        "file": str(pptx_path),
        "template_id": template_id,
        "slide_count": len(prs.slides),
        "placeholder_count": sum(len(v) for v in placeholders.values()),
        "placeholders": sorted(placeholders),
        "ok": not any(issue["level"] == "error" for issue in issues),
        "issues": issues,
    }


def validate_template_directory(directory: str | Path) -> dict[str, Any]:
    directory = Path(directory)
    reports = []
    for template_id, variant in sorted(REGISTRY_BY_ID.items()):
        if not variant.template_file:
            continue
        path = directory / variant.template_file
        if not path.exists():
            reports.append({
                "file": str(path),
                "template_id": template_id,
                "ok": False,
                "issues": [{"level": "error", "code": "missing_template_file", "template_id": template_id}],
            })
            continue
        reports.append(validate_template_file(path, template_id))
    return {
        "directory": str(directory),
        "templates": reports,
        "ok": all(report["ok"] for report in reports),
        "issue_count": sum(len(report["issues"]) for report in reports),
    }

