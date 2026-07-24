from __future__ import annotations

from io import BytesIO
import json
from pathlib import Path
import tempfile
import threading
import time

from fastapi import Response
from fastapi.testclient import TestClient
from pptx import Presentation

from deploy import aliyun_api as api


def pptx_bytes() -> bytes:
    prs = Presentation()
    for index in range(4):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(100000, 100000, 4000000, 500000)
        box.text = "Source: smoke" if index == 3 else f"Slide {index + 1}"
    stream = BytesIO()
    prs.save(stream)
    return stream.getvalue()


def envelope(metadata: dict, data: bytes = b"xlsx") -> bytes:
    encoded = json.dumps(metadata).encode("utf-8")
    return api.REQUEST_ENVELOPE_MAGIC + len(encoded).to_bytes(4, "big") + encoded + data


def wait_for(client: TestClient, job_id: str, statuses: set[str], timeout: float = 5) -> dict:
    deadline = time.time() + timeout
    state = {}
    while time.time() < deadline:
        state = client.get(f"/api/pptx-report/jobs/{job_id}").json()
        if state.get("status") in statuses:
            return state
        time.sleep(0.02)
    raise AssertionError(f"job {job_id} did not reach {statuses}; last={state}")


def main() -> None:
    original_dir = api.JOB_DIR
    original_generate = api._generate_core
    original_semaphore = api.JOB_SEMAPHORE
    with tempfile.TemporaryDirectory(prefix="surveykit-jobs-") as temp_dir:
        api.JOB_DIR = Path(temp_dir)
        api.JOB_SEMAPHORE = threading.BoundedSemaphore(1)
        client = TestClient(api.app)
        release = threading.Event()

        def blocking_generate(data, qs, metadata, progress_callback=None, **_kwargs):
            if progress_callback:
                progress_callback(20, "blocking")
            release.wait(3)
            return Response(pptx_bytes(), media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")

        api._generate_core = blocking_generate
        headers = {
            "Content-Type": "application/vnd.surveykit.pptx-request",
            "X-SurveyKit-Client-ID": "client-smoke-001",
        }
        first = client.post("/api/pptx-report/jobs", content=envelope({"title": "same"}), headers=headers)
        assert first.status_code == 202
        first_state = first.json()
        assert first_state["created_at"] and first_state["started_at"] is None
        duplicate = client.post("/api/pptx-report/jobs", content=envelope({"title": "same"}), headers=headers)
        assert duplicate.status_code == 202
        assert duplicate.json()["job_id"] == first_state["job_id"]
        assert duplicate.json()["deduplicated"] is True

        cancel = client.post(f"/api/pptx-report/jobs/{first_state['job_id']}/cancel")
        assert cancel.status_code == 202
        release.set()
        cancelled = wait_for(client, first_state["job_id"], {"cancelled"})
        assert cancelled["finished_at"] and not api._job_output_path(first_state["job_id"]).exists()

        api._generate_core = lambda data, qs, metadata, progress_callback=None, **kwargs: Response(
            pptx_bytes(), media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        ready_response = client.post("/api/pptx-report/jobs", content=envelope({"title": "ready"}, b"xlsx-2"), headers=headers)
        ready = wait_for(client, ready_response.json()["job_id"], {"ready"})
        assert ready["started_at"] and ready["finished_at"]
        assert ready["qa"]["slide_count"] == 4 and isinstance(ready["overall_score"], int)
        download = client.get(f"/api/pptx-report/jobs/{ready['job_id']}/download?delete_after=true")
        assert download.status_code == 200 and download.content.startswith(b"PK")
        assert not api._job_state_path(ready["job_id"]).exists()
        assert not api._job_output_path(ready["job_id"]).exists()

        lost_id = "a" * 32
        api._write_job_state(lost_id, {"status": "running", "service_instance_id": "old-instance"})
        raw = json.loads(api._job_state_path(lost_id).read_text(encoding="utf-8"))
        raw["service_instance_id"] = "old-instance"
        api._job_state_path(lost_id).write_text(json.dumps(raw), encoding="utf-8")
        lost = client.get(f"/api/pptx-report/jobs/{lost_id}").json()
        assert lost["status"] == "lost" and lost["finished_at"]

    api.JOB_DIR = original_dir
    api._generate_core = original_generate
    api.JOB_SEMAPHORE = original_semaphore
    print("async job lifecycle smoke: ok")


if __name__ == "__main__":
    main()