from collections import Counter
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZipFile
import os
import shutil
import sys

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.util import Inches

from pptx_report.common.capacity import split_preserving_order
from pptx_report.common.narrative import build_slide_briefs
from pptx_report.facts import build_funnel_facts, extract_data_facts, infer_data_kind
from pptx_report.model import (
    AppendixContent, ChartPageContent, ChartSpec, ChartType, CoverContent, DataFact, DataKind,
    ExecutiveFinding, ExecutiveSummaryContent, FindingsOverviewContent,
    FunnelAnalysisContent, FunnelStage, KeyFindingContent, OpportunityItem,
    OpportunityMatrixContent, RecommendationContent, RecommendationItem,
    ReportSpec, ResearchOverviewContent, SectionDividerContent, Series,
    TableData, TocContent,
)
from pptx_report.pages import (
    build_appendix, build_exec_summary, build_section_divider, build_toc,
)
from pptx_report.renderer import ReportRenderer
from pptx_report.theme import Theme
from pptx_report.report_templates import select_report_template
from pptx_report.wizard import _build_toc, _enhance_report_pages


def check_facts():
    percentage = {
        "code": "Q1", "title": "满意度", "data_kind": "percentage",
        "categories": ["非常满意", "比较满意", "不太满意", "很不满意"],
        "segments": ["Total", "年轻用户"],
        "data": {"Total": [.30, .35, .20, .15], "年轻用户": [.50, .35, .10, .05]},
        "base": {"Total": 400, "年轻用户": 25}, "ordered_scale": True,
    }
    mean = {
        "code": "Q2", "title": "平均评分", "data_kind": "mean",
        "categories": ["均值"], "segments": ["Total", "老用户"],
        "data": {"Total": [7.2], "老用户": [8.0]}, "base": {"Total": 400, "老用户": 80},
    }
    related = {**percentage, "code": "Q3", "title": "推荐意愿", "data": {"Total": [.30, .35, .20, .15], "年轻用户": [.45, .35, .10, .10]}}
    trend = {"code": "Q4", "title": "年度趋势", "data_kind": "percentage", "categories": ["2024", "2025", "2026"], "segments": ["Total"], "data": {"Total": [.40, .50, .55]}, "base": {"Total": 400}, "trend_ordered": True}
    facts = extract_data_facts([percentage, mean, related, trend])
    gap = next(f for f in facts if f.fact_type == "segment_gap" and f.category == "非常满意")
    assert round(gap.gap_pp, 1) == 20.0 and gap.significant is None
    assert any(f.fact_type == "low_base_warning" and f.segment == "年轻用户" for f in facts)
    assert any(f.fact_type == "top2box" and round(f.value, 1) == 65.0 for f in facts)
    assert any(f.fact_type == "bottom2box" and round(f.value, 1) == 35.0 for f in facts)
    mean_gap = next(f for f in facts if f.fact_type == "mean_difference" and f.question_id == "Q2")
    assert mean_gap.gap_pp is None and mean_gap.value == 8.0
    assert infer_data_kind(mean) == "mean"
    assert any(f.fact_type == "outlier" for f in facts)
    assert any(f.fact_type == "cross_question_consistency" for f in facts)
    assert any(f.fact_type == "trend_change" and round(f.gap_pp, 1) == 15.0 for f in facts)
    funnel = build_funnel_facts([("知晓", 80), ("考虑", 50), ("购买", 20)])
    assert [round(abs(f.gap_pp), 1) for f in funnel] == [30.0, 30.0]
    return facts


def check_capacity_and_briefs(facts):
    pages, audit = split_preserving_order(list(range(23)), max_items=12, source_slide_id="s1")
    assert [len(page) for page in pages] == [12, 11]
    assert [item for page in pages for item in page] == list(range(23))
    assert audit.input_blocks == audit.rendered_blocks == 23
    assert audit.truncated_blocks == 0 and not audit.removed_content
    page_plan = [
        {"slide_id": "a", "title": "用户是谁", "chapter": "用户画像", "questions": [{"code": "Q1"}]},
        {"slide_id": "b", "title": "差异在哪里", "chapter": "消费行为", "questions": [{"code": "Q2"}]},
    ]
    briefs = build_slide_briefs(page_plan, [f.to_dict() for f in facts])
    assert len(briefs) == 2
    assert all(brief.question_answered and brief.claim for brief in briefs)
    assert briefs[0].relationship_to_next and briefs[1].relationship_to_previous
    assert briefs[0].claim != briefs[1].claim


def check_templates_and_chart_semantics():
    request = {"slide_type": "key_finding", "chart_count": 1, "category_count": 4,
               "segment_count": 2, "density": "medium", "importance": "high"}
    first, _ = select_report_template(request)
    usage = Counter({first.template_id: 3})
    second, _ = select_report_template(request, previous_template_id=first.template_id, usage=usage)
    assert first.template_id != second.template_id
    for kind, unit in (("percentage", "%"), ("mean", "分"), ("currency", "元"),
                       ("count", "人"), ("score", "分"), ("nps", "")):
        spec = ChartSpec("指标", ChartType.BAR, ["A", "B"], [Series("Total", [1, 2])],
                         data_kind=DataKind(kind), unit=unit)
        restored = ChartSpec.from_dict(spec.to_dict())
        assert restored.data_kind == DataKind(kind) and restored.unit == unit


def check_ai_semantic_overrides():
    chart = ChartSpec(
        "Satisfaction",
        ChartType.BAR,
        ["Satisfied", "Not satisfied"],
        [Series("Total", [65, 35])],
        evidence_question_ids=["Q1"],
        evidence_fact_ids=["F1"],
        source_references=["Q1.Satisfaction"],
    )
    page = ChartPageContent("AI finding title", charts=[chart], data_source="Q1.Satisfaction")
    fact = DataFact(
        fact_id="F1",
        fact_type="top_rank",
        question_id="Q1",
        metric_name="percentage",
        category="Satisfied",
        value=65,
        source_reference="Q1.Satisfaction",
    )
    questions = [{"code": "Q1", "title": "Satisfaction", "categories": ["Satisfied", "Not satisfied"], "segments": ["Total"], "data": {"Total": [65, 35]}, "base": {"Total": 400}}]
    page_config = {"pages": [{
        "chapter": "蓝图章节",
        "business_implication": "Prioritize the service recovery journey",
        "evidence_fact_ids": ["F1", "invented"],
        "evidence_question_ids": ["Q1", "Q999"],
        "slide_brief": {"locked": True},
    }]}
    enhanced, briefs = _enhance_report_pages(
        [page], questions, [fact], [], "source.xlsx", page_config=page_config
    )
    assert briefs[0].claim == "AI finding title"
    assert briefs[0].business_implication == "Prioritize the service recovery journey"
    assert briefs[0].evidence_fact_ids == ["F1"]
    assert briefs[0].evidence_question_ids == ["Q1"]
    assert briefs[0].locked is True
    assert page.chapter == "蓝图章节"
    assert isinstance(enhanced[1], SectionDividerContent)
    assert enhanced[1].chapter == "蓝图章节"

def check_template_structure_sections():
    chart = ChartSpec(
        "Concept acceptance", ChartType.BAR, ["Like", "Dislike"],
        [Series("Total", [70, 30])], evidence_question_ids=["Q1"],
        evidence_fact_ids=["F1"], source_references=["Q1.Concept acceptance"],
    )
    page = ChartPageContent("Concept acceptance is strong", charts=[chart])
    fact = DataFact(
        fact_id="F1", fact_type="top_rank", question_id="Q1",
        metric_name="percentage", category="Like", value=70,
        source_reference="Q1.Concept acceptance",
    )
    finding = ExecutiveFinding(
        "Concept is accepted", "70% select Like", ["F1"],
        "Prioritize conversion", "high", ["Q1"], ["Q1.Concept acceptance"],
    )
    questions = [{
        "code": "Q1", "title": "Concept acceptance", "categories": ["Like", "Dislike"],
        "segments": ["Total"], "data": {"Total": [70, 30]}, "base": {"Total": 400},
    }]
    config = {
        "template_structure_reused": True,
        "template_report_structure": {"sections": [
            {"number": "01", "title": "项目概述", "topics": ["项目背景", "样本说明"]},
            {"number": "02", "title": "主要研究发现", "topics": ["产品概念测试结果", "目标用户画像"]},
            {"number": "03", "title": "结论与建议", "topics": ["核心结论", "行动建议"]},
        ]},
        "pages": [{"chapter": "主要研究发现"}],
    }
    enhanced, _ = _enhance_report_pages(
        [page], questions, [fact], [finding], "source.xlsx", page_config=config
    )
    section_pages = [item for item in enhanced if isinstance(item, SectionDividerContent)]
    assert [(item.chapter, item.title) for item in section_pages] == [
        ("01", "项目概述"), ("02", "主要研究发现"), ("03", "结论与建议")
    ]
    assert section_pages[1].subtitle == "产品概念测试结果 | 目标用户画像"

def check_template_subsection_sections():
    chart = ChartSpec(
        "Concept acceptance", ChartType.BAR, ["Like", "Dislike"],
        [Series("Total", [70, 30])], evidence_question_ids=["Q1"],
        evidence_fact_ids=["F1"], source_references=["Q1.Concept acceptance"],
    )
    pages = [
        ChartPageContent("Concept acceptance is strong", charts=[chart]),
        ChartPageContent("Target users have distinct needs", charts=[chart]),
        ChartPageContent("The purchase journey has key friction", charts=[chart]),
    ]
    fact = DataFact(
        fact_id="F1", fact_type="top_rank", question_id="Q1",
        metric_name="percentage", category="Like", value=70,
        source_reference="Q1.Concept acceptance",
    )
    finding = ExecutiveFinding(
        "Concept is accepted", "70% select Like", ["F1"],
        "Prioritize conversion", "high", ["Q1"], ["Q1.Concept acceptance"],
    )
    questions = [{
        "code": "Q1", "title": "Concept acceptance", "categories": ["Like", "Dislike"],
        "segments": ["Total"], "data": {"Total": [70, 30]}, "base": {"Total": 400},
    }]
    config = {
        "template_structure_reused": True,
        "template_report_structure": {"sections": [
            {"number": "01", "title": "Project overview", "topics": ["Background"]},
            {
                "number": "02", "title": "Main findings", "topics": ["Concept", "Audience", "Journey"],
                "subsections": [
                    {"number": "2.1", "title": "Concept test", "topics": ["Acceptance", "Preference"]},
                    {"number": "2.2", "title": "Target audience", "topics": ["Profile", "Behavior"]},
                    {"number": "2.3", "title": "Purchase journey", "topics": ["Awareness", "Consideration"]},
                ],
            },
            {"number": "03", "title": "Conclusions", "topics": ["Conclusions", "Actions"]},
        ]},
        "pages": [
            {"chapter": "Concept test"},
            {"chapter": "Target audience"},
            {"chapter": "Purchase journey"},
        ],
    }
    enhanced, _ = _enhance_report_pages(
        pages, questions, [fact], [finding], "source.xlsx", page_config=config
    )
    section_pages = [item for item in enhanced if isinstance(item, SectionDividerContent)]
    assert [(item.chapter, item.title) for item in section_pages] == [
        ("01", "Project overview"),
        ("02", "Main findings"),
        ("2.1", "Concept test"),
        ("2.2", "Target audience"),
        ("2.3", "Purchase journey"),
        ("03", "Conclusions"),
    ]
    assert section_pages[2].subtitle == "Acceptance | Preference"
    assert section_pages[4].subtitle == "Awareness | Consideration"
    hierarchy_output = os.environ.get("TEMPLATE_HIERARCHY_OUTPUT")
    if hierarchy_output:
        spec = ReportSpec(
            CoverContent("Template hierarchy QA", "SurveyKit", "2026-07-29"),
            TocContent([]),
            ExecutiveSummaryContent(findings=[finding]),
            enhanced,
        )
        renderer = ReportRenderer()
        renderer.render(spec, hierarchy_output)
        assert renderer.last_qa and renderer.last_qa["ok"], renderer.last_qa

def check_rendered_page_families():
    finding = ExecutiveFinding("核心用户更关注体验", "年轻用户评价更集中。", ["Q1__top_rank__001"],
                               "优先优化关键触点", "high", ["Q1"], ["Q1.满意度"])
    chart = ChartSpec("满意度", ChartType.BAR, ["满意", "不满意"],
                      [Series("Total", [65, 35])], data_kind=DataKind.PERCENTAGE,
                      evidence_question_ids=["Q1"], evidence_fact_ids=["Q1__top_rank__001"],
                      source_references=["Q1.满意度"])
    pages = [
        ResearchOverviewContent(sample_size=400, question_count=2, segment_count=2,
                                source_references=["交叉表"]),
        SectionDividerContent("用户与行为", "主要研究发现", "从画像进入行为诊断"),
        FindingsOverviewContent(findings=[finding]),
        KeyFindingContent("核心发现", finding, [chart], "数据来源：Q1.满意度"),
        FunnelAnalysisContent("转化漏斗", [FunnelStage("知晓", 80), FunnelStage("考虑", 50), FunnelStage("购买", 20)], ["考虑到购买流失最大"], "数据来源：Q3"),
        OpportunityMatrixContent("机会优先级", [OpportunityItem("产品体验", 85, 45, "优先改善"), OpportunityItem("品牌认知", 70, 75, "持续保持")], "数据来源：Q4"),
        RecommendationContent(recommendations=[RecommendationItem("改善核心体验", "对应产品体验短板", "high", ["Q1__top_rank__001"])])
    ]
    spec = ReportSpec(CoverContent("定量研究报告", "测试客户", "2026-07-23"), TocContent([]),
                      ExecutiveSummaryContent(findings=[finding]), pages)
    with TemporaryDirectory() as temp:
        output = Path(temp) / "semantic-report.pptx"
        renderer = ReportRenderer()
        renderer.render(spec, str(output))
        assert renderer.last_qa and renderer.last_qa["ok"], renderer.last_qa
        prs = Presentation(str(output))
        assert len(prs.slides) == 10
        texts = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False))
        assert "核心发现" in texts and "行动建议" in texts and "数据来源：Q4" in texts
        assert "高差异 · 待改善" in texts
        assert "当前表现（0–100）→" in texts and "差异强度指数（0–100）↑" in texts
        assert "表现基准 50" in texts and "差异基准 50" in texts
        with ZipFile(output) as archive:
            xml = "".join(archive.read(name).decode("utf-8", "ignore") for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml"))
        assert "NaN" not in xml
        if os.environ.get("REPORT_SEMANTICS_OUTPUT"):
            target = Path(os.environ["REPORT_SEMANTICS_OUTPUT"])
            target.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(output, target)


def check_optional_opportunity_matrix_and_conclusion():
    chart = ChartSpec(
        "Purchase intent", ChartType.BAR,
        ["Definitely buy", "Probably buy", "Unsure", "Probably not", "Definitely not"],
        [Series("Total", [20, 30, 25, 15, 10]), Series("Target", [40, 32, 16, 8, 4])],
        evidence_question_ids=["Q1"],
    )
    page = ChartPageContent("Purchase intent differs", charts=[chart])
    questions = [
        {
            "code": "Q1", "title": "Purchase intent",
            "categories": list(chart.categories),
            "segments": ["Total", "Target"],
            "data": {"Total": [20, 30, 25, 15, 10], "Target": [40, 32, 16, 8, 4]},
            "base": {"Total": 500, "Target": 150},
        },
        {
            "code": "Q2", "title": "Unrelated channel",
            "categories": ["A", "B", "C", "D"],
            "segments": ["Total", "Target"],
            "data": {"Total": [25, 25, 25, 25], "Target": [28, 24, 24, 24]},
            "base": {"Total": 500, "Target": 150},
        },
    ]
    facts = []
    for index, (category, value, gap) in enumerate(zip(
        chart.categories, [40, 32, 16, 8, 4], [20, 2, -9, -7, -6]
    )):
        facts.append(DataFact(
            fact_id=f"Q1_gap_{index}", fact_type="segment_gap", question_id="Q1",
            metric_name="percentage", category=category, segment="Target",
            value=value, benchmark_value=value - gap, gap_pp=gap,
        ))
    for index, category in enumerate(["A", "B", "C", "D"]):
        facts.append(DataFact(
            fact_id=f"Q2_gap_{index}", fact_type="segment_gap", question_id="Q2",
            metric_name="percentage", category=category, segment="Target",
            value=24 + index, benchmark_value=25, gap_pp=index - 1,
        ))
    finding = ExecutiveFinding(
        title="Purchase intent is polarized", description="Target users differ from total",
        evidence_fact_ids=["Q1_gap_0"], action_implication="Prioritize high-intent conversion",
        evidence_question_ids=["Q1"],
    )
    enhanced, _ = _enhance_report_pages(
        [page], questions, facts, [finding], "source.xlsx"
    )
    assert not any(isinstance(item, OpportunityMatrixContent) for item in enhanced)
    conclusion = next(item for item in enhanced if isinstance(item, FindingsOverviewContent))
    assert conclusion.title == "核心结论"
    assert sum(isinstance(item, FindingsOverviewContent) for item in enhanced) == 1

    with_matrix, _ = _enhance_report_pages(
        [page], questions, facts, [finding], "source.xlsx",
        page_config={"include_opportunity_matrix": True},
    )
    matrix = next(item for item in with_matrix if isinstance(item, OpportunityMatrixContent))
    matrix_fact_ids = [fact_id for item in matrix.opportunities for fact_id in item.fact_ids]
    assert len(matrix.opportunities) >= 4
    assert all(fact_id.startswith("Q1_gap_") for fact_id in matrix_fact_ids)
    assert matrix.title.startswith("Purchase intent")
    overview = next(item for item in enhanced if isinstance(item, ResearchOverviewContent))
    assert overview.segment_count == 2


def check_toc_follows_final_page_order():
    first = ChartPageContent("Concept result", charts=[ChartSpec(
        "Concept", ChartType.BAR, ["A", "B"], [Series("Total", [60, 40])],
        evidence_question_ids=["C1", "C2"],
    )], chapter="概念测试")
    second = ChartPageContent("Audience", charts=[ChartSpec(
        "Audience", ChartType.BAR, ["A", "B"], [Series("Total", [55, 45])],
        evidence_question_ids=["D1"],
    )], chapter="用户画像")
    toc = _build_toc([second, first], "source.xlsx")
    assert toc.sections == [
        "一、项目概述",
        "二、主要研究发现",
        "  用户画像（1题）",
        "  概念测试（2题）",
        "三、结论与建议",
    ]

def check_summary_toc_appendix_layouts():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    theme = Theme()

    toc_slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_toc(
        toc_slide,
        TocContent(["一、项目概述", "二、主要研究发现", "  用户画像", "  消费行为", "三、结论与建议"]),
        theme,
        (prs.slide_width, prs.slide_height),
    )
    toc_text = "\n".join(
        shape.text for shape in toc_slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert "02.1" in toc_text and "02.2" in toc_text
    assert "04" not in toc_text

    section_slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_section_divider(
        section_slide,
        SectionDividerContent(
            title="用户画像与场景分化",
            chapter="用户画像与场景分化",
            subtitle="还需要关注什么",
        ),
        theme,
        (prs.slide_width, prs.slide_height),
    )
    section_text = [
        shape.text for shape in section_slide.shapes
        if getattr(shape, "has_text_frame", False)
    ]
    assert section_text.count("用户画像与场景分化") == 1
    assert "还需要关注什么" not in section_text

    summary_slide = prs.slides.add_slide(prs.slide_layouts[6])
    duplicate = "Prioritize differentiated communication"
    build_exec_summary(
        summary_slide,
        ExecutiveSummaryContent(findings=[
            ExecutiveFinding(
                title="Target users need distinct messages",
                description=duplicate,
                action_implication=duplicate,
            )
        ]),
        theme,
        (prs.slide_width, prs.slide_height),
    )
    summary_text = "\n".join(
        shape.text for shape in summary_slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert summary_text.count(duplicate) == 1

    appendix_slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_appendix(
        appendix_slide,
        AppendixContent(
            title="Appendix",
            table=TableData(
                headers=["Question", "Description", "Options"],
                rows=[["Q1", "Single appendix row", 2]],
            ),
        ),
        theme,
        (prs.slide_width, prs.slide_height),
    )
    table_shape = next(shape for shape in appendix_slide.shapes if getattr(shape, "has_table", False))
    assert table_shape.height < Inches(1.2)
    assert table_shape.table.rows[1].height <= Inches(.5)
    layout_output = os.environ.get("REPORT_LAYOUT_QA_OUTPUT")
    if layout_output:
        target = Path(layout_output)
        target.parent.mkdir(parents=True, exist_ok=True)
        prs.save(target)

def main():
    facts = check_facts()
    check_capacity_and_briefs(facts)
    check_templates_and_chart_semantics()
    check_ai_semantic_overrides()
    check_template_structure_sections()
    check_template_subsection_sections()
    check_rendered_page_families()
    check_optional_opportunity_matrix_and_conclusion()
    check_toc_follows_final_page_order()
    check_summary_toc_appendix_layouts()
    print("report semantics smoke: ok")


if __name__ == "__main__":
    main()
