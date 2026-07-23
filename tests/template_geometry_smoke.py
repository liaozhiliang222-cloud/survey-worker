"""Regression check for deep-template table/chart alignment."""

from pathlib import Path
import sys

import pandas as pd
from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptx_report.model import ChartPageContent, ChartSpec, LayoutType, MultiGroupBarPageContent  # noqa: E402
from pptx_report.pages import build_chart_page, build_multi_group_bar_page  # noqa: E402
from pptx_report.renderer import ReportRenderer  # noqa: E402
from pptx_report.template import build_template_mapping  # noqa: E402
from pptx_report.theme import Theme  # noqa: E402


def main() -> None:
    template_prs = Presentation()
    template_prs.slides.add_slide(template_prs.slide_layouts[0])
    template_slide = template_prs.slides.add_slide(template_prs.slide_layouts[1])
    template_slide.shapes.title.text = "研究发现标题"
    template_slide.shapes.title.top = int(template_prs.slide_height * 0.02)
    template_slide.shapes.title.height = int(template_prs.slide_height * 0.12)
    body = template_slide.placeholders[1]
    body.text = "标题下方的洞察正文，不应被识别为标题区域。"
    body.top = int(template_prs.slide_height * 0.15)
    body.height = int(template_prs.slide_height * 0.16)
    template_mapping = build_template_mapping(template_prs)
    override_mapping = build_template_mapping(
        template_prs, role_overrides={"chart": 2, "summary": 2}
    )
    for role in ("chart", "summary"):
        selected = override_mapping["roles"][role]
        if selected["slide_index"] != 1 or not selected.get("user_confirmed"):
            raise AssertionError(f"Template role override did not select slide 2 for {role}: {selected}")
        if selected["confidence"] != 1.0:
            raise AssertionError("User-confirmed template roles must have full confidence")
    content_zones = template_mapping["roles"]["content"]["zones"]
    if content_zones["title"]["h"] > 0.13:
        raise AssertionError(f"Title zone incorrectly includes body placeholder: {content_zones}")
    if content_zones["content"]["y"] > 0.16:
        raise AssertionError(f"Content zone was pushed down by title detection: {content_zones}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    segments = ["Total", "都市中产", "都市蓝领", "都市家庭", "都市Z世代", "小镇中年", "小镇青年"]
    options = ["北京", "河北", "广东", "山东", "江苏", "天津", "上海", "辽宁", "四川", "湖北", "河南"]
    rows = []
    for row_index, option in enumerate(options):
        row = {"选项": option}
        for segment_index, segment in enumerate(segments):
            row[segment] = ((row_index * 7 + segment_index * 11) % 48 + 2) / 100
        rows.append(row)

    page = MultiGroupBarPageContent(
        title="[上海] 占比最高（44.0%）",
        groups_data=[{"title": "省份", "data": pd.DataFrame(rows)}],
        segments=segments,
        insights=["测试洞察一", "2. 测试洞察二"],
        data_source="数据来源：测试数据",
    )
    build_multi_group_bar_page(
        slide,
        page,
        Theme(),
        (prs.slide_width, prs.slide_height),
    )

    renderer = ReportRenderer()
    renderer._prs = prs
    renderer._slide_w = prs.slide_width
    renderer._slide_h = prs.slide_height
    renderer.template_path = "synthetic-template.pptx"
    renderer._template_mapping = {
        "roles": {
            "matrix": {
                "has_title_divider": True,
                "zones": {
                    "title": {"x": 0.03, "y": 0.01, "w": 0.68, "h": 0.12},
                    "content": {"x": 0.041, "y": 0.17, "w": 0.95, "h": 0.76},
                    "footer": {"x": 0.03, "y": 0.94, "w": 0.90, "h": 0.04},
                }
            }
        }
    }
    before_cleanup = len(slide.shapes)
    renderer._remove_duplicate_title_divider(slide, 0, "matrix")
    if len(slide.shapes) != before_cleanup - 1:
        raise AssertionError("Expected the renderer title divider to be removed")
    renderer._apply_template_geometry(slide, 0, "matrix")

    insight_shape = next(
        shape for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and "测试洞察一" in shape.text
    )
    for paragraph in insight_shape.text_frame.paragraphs:
        paragraph_xml = paragraph._p.xml
        if "buChar" not in paragraph_xml or "buAutoNum" in paragraph_xml:
            raise AssertionError("Expected insight paragraphs to use round bullets")
        if paragraph.text.lstrip().startswith(("1.", "2.", "3.")):
            raise AssertionError("Insight text still contains numeric segment markers")

    table_shape = next(shape for shape in slide.shapes if getattr(shape, "has_table", False))
    chart_shapes = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]
    if not chart_shapes:
        raise AssertionError("Expected overlaid bar charts in regression slide")

    data_columns_left = int(table_shape.left) + sum(
        int(table_shape.table.columns[index].width) for index in (0, 1)
    )
    first_chart_left = min(int(shape.left) for shape in chart_shapes)
    alignment_gap = first_chart_left - data_columns_left
    if not 0 <= alignment_gap <= Inches(0.20):
        raise AssertionError(
            f"Table/chart columns are misaligned by {alignment_gap / 914400:.3f} inches"
        )
    if abs(int(table_shape.width) - sum(int(column.width) for column in table_shape.table.columns)) > 2:
        raise AssertionError("Table frame width and internal column grid diverged")
    if abs(int(table_shape.height) - sum(int(row.height) for row in table_shape.table.rows)) > 2:
        raise AssertionError("Table frame height and internal row grid diverged")

    chart_slide = prs.slides.add_slide(prs.slide_layouts[6])
    doughnut_charts = [
        ChartSpec.doughnut(
            f"Demographic question {index + 1}",
            ["Option A", "Option B", "Option C"],
            [36.2, 42.1, 21.7],
            insight=f"Insight {index + 1}",
        )
        for index in range(4)
    ]
    chart_page = ChartPageContent(
        title="Four-chart page after question editing",
        charts=doughnut_charts,
        layout=LayoutType.DASHBOARD,
        insights=["Insight A", "Insight B", "Insight C"],
        data_source="Regression data source",
    )
    build_chart_page(
        chart_slide,
        chart_page,
        Theme(),
        (prs.slide_width, prs.slide_height),
    )
    dashboard_charts = [
        shape for shape in chart_slide.shapes if getattr(shape, "has_chart", False)
    ]
    if len(dashboard_charts) != 4:
        raise AssertionError("Expected four charts in edited-page regression slide")
    for shape in dashboard_charts:
        if "manualLayout" not in shape.chart._chartSpace.xml:
            raise AssertionError("Wide doughnut chart is missing manual plot-area geometry")
        if shape.chart.legend.include_in_layout is not True:
            raise AssertionError("Doughnut legend must reserve its own layout space")
        if "manualLayout" not in shape.chart._chartSpace.chart.legend.xml:
            raise AssertionError("Doughnut legend is missing fixed bottom geometry")
        label_size = shape.chart.plots[0].data_labels.font.size
        if label_size is None or label_size.pt > 7.1:
            raise AssertionError("Compact doughnut labels must fit inside the ring")
        if 'holeSize val="25"' not in shape.chart._chartSpace.xml:
            raise AssertionError("Compact doughnut hole must leave room for labels")
    min_chart_height = min(int(shape.height) for shape in dashboard_charts)
    if min_chart_height < Inches(1.85):
        raise AssertionError(
            f"Edited-page charts are too small: {min_chart_height / 914400:.3f} inches"
        )
    row_tops = sorted({int(shape.top) for shape in dashboard_charts})
    if len(row_tops) != 2:
        raise AssertionError(f"Expected a 2x2 chart grid, got row tops {row_tops}")
    first_row = [shape for shape in dashboard_charts if int(shape.top) == row_tops[0]]
    row_gap = row_tops[1] - max(int(shape.top) + int(shape.height) for shape in first_row)
    if row_gap > Inches(0.70):
        raise AssertionError(
            f"Insight space was deducted per row; row gap is {row_gap / 914400:.3f} inches"
        )
    output = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("qa_template_alignment_fixed.pptx")
    prs.save(output)
    print(f"Template alignment smoke passed: {output}")


if __name__ == "__main__":
    main()
