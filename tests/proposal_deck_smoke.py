import json
import os
import sys
import tempfile
import zipfile
from copy import deepcopy
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from fastapi.testclient import TestClient
from pptx import Presentation

from deploy import aliyun_api
from pptx_report.proposal_deck import render_proposal_deck


def build_deck():
    specs = [
        ("cover", "智能宠物饮水机用户研究方案", "cover_product_hero", []),
        ("project_context", "饮水管理升级需求与产品决策风险同时出现", "context_tension_map", [
            ("宠物健康管理精细化", "现象：饮水量、饮水质量和异常提醒成为需求｜业务含义：用户开始为健康证据付费｜研究响应：验证普遍性、迫切度与支付关联"),
            ("自动循环不等于真正省心", "现象：清洁、滤芯、噪音与宠物适应仍是负担｜业务含义：维护成本可能抵消智能价值｜研究响应：还原持续使用与流失节点"),
            ("健康价值必须转化为持续体验", "现象：卫生、静音、易清洁与提醒需要协同｜业务含义：单一技术卖点不足以支撑升级｜研究响应：识别价值组合与MVP优先级"),
            ("未经验证会放大上市风险", "现象：概念、人群、价格与体验边界尚未锁定｜业务含义：错误定位会推高获客成本｜研究响应：形成推进、优化或停止依据"),
        ]),
        ("business_decisions", "研究结束后需要支持五项核心业务决策", "decision_tree", [
            ("概念是否推进", "验证理解度、相关性、独特性与购买意愿"), ("卖点如何排序", "识别基础价值、差异化价值和弱化信息"),
            ("首发面向谁", "锁定高需求、高接受度且可触达的机会人群"), ("价格是否成立", "结合价格敏感度与购买概率判断商业空间"),
            ("体验如何优化", "聚焦清洁、噪音、宠物适应与持续使用障碍"),
        ]),
        ("decision_evidence", "每项业务决策都对应证据、判断方式与行动", "evidence_threshold_matrix", [
            ("概念是否推进", "理解度、相关性、可信度、购买意愿｜概念比较＋内部标准｜推进、优化或停止"),
            ("卖点如何排序", "偏好、差异性、选择驱动｜MaxDiff＋驱动分析｜确定主次卖点"),
            ("首发人群是谁", "需求、意愿、支付与可触达性｜细分＋机会矩阵｜确定目标人群"),
            ("价格是否成立", "价格接受与购买概率｜PSM＋Gabor-Granger｜确定价格区间"),
        ]),
        ("research_path", "定性探索与定量验证逐步降低上市风险", "dual_track_research_flow", [
            ("业务输入", "确认假设、刺激物和核心决策问题"), ("定性探索", "输出用户语言、核心张力和问卷输入"),
            ("概念优化", "形成信息量一致的可测版本"), ("定量验证", "量化规模、差异、价格和优先级"), ("策略输出", "形成产品、人群与营销行动"),
        ]),
        ("qualitative_design", "定性研究把真实饮水场景转化为可测试语言", "qualitative_design_canvas", [
            ("新手与成熟养宠家庭", "兼顾猫犬、多宠、设备经验和消费层级"), ("居家深访＋任务演示", "建议8—12人，每次60—90分钟"),
            ("现状—痛点—概念", "还原补水、拆洗、滤芯和宠物适应过程"), ("功能卡与价格情景", "统一信息量和呈现顺序"),
            ("定性小结与问卷输入", "输出自然语言、选项、概念优化点和分群假设"),
        ]),
        ("sample_design", "样本架构兼顾概念比较与重点人群差异", "quantitative_sample_architecture", [
            ("总体样本 N=700", "城市养宠家庭的主要购买决策者；执行前验证可达性"), ("概念A N=350", "猫家庭N=220；犬家庭N=130"),
            ("概念B N=350", "猫家庭N=220；犬家庭N=130"), ("重点加样", "高端宠物消费与智能设备高接受人群"),
            ("交叉分析", "猫犬×城市×设备经验×消费层级"), ("分析限制", "最低分析单元建议不低于N=80"),
        ]),
        ("questionnaire_design", "七个问卷模块沿用户决策过程逐步推进", "questionnaire_decision_journey", [
            ("使用现状", "指标：设备、角色、频率、维护投入｜决策：定义行为基线｜方法：描述统计｜输出：典型管理路径"),
            ("痛点和现有方案", "指标：清洁、噪音、卫生焦虑、适应｜决策：确认需求真实性｜方法：痛点矩阵｜输出：高频高痛机会"),
            ("概念理解", "指标：理解、相关、独特、可信｜决策：判断概念去留｜方法：概念漏斗｜输出：概念优化方向"),
            ("卖点取舍", "指标：健康、清洁、静音、提醒、外观｜决策：确定沟通优先级｜方法：MaxDiff｜输出：主次卖点"),
            ("购买意愿", "指标：Top2Box、情景、障碍｜决策：估计转化潜力｜方法：驱动分析｜输出：转化杠杆"),
            ("价格测试", "指标：价格接受、购买概率｜决策：判断商业空间｜方法：PSM＋Gabor-Granger｜输出：建议区间"),
            ("人群和触点", "指标：需求、态度、消费、渠道｜决策：识别首发人群｜方法：聚类＋机会矩阵｜输出：触达策略"),
        ]),
        ("report_example_concept", "报告将识别最具转化潜力的概念与卖点", "concept_funnel_maxdiff_example", [
            ("示例解读", "健康活水与易清洁构成基础价值；不代表实际研究结论"),
        ]),
        ("report_example_pricing", "报告将界定价格空间并识别首发机会人群", "pricing_segment_example", [
            ("示例解读", "价格上升伴随购买意愿递减；不代表实际研究结论"),
        ]),
        ("decision_outputs", "研究指标最终汇聚为七类可执行上市结论", "decision_output_map", [
            ("需求真实性", "判断痛点是否真实且值得解决"), ("最优概念方向", "比较理解、吸引、独特与可信表现"),
            ("核心卖点组合", "区分基础、差异化和低贡献信息"), ("MVP功能优先级", "明确必须、增强和延后功能"),
            ("建议价格区间", "结合敏感度和购买概率判断"), ("首发目标人群", "识别高需求高接受机会人群"),
            ("障碍与行动建议", "转成产品和营销下一步动作"),
        ]),
        ("timeline", "两周内并行推进研究执行、客户确认与交付", "timeline_gantt_risk", [
            ("交付", "研究方案、定性小结、数据底表与最终报告"), ("依赖", "概念刺激物与招募条件按时确认"), ("风险", "招募、素材变更与确认延迟需前置管理"),
        ]),
    ]
    dataset = {
        "dataset_id": "pet_water_dispenser_demo_v2", "project_type": "concept_test", "data_status": "illustrative",
        "usable_for_decision": False, "example_label": "REPORT EXAMPLE｜报告输出示例",
        "disclaimer": "本页数据为AI生成的示例数据，仅用于展示未来报告的分析形式，不代表实际研究结果。",
        "metrics": {"concept_understanding": 82, "concept_relevance": 68, "concept_uniqueness": 61, "concept_credibility": 55, "purchase_intention_t2b": 39},
        "selling_point_scores": {"健康活水": 31, "易清洁": 26, "静音运行": 21, "智能提醒": 13, "外观设计": 9},
        "segment_purchase_intention": {"高端养宠人群": 58, "智能设备尝鲜者": 51, "普通饮水机用户": 37, "普通水碗用户": 24},
        "pricing": {"acceptable_low": 299, "optimal_price": 399, "acceptable_high": 499,
                    "purchase_curve": [{"price": 299, "intention": 57}, {"price": 399, "intention": 39}, {"price": 499, "intention": 25}, {"price": 599, "intention": 14}]},
    }
    key_messages = {
        "project_context": "识别产品上市前最需要降低的四类决策风险",
        "business_decisions": "覆盖概念去留、卖点排序、人群、价格与体验优化",
        "decision_evidence": "通过相对比较、历史基准或客户内部标准进行判断",
        "research_path": "前一阶段输出直接成为后一阶段输入，而不是方法堆叠",
        "qualitative_design": "从真实使用场景提炼用户语言、概念优化点和定量假设",
        "sample_design": "兼顾总体判断、概念比较、重点加样和最低分析单元",
        "questionnaire_design": "每个模块同时连接指标、业务决策、分析方法与阶段输出",
        "report_example_concept": "以概念漏斗和MaxDiff展示未来报告的核心输出形式",
        "report_example_pricing": "价格曲线与人群机会共用同一套示例数据口径",
        "decision_outputs": "把需求、概念、卖点、功能、价格、人群和障碍连接到行动",
        "timeline": "定性与问卷设计并行推进，以支持两周内完成交付",
    }
    slides = []
    for index, (slide_type, title, visual, content_spec) in enumerate(specs, 1):
        slides.append({
            "id": f"slide_{index:02d}", "order": index,
            "slide_type": slide_type,
            "title": title,
            "subtitle": "从购买决策到持续使用的完整验证路径" if index == 1 else "",
            "key_message": key_messages.get(slide_type, ""),
            "slide_question": f"第{index}页回答的唯一研究问题", "unique_purpose": f"故事线第{index}步",
            "previous_slide_relation": "承接上一页", "next_slide_relation": "引出下一页",
            "visual_type": visual,
            "layout_variant": "v2", "relation_type": "hero" if index == 1 else "data" if index in (9, 10) else "sequence",
            "content_density": "professional", "target_canvas_occupancy": .72,
            "content": [
                {"id": f"item_{index}_{j}", "label": f"模块 {j}", "headline": headline,
                 "description": description, "priority": j,
                 "source_type": "word_proposal"}
                for j, (headline, description) in enumerate(content_spec, 1)
            ],
            "charts": ([{"chart_id": "concept", "chart_type": "funnel", "dataset_id": dataset["dataset_id"]}] if index == 9 else
                       [{"chart_id": "pricing", "chart_type": "line", "dataset_id": dataset["dataset_id"]}] if index == 10 else []),
            "data_status": "illustrative" if index in (9, 10) else "framework_only", "example_output": index in (9, 10),
            "timeline_tasks": ([
                {"task": "研究设计", "start_day": 1, "end_day": 2, "deliverable": "研究方案"}, {"task": "定性招募", "start_day": 1, "end_day": 3, "deliverable": "访问名单"},
                {"task": "定性访问", "start_day": 3, "end_day": 5, "deliverable": "定性小结"}, {"task": "问卷设计", "start_day": 2, "end_day": 5, "deliverable": "问卷终稿"},
                {"task": "编程测试", "start_day": 4, "end_day": 5, "deliverable": "测试链接"}, {"task": "定量回收", "start_day": 6, "end_day": 10, "deliverable": "有效样本"},
                {"task": "数据分析", "start_day": 9, "end_day": 11, "deliverable": "分析底表"}, {"task": "报告输出", "start_day": 11, "end_day": 14, "deliverable": "最终报告"}
            ] if index == 12 else []),
            "source_references": [{"source_type": "word_proposal", "label": "已确认的 Word 方案"}], "locked": index == 3,
            "notes": "口径：测试样稿，不代表真实调研结论。"
        })
        if slide_type == "decision_evidence":
            for entry, (decision, raw_value) in zip(slides[-1]["content"], content_spec):
                parts = raw_value.split("｜")
                entry["label"] = decision
                entry["headline"] = parts[0]
                entry["description"] = "｜".join(parts[1:])
    return {
        "schema_version": "2.0", "deck_id": "deck_pet_water", "project_id": "project_pet_water",
        "title": "AI调研方案PPT-智能宠物饮水机", "purpose": "client_proposal",
        "page_size": 12, "aspect_ratio": "16:9", "theme": "modern_insight_v2", "example_output_mode": "illustrative",
        "illustrative_dataset_id": dataset["dataset_id"], "illustrative_dataset": dataset, "content_density": "professional",
        "target_canvas_occupancy": .72, "min_body_characters": 160, "max_body_characters": 300, "slides": slides,
    }


def main():
    deck = build_deck()
    output = Path(os.environ.get("PROPOSAL_DECK_OUTPUT", ROOT / "tests" / "output" / "proposal-deck-smoke.pptx"))
    output.parent.mkdir(parents=True, exist_ok=True)
    audit = render_proposal_deck(deck, str(output))
    prs = Presentation(str(output))
    assert len(prs.slides) == len(audit["deck"]["slides"])
    assert len(prs.slides) > len(deck["slides"])
    assert round(prs.slide_width / prs.slide_height, 2) == round(16 / 9, 2)
    shape_count = sum(len(slide.shapes) for slide in prs.slides)
    text_shape_count = sum(1 for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False))
    chart_shapes = [shape for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_chart", False)]
    assert shape_count > 180, shape_count
    assert text_shape_count > 110, text_shape_count
    assert len(chart_shapes) == 4, len(chart_shapes)
    assert all(shape.chart.series for shape in chart_shapes)
    report_indices = [i for i, slide in enumerate(audit["deck"]["slides"]) if slide["visual_type"] in {"concept_funnel_maxdiff_example", "pricing_segment_example"}]
    report_slides = [prs.slides[i] for i in report_indices]
    assert all(sum(1 for shape in slide.shapes if getattr(shape, "has_chart", False)) == 2 for slide in report_slides)
    assert all(not any(getattr(shape, "has_table", False) for shape in slide.shapes) for slide in report_slides)
    visible_text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False))
    for raw_field in ["dataset_id", "metrics", "selling_point_scores", "segment_purchase_intention", "purchase_curve", "start_day", "end_day"]:
        assert raw_field not in visible_text
    path_index = next(i for i, slide in enumerate(audit["deck"]["slides"]) if slide["visual_type"] == "dual_track_research_flow")
    path_text = "\n".join(shape.text for shape in prs.slides[path_index].shapes if getattr(shape, "has_text_frame", False))
    assert "研究动作｜做什么" in path_text
    assert "阶段产出｜交付什么" in path_text
    assert "下一阶段输入｜如何被使用" in path_text
    assert "开展用户深访与场景追问" in path_text
    assert "统一概念刺激物信息量，完成可测化改写" in path_text
    assert "标准化概念刺激物与评价维度" in path_text
    gantt_index = next(i for i, slide in enumerate(audit["deck"]["slides"]) if slide["visual_type"] == "timeline_gantt_risk")
    gantt_bar_widths = {
        round(shape.width / 914400, 2)
        for shape in prs.slides[gantt_index].shapes
        if .18 <= shape.height / 914400 <= .3 and shape.width / 914400 > .3 and shape.left / 914400 >= 2.9
    }
    assert len(gantt_bar_widths) >= 3, gantt_bar_widths
    with zipfile.ZipFile(output) as package:
        assert len([name for name in package.namelist() if name.startswith("ppt/charts/chart") and name.endswith(".xml")]) == 4
        assert len([name for name in package.namelist() if name.startswith("ppt/embeddings/")]) == 4
        assert not [name for name in package.namelist() if name.startswith("ppt/media/")]
    split_issues = [issue for issue in audit["issues"] if issue.get("code") == "slide_split_for_capacity"]
    assert split_issues
    assert sum(len(slide["content"]) for slide in audit["deck"]["slides"]) == sum(len(slide.get("content", [])) for slide in deck["slides"])
    render_audit = audit["deck"]["render_audit"]
    assert render_audit["complete"] is True
    assert render_audit["input_blocks"] == render_audit["rendered_blocks"]
    assert render_audit["truncated_blocks"] == 0
    assert render_audit["removed_content"] == []
    assert all(slide["slide_brief"]["slide_id"] == slide["slide_id"] for slide in audit["deck"]["slides"])
    assert all(slide["slide_brief"]["question_answered"] for slide in audit["deck"]["slides"])
    assert all(slide["slide_brief"]["claim"] for slide in audit["deck"]["slides"])
    assert all(slide["slide_brief"]["template_id"] == slide["template_id"] for slide in audit["deck"]["slides"])
    assert any(slide["slide_brief"]["locked"] for slide in audit["deck"]["slides"])
    assert audit["ok"] is True
    assert not [issue for issue in audit["issues"] if issue["level"] == "error"]
    assert deck["illustrative_dataset"]["usable_for_decision"] is False
    assert sum(deck["illustrative_dataset"]["selling_point_scores"].values()) == 100
    assert len({task["start_day"] for task in deck["slides"][11]["timeline_tasks"]}) > 1
    assert all(str(value) not in " ".join(item["headline"] + item["description"] for item in deck["slides"][10]["content"])
               for value in deck["illustrative_dataset"]["metrics"].values())

    with tempfile.TemporaryDirectory() as temp_dir:
        edited = deepcopy(deck)
        edited["illustrative_dataset"]["metrics"]["concept_understanding"] = 80
        edited_path = Path(temp_dir) / "edited.pptx"
        render_proposal_deck(edited, str(edited_path))
        edited_prs = Presentation(str(edited_path))
        edited_charts = [shape.chart for slide in edited_prs.slides for shape in slide.shapes if getattr(shape, "has_chart", False)]
        assert any(80 in list(series.values) for series in edited_charts[0].series)

        framework = deepcopy(deck)
        framework["example_output_mode"] = "framework_only"
        framework["illustrative_dataset"] = None
        for slide in framework["slides"]:
            if slide["example_output"]:
                slide["data_status"] = "framework_only"
        framework_path = Path(temp_dir) / "framework.pptx"
        framework_audit = render_proposal_deck(framework, str(framework_path))
        assert framework_audit["ok"] is True
        framework_prs = Presentation(str(framework_path))
        framework_text = "\n".join(shape.text for slide in framework_prs.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False))
        assert "82%" not in framework_text

    client = TestClient(aliyun_api.app)
    unauthorized = client.post("/api/pptx-report/proposal-deck", json=deck)
    assert unauthorized.status_code == 403
    validated = client.post("/api/pptx-report/proposal-deck/validate", json=deck,
                            headers={"X-Project-Id": deck["project_id"]})
    assert validated.status_code == 200, validated.text
    generated = client.post("/api/pptx-report/proposal-deck", json=deck,
                            headers={"X-Project-Id": deck["project_id"]})
    assert generated.status_code == 200, generated.text
    assert generated.content[:2] == b"PK"

    original_renderer = aliyun_api.render_proposal_deck
    aliyun_api.render_proposal_deck = lambda *_args, **_kwargs: (_ for _ in ()).throw(RuntimeError("disk full"))
    try:
        failed = client.post("/api/pptx-report/proposal-deck", json=deck,
                             headers={"X-Project-Id": deck["project_id"]})
        assert failed.status_code == 500
        assert "PPT" in failed.text
    finally:
        aliyun_api.render_proposal_deck = original_renderer

    print(json.dumps({"output": str(output), "slides": len(prs.slides), "shapes": shape_count,
                      "text_shapes": text_shape_count}, ensure_ascii=False))


if __name__ == "__main__":
    main()
