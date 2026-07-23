"""Multi-layout regression check for uploaded-template header geometry."""

from pathlib import Path
import sys

import pandas as pd
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptx_report.model import (  # noqa: E402
    ChartPageContent,
    ChartSpec,
    LayoutType,
    MultiGroupBarPageContent,
)
from pptx_report.pages import build_chart_page, build_multi_group_bar_page  # noqa: E402
from pptx_report.renderer import ReportRenderer  # noqa: E402
from pptx_report.template import build_template_mapping  # noqa: E402
from pptx_report.theme import Theme  # noqa: E402


def _synthetic_template(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.slides.add_slide(prs.slide_layouts[6])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(
        int(prs.slide_width * 0.04), int(prs.slide_height * 0.044),
        int(prs.slide_width * 0.90), int(prs.slide_height * 0.071),
    )
    title.text_frame.paragraphs[0].add_run().text = "Research finding title"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
    body = slide.shapes.add_textbox(
        int(prs.slide_width * 0.04), int(prs.slide_height * 0.133),
        int(prs.slide_width * 0.90), int(prs.slide_height * 0.107),
    )
    body.text_frame.paragraphs[0].add_run().text = "Insight line one"
    body.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
    data = ChartData()
    data.categories = ["A", "B"]
    data.add_series("Share", [60, 40])
    slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, Inches(1), Inches(2), Inches(4), Inches(3), data
    )
    prs.save(path)


def _chart_pages() -> list[ChartPageContent]:
    doughnut = lambda title: ChartSpec.doughnut(title, ["A", "B", "C"], [60, 25, 15])
    bar = lambda title: ChartSpec.bar(
        title, ["Option A", "Option B", "Option C"], {"Total": [72, 55, 41]}
    )
    return [
        ChartPageContent(
            title="Single chart title",
            charts=[bar("Single comparison")],
            layout=LayoutType.SINGLE,
            insights=["Primary insight for the single-chart page"],
            data_source="Regression source 1",
        ),
        ChartPageContent(
            title=(
                "Long research finding title used to verify that imported template "
                "headers remain clear and do not collide with the divider"
            ),
            charts=[doughnut("Composition"), bar("Preference")],
            layout=LayoutType.DUAL,
            insights=["Dual-page insight one", "Dual-page insight two"],
            data_source="Regression source 2",
        ),
        ChartPageContent(
            title="Three-chart dashboard title",
            charts=[doughnut("Gender"), doughnut("Age"), bar("Household")],
            layout=LayoutType.DASHBOARD,
            insights=["Dashboard insight one", "Dashboard insight two"],
            data_source="Regression source 3",
        ),
        ChartPageContent(
            title="Four-chart dashboard title",
            charts=[doughnut(f"Question {index + 1}") for index in range(4)],
            layout=LayoutType.DASHBOARD,
            insights=["Four-chart insight one", "Four-chart insight two"],
            data_source="Regression source 4",
        ),
        ChartPageContent(
            title="Mixed chart and insight sidebar title",
            charts=[bar("Decision factors")],
            layout=LayoutType.MIXED,
            insights=["Sidebar insight one", "Sidebar insight two", "Sidebar insight three"],
            data_source="Regression source 5",
        ),
    ]


def _matrix_page() -> MultiGroupBarPageContent:
    segments = ["Total", "Segment A", "Segment B"]
    groups = []
    for group_index in range(2):
        rows = []
        for option_index in range(4):
            rows.append({
                "\u9009\u9879": f"Option {group_index + 1}-{option_index + 1}",
                "Total": 0.25 + option_index * 0.08,
                "Segment A": 0.31 + option_index * 0.06,
                "Segment B": 0.22 + option_index * 0.07,
            })
        groups.append({"title": f"Question group {group_index + 1}", "data": pd.DataFrame(rows)})
    return MultiGroupBarPageContent(
        title="Matrix and multi-group chart title",
        groups_data=groups,
        segments=segments,
        insights=["Matrix insight one", "Matrix insight two"],
        data_source="Regression source 6",
    )


def _assert_header_geometry(renderer, slide, start: int, role: str, page_title: str) -> None:
    role_map = renderer._template_mapping["roles"][role]
    new_shapes = list(slide.shapes)[start:]
    title_shape = next(
        shape for shape in new_shapes
        if getattr(shape, "has_text_frame", False) and page_title in shape.text
    )
    title_sizes = [
        run.font.size.pt
        for paragraph in title_shape.text_frame.paragraphs
        for run in paragraph.runs if run.font.size is not None
    ]
    expected_cap = role_map.get("title_font_size_pt") or 28
    if title_sizes and max(title_sizes) > expected_cap + 0.1:
        raise AssertionError(f"Title font was not capped: {title_sizes}, cap={expected_cap}")
    header_art_left = role_map.get("header_art_left")
    if header_art_left:
        safe_title_right = int((float(header_art_left) - 0.015) * renderer._slide_w)
        if int(title_shape.left + title_shape.width) > safe_title_right:
            raise AssertionError("Title intrudes into right-side header artwork")

    dividers = [
        shape for shape in new_shapes
        if not str(getattr(shape, "text", "") or "").strip()
        and not getattr(shape, "has_chart", False)
        and not getattr(shape, "has_table", False)
        and int(shape.width) >= renderer._slide_w * 0.45
        and abs(int(shape.height)) <= renderer._slide_h * 0.01
    ]
    body_tops = [
        int(shape.top) for shape in new_shapes
        if getattr(shape, "has_text_frame", False)
        and str(getattr(shape, "text", "") or "").strip()
        and renderer._slide_h * 0.12 <= int(shape.top) < renderer._slide_h * 0.90
    ]
    if dividers:
        required_title_gap = int(renderer._slide_h * 0.010)
        if min(int(shape.top) for shape in dividers) < int(title_shape.top + title_shape.height) + required_title_gap:
            raise AssertionError("Title divider is too close to the title text box")
    if dividers and body_tops:
        if max(int(shape.top + shape.height) for shape in dividers) >= min(body_tops):
            raise AssertionError("Title divider overlaps body text")
    header_art_bottom = role_map.get("header_art_bottom")
    if dividers and header_art_bottom:
        required_top = int((float(header_art_bottom) + 0.012) * renderer._slide_h)
        if min(int(shape.top) for shape in dividers) < required_top:
            raise AssertionError(
                f"Divider is too close to header artwork: {role_map}, required={required_top}"
            )
    for shape in new_shapes:
        if int(shape.left) < -renderer._slide_w * 0.02 or int(shape.top) < -renderer._slide_h * 0.02:
            raise AssertionError(f"Shape starts outside slide: {shape.name}")
        if int(shape.left + shape.width) > renderer._slide_w * 1.02:
            raise AssertionError(f"Shape exceeds right edge: {shape.name}")
        if int(shape.top + shape.height) > renderer._slide_h * 1.02:
            raise AssertionError(f"Shape exceeds bottom edge: {shape.name}")


def main() -> None:
    output = Path(sys.argv[2]) if len(sys.argv) > 2 else Path("qa_template_title_fixed.pptx")
    if len(sys.argv) > 1:
        template_path = Path(sys.argv[1])
    else:
        template_path = output.with_name("synthetic_title_template.pptx")
        _synthetic_template(template_path)

    mapping = build_template_mapping(Presentation(str(template_path)))
    chart_role = mapping["roles"].get("chart") or mapping["roles"]["content"]
    title_zone = chart_role["zones"]["title"]
    if not title_zone or title_zone["h"] > 0.14 or title_zone["y"] + title_zone["h"] > 0.17:
        raise AssertionError(f"Unsafe title zone: {chart_role}")
    matrix_role = mapping["roles"]["matrix"]
    if matrix_role["zones"]["content"]["w"] < 0.60:
        raise AssertionError(f"Unsafe matrix layout did not fall back: {matrix_role}")

    renderer = ReportRenderer(theme=Theme(), template_path=str(template_path))
    renderer._prs = renderer._build_presentation()
    dims = (renderer._slide_w, renderer._slide_h)

    for page in _chart_pages():
        slide = renderer._blank_slide("chart")
        if chart_role.get("header_art_bottom") and not any(int(shape.shape_type) == 13 for shape in slide.shapes):
            raise AssertionError("Slide-level header logo copy is missing")
        start = len(slide.shapes)
        build_chart_page(slide, page, renderer.theme, dims)
        renderer._remove_duplicate_title_divider(slide, start, "chart")
        renderer._apply_template_geometry(slide, start, "chart")
        _assert_header_geometry(renderer, slide, start, "chart", page.title)
        if len(page.charts) == 4:
            compact_charts = [shape for shape in list(slide.shapes)[start:] if getattr(shape, "has_chart", False)]
            for chart_shape in compact_charts:
                xml = chart_shape.chart._chartSpace.xml
                if 'h val="0.68"' not in xml:
                    raise AssertionError("Compact doughnut plot area is still too small")
        if page.layout == LayoutType.MIXED:
            sidebar_text = next(
                shape for shape in list(slide.shapes)[start:]
                if getattr(shape, "has_text_frame", False) and "Sidebar insight one" in shape.text
            )
            if int(sidebar_text.left) < renderer._slide_w * 0.50:
                raise AssertionError("Mixed-layout insights were not moved into the sidebar")

    matrix_page = _matrix_page()
    slide = renderer._blank_slide("matrix")
    start = len(slide.shapes)
    build_multi_group_bar_page(slide, matrix_page, renderer.theme, dims)
    renderer._remove_duplicate_title_divider(slide, start, "matrix")
    renderer._apply_template_geometry(slide, start, "matrix")
    _assert_header_geometry(renderer, slide, start, "matrix", matrix_page.title)
    matrix_shapes = list(slide.shapes)[start:]
    matrix_insight = next(
        shape for shape in matrix_shapes
        if getattr(shape, "has_text_frame", False) and "Matrix insight one" in shape.text
    )
    matrix_table = next(shape for shape in matrix_shapes if getattr(shape, "has_table", False))
    if int(matrix_table.top) < int(matrix_insight.top + matrix_insight.height) + int(renderer._slide_h * 0.005):
        raise AssertionError("Matrix table is too close to or overlaps the insight box")

    renderer._prs.save(output)
    print(f"Template title multi-layout smoke passed: {output}")
    print(
        f"slides={len(renderer._prs.slides)}, title={chart_role.get('title_font_size_pt')} pt, "
        f"title_zone={title_zone}, header_art_bottom={chart_role.get('header_art_bottom')}"
    )


if __name__ == "__main__":
    main()