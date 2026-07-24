from pathlib import Path
from tempfile import TemporaryDirectory
import sys

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches

from pptx_report.common.qa import inspect_presentation


def main() -> None:
    with TemporaryDirectory(prefix="surveykit-object-qa-") as temp_dir:
        output = Path(temp_dir) / "qa-smoke.pptx"
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        for slide_number in range(1, 5):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(9.5), Inches(0.7))
            title.text = f"QA slide {slide_number}"
            if slide_number == 1:
                body = slide.shapes.add_textbox(Inches(0.8), Inches(1.10), Inches(8.5), Inches(0.85))
                body.text = "This text box is crossed by duplicate divider lines"
                slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.6), Inches(1.45), Inches(12.5), Inches(1.45))
                slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.6), Inches(1.48), Inches(12.5), Inches(1.48))
        prs.save(output)

        qa = inspect_presentation(output).to_dict()
        codes = [issue["code"] for issue in qa["issues"]]
        assert "duplicate_divider_lines" in codes, qa
        assert "divider_intersects_text_box" in codes, qa
        assert "source_missing" in codes, qa
        assert qa["warning_count"] == len(qa["issues"])
        assert qa["error_count"] == 0 and qa["ok"] is True
        assert qa["slides_with_issues"] == [1, 4]
        assert qa["score"] < 100
    print("object QA smoke: ok")


if __name__ == "__main__":
    main()
