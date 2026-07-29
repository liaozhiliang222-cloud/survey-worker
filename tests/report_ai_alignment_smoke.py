"""Regression checks for chart/evidence/AI alignment."""

from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))
from pptx import Presentation  # noqa: E402

from pptx_report.build_jd_report import _repair_single_choice_dimension_outliers  # noqa: E402
from pptx_report.wizard import _build_multi_group_bar_page  # noqa: E402
from pptx_report.model import DataFact  # noqa: E402
from pptx_report.wizard import _research_model_semantics  # noqa: E402
from pptx_report.pages import build_multi_group_bar_page  # noqa: E402
from pptx_report.theme import theme_from_key  # noqa: E402
from pptx_report.wizard import _build_narrative_findings  # noqa: E402


def test_isolated_dimension_leak_is_repaired() -> None:
    question = {
        "categories": ["3+", "2", "1"],
        "segments": ["Total", "high_intent", "commuter"],
        "data": {
            "Total": [0.012, 0.062, 0.926],
            "high_intent": [0.694, 0.049, 0.944],
            "commuter": [0.418, 0.054, 0.941],
        },
    }
    _repair_single_choice_dimension_outliers(question)
    assert all(abs(a - b) < 1e-9 for a, b in zip(question["data"]["high_intent"], [0.007, 0.049, 0.944]))
    assert all(abs(a - b) < 1e-9 for a, b in zip(question["data"]["commuter"], [0.005, 0.054, 0.941]))
    assert len(question["data_quality_warnings"]) == 2


def test_multi_select_distribution_is_not_normalized() -> None:
    question = {
        "categories": ["A", "B", "C"],
        "segments": ["Total", "segment_a"],
        "data": {
            "Total": [0.8, 0.7, 0.4],
            "segment_a": [0.9, 0.8, 0.5],
        },
    }
    original = list(question["data"]["segment_a"])
    _repair_single_choice_dimension_outliers(question)
    assert question["data"]["segment_a"] == original
    assert "data_quality_warnings" not in question


def test_psm_semantics_forbid_single_curve_threshold_claims() -> None:
    too_cheap = _research_model_semantics({
        "title": "\u4ef7\u683c\u4f4e\u5230\u591a\u5c11\u65f6\u4f1a\u89c9\u5f97\u592a\u4fbf\u5b9c\u4ee5\u81f3\u4e8e\u6000\u7591\u54c1\u8d28",
    })
    assert too_cheap["analysis_model"] == "psm"
    assert too_cheap["metric_role"] == "too_cheap"
    assert any("\u4e0d\u662f\u8d2d\u4e70\u63a5\u53d7\u7387" in rule for rule in too_cheap["interpretation_rules"])
    assert any("\u4ea4\u70b9" in rule for rule in too_cheap["interpretation_rules"])


def test_kano_semantics_require_mean_crosshair() -> None:
    kano = _research_model_semantics({"title": "KANO Better/Worse \u7cfb\u6570\u77e9\u9635"})
    assert kano["analysis_model"] == "kano"
    assert any("\u5e73\u5747\u503c" in rule for rule in kano["interpretation_rules"])



def test_sub_one_percent_is_not_scaled_twice() -> None:
    question = {
        "code": "S2A1",
        "title": "Vehicle ownership count",
        "categories": ["3+", "2", "1"],
        "segments": ["Total", "high_intent"],
        "data": {
            "Total": [0.012, 0.062, 0.926],
            "high_intent": [0.00694, 0.049, 0.94406],
        },
        "base": {"Total": 500, "high_intent": 120},
    }
    page = _build_multi_group_bar_page(
        [question], ["Total", "high_intent"], "synthetic", 1, 1
    )
    frame = page.groups_data[0]["data"]
    row = frame.loc[frame["\u9009\u9879"] == "3+"].iloc[0]
    assert abs(float(row["high_intent"]) - 0.00694) < 1e-9


def test_system_findings_prefer_validated_slide_brief() -> None:
    fact = DataFact(
        fact_id="F1", fact_type="top_rank", question_id="Q1",
        metric_name="purchase_intent", value=0.72,
    )
    findings = _build_narrative_findings({
        "pages": [{
            "slide_brief": {
                "claim": "Decision certainty is the main conversion barrier",
                "business_implication": "Prioritize proof and reassurance",
                "evidence_fact_ids": ["F1"],
                "evidence_question_ids": ["Q1"],
            },
        }],
    }, [fact], [])
    assert len(findings) == 1
    assert findings[0].title == "Decision certainty is the main conversion barrier"
    assert findings[0].evidence_fact_ids == ["F1"]

    question = {
        "code": "S2A1",
        "title": "Vehicle ownership count",
        "categories": ["3+", "2", "1"],
        "segments": ["Total", "high_intent"],
        "data": {
            "Total": [0.012, 0.062, 0.926],
            "high_intent": [0.00694, 0.049, 0.94406],
        },
        "base": {"Total": 500, "high_intent": 120},
    }
    page = _build_multi_group_bar_page([question], ["Total", "high_intent"], "synthetic", 1, 1)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build_multi_group_bar_page(
        slide, page, theme_from_key("blue"),
        (presentation.slide_width, presentation.slide_height),
    )
    chart = next(shape.chart for shape in slide.shapes if getattr(shape, "has_chart", False))
    assert chart.value_axis.minimum_scale == 0.0
def main() -> None:
    test_isolated_dimension_leak_is_repaired()
    test_multi_select_distribution_is_not_normalized()
    test_psm_semantics_forbid_single_curve_threshold_claims()
    test_kano_semantics_require_mean_crosshair()
    test_sub_one_percent_is_not_scaled_twice()
    print("report AI alignment smoke passed")
if __name__ == "__main__":
    test_system_findings_prefer_validated_slide_brief()
    main()
