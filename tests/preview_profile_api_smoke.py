"""API smoke checks for persisted template roles and preview capability guard."""

from __future__ import annotations

from io import BytesIO
import json
from pathlib import Path
import sys
import tempfile

from fastapi.testclient import TestClient
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
from deploy import aliyun_api as api  # noqa: E402
from pptx_report.common.qa import inspect_presentation  # noqa: E402


def template_bytes() -> bytes:
    prs = Presentation()
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = "Customer Research Report"
    chart = prs.slides.add_slide(prs.slide_layouts[1])
    chart.shapes.title.text = "Key finding"
    chart.placeholders[1].text = "Evidence and implications"
    stream = BytesIO()
    prs.save(stream)
    return stream.getvalue()


def main() -> None:
    with tempfile.TemporaryDirectory(prefix="surveykit-template-profile-") as temp_dir:
        api.TEMPLATE_DIR = Path(temp_dir)
        client = TestClient(api.app)
        uploaded = client.post(
            "/api/pptx-report/templates",
            content=template_bytes(),
            headers={
                "Content-Type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "X-Template-Name": "profile-smoke.pptx",
            },
        )
        assert uploaded.status_code == 200, uploaded.text
        payload = uploaded.json()
        template_id = payload["template_id"]
        assert payload["profile"]["version"] == 1
        assert payload["recommended_roles"]

        updated = client.put(
            f"/api/pptx-report/templates/{template_id}/profile",
            json={"name": "confirmed-profile", "roles": {"chart": 2, "summary": 2}},
        )
        assert updated.status_code == 200, updated.text
        profile = updated.json()
        assert profile["version"] == 2
        assert profile["roles"]["chart"]["slide_index"] == 1
        assert profile["roles"]["chart"]["user_confirmed"] is True

        loaded = client.get(f"/api/pptx-report/templates/{template_id}/profile")
        assert loaded.status_code == 200
        assert loaded.json()["roles"]["summary"]["slide_index"] == 1

        original_converter = api._office_converter
        api._office_converter = lambda: None
        try:
            metadata = json.dumps({"pages": [1]}).encode("utf-8")
            envelope = api.REQUEST_ENVELOPE_MAGIC + len(metadata).to_bytes(4, "big") + metadata + b"xlsx"
            preview = client.post(
                "/api/pptx-report/preview-render",
                content=envelope,
                headers={"Content-Type": "application/vnd.surveykit.pptx-request"},
            )
            assert preview.status_code == 503, preview.text
            assert "LibreOffice" in preview.json()["error"]["message"]
        finally:
            api._office_converter = original_converter

    with tempfile.TemporaryDirectory(prefix="surveykit-object-qa-") as temp_dir:
        qa_path = Path(temp_dir) / "qa.pptx"
        logo_path = Path(temp_dir) / "logo.png"
        Image.new("RGB", (80, 80), "red").save(logo_path)
        qa_prs = Presentation()
        slide = qa_prs.slides.add_slide(qa_prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(1.2))
        title.text = "First line\nSecond line\nThird line"
        slide.shapes.add_picture(str(logo_path), Inches(5.5), Inches(0.2), Inches(1), Inches(1))
        qa_prs.save(qa_path)
        qa_result = inspect_presentation(qa_path, require_sources=False)
        codes = {issue["code"] for issue in qa_result.issues}
        assert "title_over_two_lines" in codes
        assert "template_logo_overlaps_title" in codes
    print("template profile and preview API smoke: ok")


if __name__ == "__main__":
    main()
