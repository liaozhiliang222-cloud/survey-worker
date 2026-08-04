"""Smoke tests for one-slide editable research model chart exports."""

from io import BytesIO
from pathlib import Path
import sys
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from pptx_report.model_chart_export import render_model_chart_pptx


SAMPLES = {
    "psm": {
        "sampleCount": 200,
        "acceptable": "¥120 - ¥260",
        "pmc": 120,
        "pme": 260,
        "ipp": 190,
        "opp": 180,
        "curve": [
            {"price": 100, "tooCheap": 80, "cheap": 90, "expensive": 10, "tooExpensive": 2},
            {"price": 180, "tooCheap": 35, "cheap": 60, "expensive": 45, "tooExpensive": 20},
            {"price": 280, "tooCheap": 5, "cheap": 15, "expensive": 85, "tooExpensive": 75},
        ],
    },
    "kano": [
        {"name": "配送速度", "classification": "期望属性", "better": 0.72, "worse": -0.66},
        {"name": "包装设计", "classification": "魅力属性", "better": 0.65, "worse": -0.25},
        {"name": "包装质感", "classification": "魅力属性", "better": 0.38, "worse": -0.69},
        {"name": "售后保障", "classification": "必备属性", "better": 0.42, "worse": -0.81},
    ],
    "maxdiff": [
        {"item": "口味更好", "score": 0.45},
        {"item": "价格更划算", "score": 0.28},
        {"item": "售后更安心", "score": -0.18},
    ],
    "driver": {
        "r2": 0.62,
        "n": 320,
        "dvName": "整体满意度",
        "results": [
            {"name": "产品质量", "importance": 0.34},
            {"name": "售后服务", "importance": 0.26},
            {"name": "价格合理性", "importance": 0.18},
        ],
    },
}


def main() -> None:
    for model_type, data in SAMPLES.items():
        blob = render_model_chart_pptx(model_type, data)
        assert blob.startswith(b"PK")
        presentation = Presentation(BytesIO(blob))
        assert len(presentation.slides) == 1
        shapes = list(presentation.slides[0].shapes)
        charts = [shape for shape in shapes if getattr(shape, "has_chart", False)]
        assert len(charts) == 1, model_type
        assert all(shape.shape_type != MSO_SHAPE_TYPE.PICTURE for shape in shapes)
        text = "\n".join(shape.text for shape in shapes if getattr(shape, "has_text_frame", False))
        assert "可编辑" in text
        assert charts[0].chart.part.chart_workbook.xlsx_part is not None
        if model_type == "kano":
            chart = charts[0].chart
            assert chart.has_legend is False
            assert len(chart.series) == 1
            series = chart.series[0]
            assert series.name == "better"
            assert len(series.points) == len(SAMPLES["kano"])
            assert series.format.line.fill.type == MSO_FILL_TYPE.BACKGROUND

            data_labels = series._ser.find(qn("c:dLbls"))
            assert data_labels is not None
            assert len(data_labels.findall(qn("c:dLbl"))) == len(SAMPLES["kano"])
            assert data_labels.find(qn("c:showSerName")).get("val") == "0"
            assert data_labels.find(qn("c:showVal")).get("val") == "0"
            assert data_labels.find(qn("c:dLblPos")).get("val") == "r"
            c15_ns = "http://schemas.microsoft.com/office/drawing/2012/chart"
            data_range = series._ser.find(f".//{{{c15_ns}}}datalabelsRange")
            assert data_range is not None
            formula = data_range.find(f"{{{c15_ns}}}f")
            assert formula.text == "Sheet1!$C$2:$C$5"

            xlsx_part = chart.part.chart_workbook.xlsx_part
            with ZipFile(BytesIO(xlsx_part.blob), "r") as workbook:
                sheet_xml = workbook.read("xl/worksheets/sheet1.xml").decode("utf-8")
                strings_xml = workbook.read("xl/sharedStrings.xml").decode("utf-8")
            assert 'ref="A1:C5"' in sheet_xml
            for cell in ("A1", "C1", "C2", "C3", "C4", "C5"):
                assert f'r="{cell}"' in sheet_xml
            assert "worse" in strings_xml
            assert "better" in strings_xml
            assert "KANO\u5c5e\u6027\u6807\u7b7e" in strings_xml
            for row in SAMPLES["kano"]:
                assert row["name"] in strings_xml

            assert chart.category_axis.has_major_gridlines is False
            assert chart.value_axis.has_major_gridlines is False
            assert chart.category_axis.has_minor_gridlines is False
            assert chart.value_axis.has_minor_gridlines is False
            assert abs(chart.category_axis.crosses_at - 0.6025) < 1e-6
            assert abs(chart.value_axis.crosses_at - 0.5425) < 1e-6
            plot_area = chart._chartSpace.chart.plotArea
            plot_border = plot_area.find(qn("c:spPr"))
            assert plot_border is not None
            assert plot_border.find(qn("a:ln")) is not None
            for label in ("魅力属性", "期望属性", "必备属性", "无差异属性", "Worse 均值", "Better 均值"):
                assert label in text

    export_source = (ROOT / "pptx_report" / "model_chart_export.py").read_text(encoding="utf-8")
    assert "_add_matrix_line" not in export_source

    api_source = (ROOT / "deploy" / "aliyun_api.py").read_text(encoding="utf-8")
    assert '/api/pptx-report/model-chart' in api_source
    assert 'X-SurveyKit-Editable-Chart' in api_source
    print("editable model chart PPTX smoke: ok")


if __name__ == "__main__":
    main()
