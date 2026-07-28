from pathlib import Path
from tempfile import TemporaryDirectory
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptx import Presentation
from pptx.util import Inches, Pt

from pptx_report.template import analyze_template


def add_text(slide, text, x, y, w, h, size):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)


def add_divider(prs, number, title, topics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, number, 0.6, 1.2, 2.0, 1.0, 44 if "." not in number else 12)
    add_text(slide, title, 4.5, 2.4, 7.0, 0.8, 34)
    add_text(slide, " | ".join(topics), 4.5, 3.3, 7.5, 0.6, 16)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    while prs.slides:
        rel_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rel_id)
        del prs.slides._sldIdLst[0]
    add_divider(prs, "01", "项目概述", ["项目背景", "样本说明"])
    add_divider(prs, "02", "主要研究发现", ["产品概念测试结果", "目标用户画像"])
    add_divider(prs, "2.1", "产品概念测试结果", ["市场接受度", "功能偏好"])
    add_divider(prs, "2.2", "目标用户画像", ["基础特征", "使用行为特征"])
    add_divider(prs, "03", "结论与建议", ["核心结论", "行动建议"])
    with TemporaryDirectory() as temp_dir:
        path = Path(temp_dir) / "template.pptx"
        prs.save(path)
        structure = analyze_template(str(path))["report_structure"]
    assert structure["confidence"] >= 0.85
    assert [item["title"] for item in structure["sections"]] == ["项目概述", "主要研究发现", "结论与建议"]
    assert [item["title"] for item in structure["sections"][1]["subsections"]] == ["产品概念测试结果", "目标用户画像"]
    assert structure["sections"][1]["subsections"][0]["topics"] == ["市场接受度", "功能偏好"]
    print("template report structure smoke passed")


if __name__ == "__main__":
    main()